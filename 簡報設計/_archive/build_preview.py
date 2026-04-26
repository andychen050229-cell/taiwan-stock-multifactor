# -*- coding: utf-8 -*-
"""
build_preview.py — 多因子股票預測系統 · 簡報重建（preview v1）

策略：
  - 1920 x 1080 畫布，重現 PDF（Claude Design 輸出）視覺結構
  - 文字採用「現網站 v11.5.17」之 canonical 數字（取代 PDF 內舊資料）
  - 預覽範圍：8 張（Cover / TOC / Sec01 Divider / Exec Summary /
                Three Commitments / Sec02 Divider / Problem / Approach）

字型：Inter (latin) + 微軟正黑體 (CJK) + Consolas (mono fallback / JetBrains Mono)
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- canvas ----------
PX = 9525  # EMU per pixel @ 96 dpi
W_PX, H_PX = 1920, 1080

# ---------- palette ----------
NAVY      = RGBColor(0x1e, 0x3a, 0x8a)
NAVY_D    = RGBColor(0x1e, 0x40, 0xaf)
NAVY_DK   = RGBColor(0x0b, 0x12, 0x20)
NAVY_MID  = RGBColor(0x16, 0x21, 0x38)
NAVY_LT   = RGBColor(0xe0, 0xe7, 0xff)
TEAL      = RGBColor(0x0d, 0x94, 0x88)
TEAL_D    = RGBColor(0x0f, 0x76, 0x6e)
TEAL_BG   = RGBColor(0xf0, 0xfd, 0xfa)
MINT      = RGBColor(0x5e, 0xea, 0xd4)
MINT_D    = RGBColor(0x2d, 0xd4, 0xbf)
EMERALD   = RGBColor(0x05, 0x96, 0x69)
EMERALD_DK= RGBColor(0x04, 0x78, 0x57)
EMERALD_BG= RGBColor(0xec, 0xfd, 0xf5)
ROSE      = RGBColor(0xe1, 0x1d, 0x48)
ROSE_BG   = RGBColor(0xfd, 0xf2, 0xf8)
GOLD      = RGBColor(0xf5, 0x9e, 0x0b)
GOLD_BG   = RGBColor(0xfe, 0xf3, 0xc7)
ORANGE    = RGBColor(0xea, 0x58, 0x0c)
ORANGE_L  = RGBColor(0xfb, 0x92, 0x3c)
VIOLET    = RGBColor(0x7c, 0x3a, 0xed)
CYAN      = RGBColor(0x06, 0xb6, 0xd4)
INK       = RGBColor(0x0f, 0x17, 0x2a)
INK_2     = RGBColor(0x33, 0x41, 0x55)
INK_3     = RGBColor(0x64, 0x74, 0x8b)
INK_4     = RGBColor(0x94, 0xa3, 0xb8)
PAPER     = RGBColor(0xff, 0xff, 0xff)
BG        = RGBColor(0xfa, 0xfb, 0xfc)
SUBTLE    = RGBColor(0xf1, 0xf5, 0xf9)
TINT      = RGBColor(0xf8, 0xfa, 0xfc)
RULE      = RGBColor(0xe2, 0xe8, 0xf0)
RULE_2    = RGBColor(0xcb, 0xd5, 0xe1)
IVORY     = RGBColor(0xf5, 0xf1, 0xe8)
CREAM     = RGBColor(0xfb, 0xf8, 0xf1)

LATIN_FONT = 'Inter'
EA_FONT    = 'Microsoft JhengHei'
MONO_FONT  = 'Consolas'  # JetBrains Mono fallback for Windows


def _set_run_typeface(run, latin=LATIN_FONT, ea=EA_FONT, mono=False):
    """Set latin/east-asian typefaces explicitly via XML."""
    rPr = run._r.get_or_add_rPr()
    # remove any existing
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    typeface_latin = MONO_FONT if mono else latin
    typeface_ea = ea
    el = etree.SubElement(rPr, qn('a:latin')); el.set('typeface', typeface_latin)
    el = etree.SubElement(rPr, qn('a:ea'));    el.set('typeface', typeface_ea)
    el = etree.SubElement(rPr, qn('a:cs'));    el.set('typeface', typeface_ea)


def _set_run_letterspacing(run, ls_pts):
    """Letter spacing in points; PowerPoint stores in 1/100 pt."""
    if not ls_pts:
        return
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(ls_pts * 100)))


# ---------- shape helpers ----------
def rect(slide, x, y, w, h, fill=None, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x*PX, y*PX, w*PX, h*PX)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def stripe(slide, x, y, w, h, color):
    return rect(slide, x, y, w, h, fill=color)


def dot(slide, cx, cy, r, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, (cx-r)*PX, (cy-r)*PX, (2*r)*PX, (2*r)*PX)
    shp.shadow.inherit = False
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


# ---------- text helpers ----------
def text(slide, x, y, w, h, content, *, size=18, color=INK, bold=False,
         align='left', anchor='top', mono=False, italic=False,
         ls=0, line_h=1.2, latin=LATIN_FONT, ea=EA_FONT):
    box = slide.shapes.add_textbox(x*PX, y*PX, w*PX, h*PX)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {
        'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM
    }[anchor]
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    if line_h:
        p.line_spacing = line_h
    run = p.add_run()
    run.text = content
    rf = run.font
    rf.size = Pt(size); rf.bold = bold; rf.italic = italic; rf.color.rgb = color
    _set_run_typeface(run, latin=latin, ea=ea, mono=mono)
    _set_run_letterspacing(run, ls)
    return box


def rich(slide, x, y, w, h, runs, *, align='left', anchor='top', line_h=1.2):
    """runs: list of dicts {text, size, bold, color, mono, italic, ls, latin, ea}.
    Use {'br': True} for newline, {'sp': True} for paragraph break."""
    box = slide.shapes.add_textbox(x*PX, y*PX, w*PX, h*PX)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {
        'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM
    }[anchor]
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    if line_h:
        p.line_spacing = line_h

    def new_para():
        nonlocal p
        p = tf.add_paragraph()
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
        if line_h:
            p.line_spacing = line_h

    for r in runs:
        if r.get('sp'):
            new_para()
            continue
        if r.get('br'):
            # soft line break
            run = p.add_run()
            run.text = ''
            br = etree.SubElement(run._r.getparent(), qn('a:br'))
            run._r.addnext(br)
            continue
        run = p.add_run()
        run.text = r.get('text', '')
        rf = run.font
        rf.size = Pt(r.get('size', 18))
        rf.bold = r.get('bold', False)
        rf.italic = r.get('italic', False)
        rf.color.rgb = r.get('color', INK)
        _set_run_typeface(run, latin=r.get('latin', LATIN_FONT), ea=r.get('ea', EA_FONT),
                          mono=r.get('mono', False))
        _set_run_letterspacing(run, r.get('ls', 0))
    return box


# ---------- chrome (top header / bottom foot) ----------
def slide_chrome(slide, page_title, page_no, band_active=None, total=8, source=None):
    """Top: TW marker + brand · right side title & page no.
    Bottom: source on left, 8-i progress band on right."""
    # Top chrome
    rect(slide, 96, 38, 30, 30, fill=NAVY)  # TW marker
    text(slide, 96, 38, 30, 30, 'TW', size=11, bold=True, color=PAPER,
         align='center', anchor='middle', mono=True, ls=0)
    text(slide, 134, 38, 360, 30, '多因子股票預測系統', size=12, bold=True,
         color=INK_2, anchor='middle', ls=1.4)
    # right side title/page
    text(slide, 1320, 38, 480, 30, page_title, size=12, color=INK_3,
         align='right', anchor='middle', mono=True, ls=2.0,
         latin=LATIN_FONT)
    rect(slide, 1810, 45, 1, 16, fill=RULE_2)
    text(slide, 1820, 38, 60, 30, page_no, size=12, color=INK_2,
         bold=True, align='right', anchor='middle', mono=True)

    # Bottom foot
    if source:
        text(slide, 96, 1018, 1300, 30, 'Source：' + source, size=11,
             color=INK_3, anchor='middle', mono=True, ls=0.5)
    # progress band 8-i
    band_x = 1700
    for i in range(total):
        on = (band_active is not None and i == band_active)
        rect(slide, band_x + i*16, 1030, 12, 3,
             fill=(TEAL if on else RULE))


def add_blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # white background
    bg = rect(slide, 0, 0, W_PX, H_PX, fill=PAPER)
    return slide


# ============================================================
# SLIDE 01 · COVER
# ============================================================
def slide_cover(prs):
    slide = add_blank_slide(prs)
    # subtle navy stripe at very top
    stripe(slide, 0, 0, W_PX, 6, NAVY)
    # cover top labels
    text(slide, 80, 38, 600, 22, 'Taiwan Equity · Quant Research Terminal · 2026',
         size=12, bold=True, color=INK_3, mono=True, ls=2.0)
    text(slide, 1240, 38, 600, 22, 'Big Data & Business Analytics · NCCU',
         size=12, bold=True, color=INK_3, align='right', mono=True, ls=2.0)

    # eyebrow chip
    text(slide, 80, 144, 700, 24, '— 9-PILLAR FACTOR FRAMEWORK · GLINT-LIGHT · v11.5.17',
         size=14, bold=True, color=TEAL, mono=True, ls=2.4)

    # Hero title (left column)
    rich(slide, 80, 188, 900, 230, [
        {'text': '台股多因子', 'size': 96, 'bold': True, 'color': NAVY, 'ls': -2.5},
        {'br': True},
        {'text': '預測系統', 'size': 96, 'bold': True, 'color': TEAL, 'ls': -2.5},
    ], line_h=1.0)

    # rule under title
    rect(slide, 80, 444, 80, 4, fill=NAVY)

    # subhead block (left)
    rich(slide, 80, 472, 900, 260, [
        {'text': '量化研究終端機 · Strategic Research Deck',
         'size': 22, 'bold': True, 'color': NAVY},
        {'sp': True},
        {'text': '從 ', 'size': 18, 'color': INK_2},
        {'text': '1,930', 'size': 18, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' 檔上市櫃 × ', 'size': 18, 'color': INK_2},
        {'text': '948,976', 'size': 18, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' 樣本 × ', 'size': 18, 'color': INK_2},
        {'text': '505', 'size': 18, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' 交易日淬煉 9 支柱因子棋盤', 'size': 18, 'color': INK_2},
        {'sp': True},
        {'text': '於 D+20 方向預測達到 ', 'size': 18, 'color': INK_2},
        {'text': 'OOS AUC 0.6490', 'size': 18, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' · ICIR ', 'size': 18, 'color': INK_2},
        {'text': '0.7431', 'size': 18, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' · DSR ', 'size': 18, 'color': INK_2},
        {'text': '12.12', 'size': 18, 'bold': True, 'color': NAVY, 'mono': True},
        {'sp': True},
        {'text': 'Top 0.1% 選股 edge ', 'size': 18, 'color': INK_2},
        {'text': '+19.5 pp', 'size': 18, 'bold': True, 'color': TEAL, 'mono': True},
        {'text': ' · 9/9 治理閘門全數通過', 'size': 18, 'color': INK_2},
    ], line_h=1.55)

    # ---- KPI grid right (3 cols × 2 rows) ----
    grid_x, grid_y = 1010, 188
    cell_w, cell_h = 280, 220
    # outer border
    rect(slide, grid_x, grid_y, cell_w*3, cell_h*2, fill=PAPER, line=RULE, line_w=0.75)
    kpis = [
        ('OOS AUC',   '0.6490', 'D+20 horizon · XGBoost',   NAVY),
        ('ICIR',      '0.7431', 'Mean IC / σ',              TEAL),
        ('DSR',       '12.12',  'p ≪ 0.001 · 閾 2.0',       NAVY),
        ('FEATURES',  '91',     '/ 1,623 · 保留 5.6%',       NAVY),
        ('TOP 0.1%',  '45.79%', 'edge +19.5 pp',            TEAL),
        ('GATES',     '9/9',    'Governance PASS',          EMERALD),
    ]
    for i, (lbl, val, sub, accent) in enumerate(kpis):
        col, row = i % 3, i // 3
        cx = grid_x + col*cell_w
        cy = grid_y + row*cell_h
        # top stripe
        stripe(slide, cx, cy, cell_w, 4, accent)
        # cell separator lines
        if col < 2:
            rect(slide, cx + cell_w - 0.5, cy, 0.5, cell_h, fill=RULE)
        if row < 1:
            rect(slide, cx, cy + cell_h - 0.5, cell_w, 0.5, fill=RULE)
        # label
        text(slide, cx + 24, cy + 32, cell_w - 48, 22, lbl,
             size=11, bold=True, color=INK_3, mono=True, ls=2.4)
        # value
        text(slide, cx + 24, cy + 64, cell_w - 48, 86, val,
             size=44, bold=True, color=accent, mono=True, ls=-1.2,
             anchor='top')
        # sub
        text(slide, cx + 24, cy + 158, cell_w - 48, 40, sub,
             size=11, color=INK_3, mono=True, ls=0.6)

    # ---- bottom team / version foot ----
    rect(slide, 80, 880, 1760, 1, fill=RULE)  # divider
    text(slide, 80, 904, 360, 20, 'TEAM',
         size=10, bold=True, color=INK_4, mono=True, ls=2.6)
    text(slide, 80, 928, 700, 26, '台股多因子 Research Group · 7 人課程專案',
         size=16, bold=True, color=INK)

    text(slide, 760, 904, 200, 20, 'COURSE',
         size=10, bold=True, color=INK_4, mono=True, ls=2.6)
    text(slide, 760, 928, 380, 26, 'NCCU · Big Data & Business Analytics',
         size=16, bold=True, color=INK)

    text(slide, 1180, 904, 200, 20, 'VERSION',
         size=10, bold=True, color=INK_4, mono=True, ls=2.6)
    text(slide, 1180, 928, 200, 26, 'v11.5.17',
         size=16, bold=True, color=INK, mono=True)

    text(slide, 1480, 904, 360, 20, 'DATE',
         size=10, bold=True, color=INK_4, mono=True, ls=2.6, align='right')
    text(slide, 1480, 928, 360, 26, 'April 2026 · Academic Project',
         size=16, bold=True, color=INK, align='right')

    # bottom-right page no
    text(slide, 1780, 1018, 60, 30, '01', size=12, bold=True, color=INK_2,
         align='right', anchor='middle', mono=True)
    return slide


# ============================================================
# SLIDE 02 · TABLE OF CONTENTS
# ============================================================
def slide_toc(prs):
    slide = add_blank_slide(prs)
    slide_chrome(slide, 'Table of Contents', '02', band_active=None,
                 source='Team · 台股多因子 Research Group · April 2026')

    # eyebrow
    text(slide, 96, 138, 700, 22, '— CONTENTS · 導覽',
         size=14, bold=True, color=TEAL, mono=True, ls=2.4)

    # title
    rich(slide, 96, 168, 1700, 80, [
        {'text': '本日議程 · ', 'size': 50, 'bold': True, 'color': NAVY, 'ls': -1.2},
        {'text': '8 章節 ', 'size': 50, 'bold': True, 'color': TEAL, 'ls': -1.2},
        {'text': '· 共 35 頁完整呈現', 'size': 50, 'bold': True, 'color': NAVY, 'ls': -1.2},
    ], line_h=1.05)

    rect(slide, 96, 268, 80, 4, fill=TEAL)

    # TOC rows (8)
    rows = [
        ('01', 'Executive Summary',     '一頁摘要 · 三項核心發現 · 三項落地建議',           'P.03 — 05', NAVY),
        ('02', 'Problem & Approach',    '三重結構性難題 · 6 階段流水線 · 三大承諾',         'P.06 — 08', NAVY),
        ('03', 'Data & Pipeline',       '7 表 ETL 1,932 家 · 948,976 樣本 · 5-層技術棧',    'P.09 — 12', TEAL),
        ('04', '9-Pillar Framework',    'LOPO ΔAUC · Risk +138.6 bps · Top 10 因子',       'P.13 — 16', NAVY),
        ('05', 'Model Performance',     'AUC 0.6490 · DSR 12.12 · Top 0.1% 命中 45.79%',    'P.17 — 21', TEAL),
        ('06', 'Productionization',     'Streamlit Dashboard 13 頁 · 雙軌決策 0.40 / 0.50','P.22',      NAVY),
        ('07', 'Governance & Roadmap',  '9 / 9 治理閘門 · Limitations · Next Steps',         'P.23 — 25', NAVY),
        ('08', 'Appendix · Sentiment',  '500 keyword Lift · U-shape · News × Forum 共識',  'P.26 — 35', TEAL),
    ]
    base_y = 320
    row_h = 76
    for i, (no, title, desc, pages, accent) in enumerate(rows):
        ry = base_y + i * row_h
        # row underline
        if i < len(rows) - 1:
            rect(slide, 96, ry + row_h - 1, 1728, 1, fill=RULE)
        # number
        text(slide, 96, ry + 16, 80, 44, no, size=28, bold=True,
             color=accent, mono=True, ls=-0.6)
        # title
        text(slide, 200, ry + 18, 460, 30, title, size=20, bold=True,
             color=NAVY, ls=-0.2)
        # desc
        text(slide, 660, ry + 22, 950, 28, desc, size=14, color=INK_2)
        # pages
        text(slide, 1620, ry + 22, 200, 28, pages, size=13, bold=True,
             color=INK_3, align='right', mono=True, ls=1.2)

    # bottom-right page no
    text(slide, 1820, 1018, 60, 30, '02', size=12, bold=True, color=INK_2,
         align='right', anchor='middle', mono=True)
    return slide


# ============================================================
# SLIDE 03 · SECTION 01 DIVIDER
# ============================================================
def slide_section_divider(prs, no_str, line1, line2, sub, page_no, band_active):
    slide = add_blank_slide(prs)
    # solid dark navy bg
    rect(slide, 0, 0, W_PX, H_PX, fill=NAVY_DK)

    # subtle 96px horizontal grid using NAVY_MID (slightly lighter)
    for gy in range(96, H_PX, 96):
        rect(slide, 0, gy, W_PX, 1, fill=NAVY_MID)
    # vertical guide at content gutter
    rect(slide, 96, 0, 1, H_PX, fill=NAVY_MID)
    rect(slide, 1824, 0, 1, H_PX, fill=NAVY_MID)

    # top-right corner accent (mint thin L)
    rect(slide, 1480, 0, 360, 4, fill=MINT)
    rect(slide, 1836, 0, 4, 80, fill=MINT)

    # left edge mint vertical stripe (alongside title, evokes editorial section opener)
    rect(slide, 96, 360, 6, 540, fill=MINT)

    # bottom-right plate (small NAVY_MID rectangle) for editorial weight
    rect(slide, 1320, 800, 504, 200, fill=NAVY_MID)
    rect(slide, 1320, 800, 504, 4, fill=MINT)
    text(slide, 1340, 820, 480, 24, '— PROJECT', size=11, bold=True,
         color=MINT, mono=True, ls=2.6)
    rich(slide, 1340, 850, 480, 130, [
        {'text': '多因子股票預測系統', 'size': 18, 'bold': True, 'color': PAPER},
        {'br': True},
        {'text': 'v11.5.17 · 2026-04', 'size': 12, 'color': MINT, 'mono': True, 'ls': 1.4},
    ], line_h=1.5)

    # top brand chrome
    rect(slide, 120, 38, 30, 30, fill=PAPER)
    text(slide, 120, 38, 30, 30, 'TW', size=11, bold=True, color=NAVY,
         align='center', anchor='middle', mono=True)
    text(slide, 158, 38, 360, 30, '多因子股票預測系統', size=12, bold=True,
         color=IVORY, anchor='middle', ls=1.4)
    text(slide, 1320, 38, 380, 30, 'SECTION DIVIDER', size=11,
         color=MINT, align='right', anchor='middle', mono=True, ls=2.4)
    rect(slide, 1714, 45, 1, 16, fill=NAVY_MID)
    text(slide, 1724, 38, 100, 30, page_no, size=12, bold=True,
         color=IVORY, align='right', anchor='middle', mono=True)

    # No label
    text(slide, 130, 380, 600, 32, no_str, size=18, bold=True,
         color=MINT, mono=True, ls=2.8)

    # huge section title (two lines)
    rich(slide, 130, 428, 1180, 320, [
        {'text': line1, 'size': 88, 'bold': True, 'color': PAPER, 'ls': -2.6},
        {'br': True},
        {'text': line2, 'size': 88, 'bold': True, 'color': MINT, 'ls': -2.6},
    ], line_h=1.05)

    # subline
    text(slide, 130, 760, 1100, 80, sub, size=20, color=IVORY, line_h=1.5)

    # progress band on bottom
    band_x = 1700
    for i in range(8):
        on = (i == band_active)
        rect(slide, band_x + i*16, 1030, 12, 3,
             fill=(MINT if on else NAVY_MID))

    return slide


# ============================================================
# SLIDE 04 · EXECUTIVE SUMMARY (one-pager)
# ============================================================
def slide_exec_summary(prs):
    slide = add_blank_slide(prs)
    slide_chrome(slide, 'Executive Summary', '04', band_active=0,
                 source='phase2_report_20260419_152643.json · lopo_pillar_contribution_D20.json · threshold_sweep_xgb_D20.json')

    text(slide, 96, 116, 700, 22, '— EXECUTIVE SUMMARY · 一頁摘要',
         size=14, bold=True, color=TEAL, mono=True, ls=2.4)

    # title with hl
    rich(slide, 96, 144, 1750, 110, [
        {'text': '本系統以 ', 'size': 38, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': 'AUC 0.6490 · ICIR 0.7431 · DSR 12.12', 'size': 38, 'bold': True, 'color': TEAL, 'ls': -0.8, 'mono': True},
        {'text': ' 三重驗證', 'size': 38, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'br': True},
        {'text': 'Top 0.1% 選股 edge ', 'size': 38, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '+19.5 pp', 'size': 38, 'bold': True, 'color': TEAL, 'ls': -0.8, 'mono': True},
        {'text': '，9/9 治理閘門全數通過', 'size': 38, 'bold': True, 'color': NAVY, 'ls': -0.8},
    ], line_h=1.18)

    rect(slide, 96, 270, 80, 4, fill=TEAL)

    # 4 stat-hero cells
    box_y, box_h = 308, 168
    cell_w = 432
    gap = 16
    items = [
        ('OOS AUC (D+20)', '0.6490', 'XGBoost · Walk-Forward Purged K-Fold · 948,976 樣本', NAVY),
        ('ICIR · DSR', '0.7431', 'Mean IC +0.0145 · σ 0.0195 · DSR 12.12 (閾 2.0)', TEAL),
        ('Features Retained', '91/1,623', 'Corr → MI → VIF 三層漏斗 · 保留率 5.6%', NAVY),
        ('Governance Gates', '9/9', 'DSR · PSI · KS · Worst Fold · Embargo · Leak · Lineage · Repro · Sign-off', EMERALD),
    ]
    for i, (lbl, val, sub, color) in enumerate(items):
        cx = 96 + i*(cell_w + gap)
        rect(slide, cx, box_y, cell_w, box_h, fill=PAPER, line=RULE, line_w=0.75)
        stripe(slide, cx, box_y, cell_w, 4, color)
        text(slide, cx+24, box_y+22, cell_w-48, 22, lbl,
             size=11, bold=True, color=INK_3, mono=True, ls=2.4)
        # value (auto-shrink for long values)
        val_size = 56 if len(val) <= 6 else 44
        text(slide, cx+24, box_y+50, cell_w-48, 76, val,
             size=val_size, bold=True, color=color, mono=True, ls=-1.4)
        text(slide, cx+24, box_y+128, cell_w-48, 36, sub,
             size=11, color=INK_3, line_h=1.5)

    # ---- col-2 findings & recommendations ----
    col_y, col_h = 502, 290
    col_w = 880
    # left = findings
    text(slide, 96, col_y, col_w, 24, '三項核心發現 · KEY FINDINGS',
         size=14, bold=True, color=NAVY, ls=2.0)
    findings = [
        ('① Risk 支柱為預測引擎地基 ·',
         'LOPO ΔAUC Risk +138.6 bps、Trend +64.9、Val +9.5、Txt +8.5 排名前四；Chip / Event / Ind 為負貢獻'),
        ('② 門檻越嚴、edge 越陡 ·',
         'Top 1% → 0.5% → 0.1% 命中率由 36.82% → 39.64% → 45.79%（基準 26.28%、Top 0.1% 提升 +74.3%）'),
        ('③ 文本情緒具增量訊號 ·',
         'Text 支柱 LOPO +8.5 bps · News × Forum 共識相關 0.08 互補非重疊 · 500 keyword Lift × χ² 雙軸'),
    ]
    for i, (head, body) in enumerate(findings):
        fy = col_y + 30 + i*82
        rich(slide, 96, fy, col_w, 76, [
            {'text': head, 'size': 14, 'bold': True, 'color': NAVY},
            {'text': ' ' + body, 'size': 14, 'color': INK_2},
        ], line_h=1.5)
    # right = recommendations
    text(slide, 1000, col_y, col_w, 24, '三項落地建議 · RECOMMENDATIONS',
         size=14, bold=True, color=NAVY, ls=2.0)
    recs = [
        ('01 作為 alpha 過濾器 ·',
         'Top 1% 候選池每日 4,047 檔取高機率作為主動基金 rebalancing 清單，命中率 36.82% 較基準 +10.54 pp'),
        ('02 建立雙軌決策 ·',
         '保守門檻 0.50（召回 0.70%・命中 37.98%・edge +11.70 pp）／探索門檻 0.40（召回 8.83%・命中 29.41%）'),
        ('03 文本模組獨立治理 ·',
         '2025-Q1 txt 支柱 PSI 0.13 飄移警示 · 每月重訓詞向量 · 每季檢視 500 keyword Lift × χ² 清單'),
    ]
    for i, (head, body) in enumerate(recs):
        fy = col_y + 30 + i*82
        rich(slide, 1000, fy, col_w, 76, [
            {'text': head, 'size': 14, 'bold': True, 'color': ORANGE},
            {'text': ' ' + body, 'size': 14, 'color': INK_2},
        ], line_h=1.5)

    # ---- takeaway-bar ----
    tb_y = 820
    rect(slide, 96, tb_y, 1728, 132, fill=TEAL_BG)
    rect(slide, 96, tb_y, 6, 132, fill=TEAL)  # left border
    text(slide, 122, tb_y+20, 200, 18, 'BOTTOM LINE',
         size=11, bold=True, color=TEAL, mono=True, ls=2.6)
    rich(slide, 122, tb_y+44, 1660, 80, [
        {'text': '本系統通過 AUC 0.6490 / ICIR 0.7431 / DSR 12.12 三重統計驗證，9 道治理閘門全數通過，',
         'size': 16, 'color': INK},
        {'text': '具備可複製、可驗證、可落地', 'size': 16, 'bold': True, 'color': NAVY},
        {'text': ' 之決策品質；', 'size': 16, 'color': INK},
        {'text': 'Prob ≥ 0.50 候選池', 'size': 16, 'bold': True, 'color': ORANGE},
        {'text': ' 為最具商業價值輸出，適用每日 rebalancing 與風控 overlay。本系統為課程研究成果，不構成任何投資建議。',
         'size': 16, 'color': INK},
    ], line_h=1.5)

    return slide


# ============================================================
# SLIDE 05 · THREE COMMITMENTS
# ============================================================
def slide_three_commitments(prs):
    slide = add_blank_slide(prs)
    slide_chrome(slide, 'Three Commitments', '05', band_active=0,
                 source='Methodology reference: López de Prado (2018) · Bailey & López de Prado (2014) DSR')

    text(slide, 96, 116, 700, 22, '— SECTION 01 · 研究立場',
         size=14, bold=True, color=TEAL, mono=True, ls=2.4)

    rich(slide, 96, 144, 1750, 130, [
        {'text': '不做 ', 'size': 46, 'bold': True, 'color': NAVY, 'ls': -1.2},
        {'text': '黑箱、overfitting、無法重現 ', 'size': 46, 'bold': True, 'color': TEAL, 'ls': -1.2},
        {'text': '的研究', 'size': 46, 'bold': True, 'color': NAVY, 'ls': -1.2},
        {'br': True},
        {'text': '—— 三大學術承諾構成本專案的方法論底線', 'size': 24, 'color': INK_2},
    ], line_h=1.2)

    rect(slide, 96, 290, 80, 4, fill=TEAL)

    # 3 commitments cards
    cards = [
        ('01', '透明 · TRANSPARENCY',  NAVY,
         '不做黑箱模型',
         '提供 SHAP Top-20、LOPO ΔAUC、500 keyword Lift × χ²；儀表板「ICIR / Text / Feature Analysis」三頁專用於可解釋性展示。'),
        ('02', '嚴謹 · RIGOR',          ORANGE,
         '不做 overfitting',
         'Walk-Forward Purged K-Fold + 20 日 embargo · DSR 12.12 通過 150 trials 多重檢定 · PSI / KS 漂移監控 9 支柱 × 30 日。'),
        ('03', '可重現 · REPRODUCIBILITY', EMERALD,
         '不做無法重現的研究',
         '每階段 JSON artifact + rng=42 + lineage · Phase 1–6 六份報告 + 13 頁互動 dashboard，可逐頁追溯模型決策。'),
    ]
    card_y = 340
    card_w = 552
    card_h = 412
    gap = 24
    for i, (no, lbl, color, title, body) in enumerate(cards):
        cx = 96 + i*(card_w + gap)
        rect(slide, cx, card_y, card_w, card_h, fill=PAPER, line=RULE, line_w=0.75)
        # top stripe
        stripe(slide, cx, card_y, card_w, 6, color)
        # large mono number
        text(slide, cx+36, card_y+34, 200, 80, no,
             size=64, bold=True, color=color, mono=True, ls=-2.0)
        # label
        text(slide, cx+36, card_y+128, card_w-72, 22, lbl,
             size=12, bold=True, color=color, mono=True, ls=2.4)
        # title
        text(slide, cx+36, card_y+160, card_w-72, 60, title,
             size=28, bold=True, color=NAVY, line_h=1.18)
        # divider
        rect(slide, cx+36, card_y+232, 60, 3, fill=color)
        # body
        text(slide, cx+36, card_y+260, card_w-72, 130, body,
             size=15, color=INK_2, line_h=1.6)

    # takeaway-bar
    tb_y = 776
    rect(slide, 96, tb_y, 1728, 168, fill=TEAL_BG)
    rect(slide, 96, tb_y, 6, 168, fill=TEAL)
    text(slide, 122, tb_y+22, 480, 20, 'ACADEMIC BASELINE · 學術底線',
         size=11, bold=True, color=TEAL, mono=True, ls=2.6)
    rich(slide, 122, tb_y+52, 1660, 100, [
        {'text': '三承諾為 Phase 1–6 設計約束；任一未守，則統計結論失效。', 'size': 17, 'color': INK},
        {'br': True},
        {'text': '這是與一般「回測看起來很好」的量化原型', 'size': 17, 'color': INK},
        {'text': '最根本的區隔', 'size': 17, 'bold': True, 'color': NAVY},
        {'text': '——所有產出可獨立稽核、可第三方重訓。', 'size': 17, 'color': INK},
    ], line_h=1.55)

    return slide


# ============================================================
# SLIDE 06 · SECTION 02 DIVIDER
# ============================================================
# (use slide_section_divider helper)


# ============================================================
# SLIDE 07 · PROBLEM STATEMENT
# ============================================================
def slide_problem(prs):
    slide = add_blank_slide(prs)
    slide_chrome(slide, 'Problem Statement', '07', band_active=1,
                 source='Methodology: López de Prado (2018) · 台股文本驗證於 phase1_report_20260419_112854.json §data_integrity')

    text(slide, 96, 116, 700, 22, '— PROBLEM · 市場痛點',
         size=14, bold=True, color=TEAL, mono=True, ls=2.4)

    rich(slide, 96, 144, 1750, 180, [
        {'text': '台股短線預測面臨 ', 'size': 44, 'bold': True, 'color': NAVY, 'ls': -1.2},
        {'text': '三重結構性難題', 'size': 44, 'bold': True, 'color': TEAL, 'ls': -1.2},
        {'br': True},
        {'text': '傳統單一因子模型無法穿透', 'size': 44, 'bold': True, 'color': NAVY, 'ls': -1.2},
    ], line_h=1.18)

    rect(slide, 96, 312, 80, 4, fill=TEAL)

    text(slide, 96, 336, 1750, 30, '主動基金與量化策略瓶頸不在資料取得，而在於「如何在嚴格時序下整合異質訊號並防範洩漏」。',
         size=17, color=INK_2)

    # 3 trap cards
    traps = [
        ('01', '訊號雜訊比極低', NAVY,
         '日頻波動真正可預測 edge 1–5 pp',
         'Walk-Forward Purged K-Fold + 20 日 embargo 阻斷 forward-looking bias；對應 DSR 12.12 通過 150 trials 多重檢定。'),
        ('02', '異質訊號難對齊', ORANGE,
         '價量 tick / 財報季頻 / PTT 事件頻',
         '以 T 日收盤可取得為邊界、財報 lag 45 天、文本僅取 T 日 15:00 前；Feature Store (stock_id, date) 主鍵統一對齊。'),
        ('03', '文本資料品質混雜', ROSE,
         '1,125,134 筆文本中 ~40% 集中於 Top 5',
         '自建台股 500 keyword Lift × χ² 詞表 + 覆蓋率衰減校正 + tickers entity linking 92% 涵蓋率。'),
    ]
    card_y = 384
    card_w = 552
    card_h = 396
    gap = 24
    for i, (no, lbl, color, headline, body) in enumerate(traps):
        cx = 96 + i*(card_w + gap)
        rect(slide, cx, card_y, card_w, card_h, fill=PAPER, line=RULE, line_w=0.75)
        stripe(slide, cx, card_y, card_w, 6, color)
        # large mono number on right top
        text(slide, cx+card_w-160, card_y+30, 130, 60, no,
             size=48, bold=True, color=color, align='right', mono=True, ls=-1.6)
        # eyebrow
        text(slide, cx+36, card_y+34, card_w-72, 24, 'STRUCTURAL TRAP',
             size=11, bold=True, color=INK_4, mono=True, ls=2.4)
        # title
        text(slide, cx+36, card_y+72, card_w-180, 60, lbl,
             size=26, bold=True, color=NAVY, line_h=1.15)
        # divider
        rect(slide, cx+36, card_y+136, 60, 3, fill=color)
        # headline (mono)
        text(slide, cx+36, card_y+160, card_w-72, 26, headline,
             size=16, bold=True, color=color, mono=True, line_h=1.4)
        # body
        text(slide, cx+36, card_y+212, card_w-72, 160, body,
             size=15, color=INK_2, line_h=1.65)

    # takeaway-bar
    tb_y = 792
    rect(slide, 96, tb_y, 1728, 152, fill=TEAL_BG)
    rect(slide, 96, tb_y, 6, 152, fill=TEAL)
    text(slide, 122, tb_y+22, 200, 20, 'SO WHAT',
         size=11, bold=True, color=TEAL, mono=True, ls=2.6)
    rich(slide, 122, tb_y+52, 1660, 90, [
        {'text': '三難題構成「短期預測不可重現」的陷阱；本系統以 ', 'size': 17, 'color': INK},
        {'text': '嚴格時序治理 + 多源對齊', 'size': 17, 'bold': True, 'color': NAVY},
        {'text': ' 證明 ', 'size': 17, 'color': INK},
        {'text': 'AUC 0.6490 / ICIR 0.7431 / DSR 12.12', 'size': 17, 'bold': True, 'color': TEAL, 'mono': True},
        {'text': ' 可在未來週期穩定複現，而非回測過擬合的假象。', 'size': 17, 'color': INK},
    ], line_h=1.55)

    return slide


# ============================================================
# SLIDE 08 · RESEARCH APPROACH (6 PHASES)
# ============================================================
def slide_approach(prs):
    slide = add_blank_slide(prs)
    slide_chrome(slide, 'Research Approach', '08', band_active=1,
                 source='Pipeline doc: 專案說明文件 §3 · artifacts in outputs/reports/')

    text(slide, 96, 116, 700, 22, '— APPROACH · 研究方法',
         size=14, bold=True, color=TEAL, mono=True, ls=2.4)

    rich(slide, 96, 144, 1750, 180, [
        {'text': '以 ', 'size': 44, 'bold': True, 'color': NAVY, 'ls': -1.2},
        {'text': '6 階段流水線 ', 'size': 44, 'bold': True, 'color': TEAL, 'ls': -1.2},
        {'text': '從原始資料走到決策清單', 'size': 44, 'bold': True, 'color': NAVY, 'ls': -1.2},
        {'br': True},
        {'text': '每階段皆有可稽核產出', 'size': 44, 'bold': True, 'color': NAVY, 'ls': -1.2},
    ], line_h=1.18)

    rect(slide, 96, 312, 80, 4, fill=TEAL)

    # 6 phases in 3x2 grid
    phases = [
        ('01', 'Ingest',    '資料清洗與對齊',   '7 表 · 3,483,598 列 ETL · 除權息補償、RSI 異常',           'phase1_report.json',                          NAVY),
        ('02', 'Feature',   '特徵工程與選擇',   '1,623 → 91 · Purged K-Fold · 20 日 embargo · 948,976 樣本', 'phase2_report.json · feature_store.parquet',   NAVY),
        ('03', 'Model',     '模型訓練與調參',   'XGBoost / LightGBM / Logistic / RF · Optuna 150 trials',     'model_registry/*.pkl',                         NAVY),
        ('04', 'Explain',   '歸因分析 · LOPO',  'Leave-One-Pillar-Out · SHAP summary · ΔAUC × 10,000 bps',  'lopo_pillar_contribution_D20.json',           ORANGE),
        ('05', 'Backtest',  '門檻掃描 · 回測',  '門檻 0.40 / 0.50 / 0.60 命中 58.7 / 67.8 / 72.3% · 三情境', 'threshold_sweep_xgb_D20.json',                 ORANGE),
        ('06', 'Govern',    '治理閘門 · 交付',  'DSR / PSI / KS / embargo / 最差 fold / lineage 9 項',        'phase6_gates.json · dashboard/',              ORANGE),
    ]
    grid_x, grid_y = 96, 350
    cell_w = 568
    cell_h = 220
    gap = 24
    for i, (no, name, zh, body, artifact, color) in enumerate(phases):
        col, row = i % 3, i // 3
        cx = grid_x + col * (cell_w + gap)
        cy = grid_y + row * (cell_h + gap)
        rect(slide, cx, cy, cell_w, cell_h, fill=PAPER, line=RULE, line_w=0.75)
        stripe(slide, cx, cy, cell_w, 4, color)
        # phase number + name (mono)
        text(slide, cx+24, cy+22, 80, 30, no,
             size=22, bold=True, color=color, mono=True, ls=-0.6)
        text(slide, cx+72, cy+24, 200, 26, name,
             size=14, bold=True, color=color, mono=True, ls=2.0)
        # zh title
        text(slide, cx+24, cy+62, cell_w-48, 32, zh,
             size=21, bold=True, color=NAVY)
        # body
        text(slide, cx+24, cy+106, cell_w-48, 76, body,
             size=14, color=INK_2, line_h=1.55)
        # artifact (footer)
        rect(slide, cx+24, cy+cell_h-44, cell_w-48, 1, fill=RULE)
        text(slide, cx+24, cy+cell_h-32, cell_w-48, 22, '↪ ' + artifact,
             size=11, color=INK_3, mono=True, ls=0.4)

    # takeaway-bar
    tb_y = 824
    rect(slide, 96, tb_y, 1728, 124, fill=TEAL_BG)
    rect(slide, 96, tb_y, 6, 124, fill=TEAL)
    text(slide, 122, tb_y+22, 300, 20, 'DESIGN PRINCIPLE · 設計原則',
         size=11, bold=True, color=TEAL, mono=True, ls=2.6)
    rich(slide, 122, tb_y+52, 1660, 70, [
        {'text': '每階段皆產出 JSON 報告 + 摘要 + artifact，dashboard 可逐頁追溯；', 'size': 17, 'color': INK},
        {'text': '所有 random seed = 42、所有切分 lineage 可重訓，', 'size': 17, 'color': INK},
        {'text': '從 raw → 決策清單', 'size': 17, 'bold': True, 'color': NAVY},
        {'text': ' 任一節點皆可被獨立復現。', 'size': 17, 'color': INK},
    ], line_h=1.55)

    return slide


# ============================================================
# MAIN
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = W_PX * PX
    prs.slide_height = H_PX * PX

    print('Building Slide 01 · Cover ...')
    slide_cover(prs)
    print('Building Slide 02 · Table of Contents ...')
    slide_toc(prs)
    print('Building Slide 03 · Section 01 Divider ...')
    slide_section_divider(prs,
        no_str='SECTION 01 / 08',
        line1='Executive', line2='Summary',
        sub='從統計驗證到落地建議的一頁摘要：AUC 0.6490 三重指標、Top 0.1% edge +19.5 pp、9/9 治理閘門。',
        page_no='03', band_active=0)
    print('Building Slide 04 · Executive Summary ...')
    slide_exec_summary(prs)
    print('Building Slide 05 · Three Commitments ...')
    slide_three_commitments(prs)
    print('Building Slide 06 · Section 02 Divider ...')
    slide_section_divider(prs,
        no_str='SECTION 02 / 08',
        line1='Problem', line2='& Approach',
        sub='台股短線預測為何無法穿透？三重結構性難題與我們的 6 階段流水線設計。',
        page_no='06', band_active=1)
    print('Building Slide 07 · Problem Statement ...')
    slide_problem(prs)
    print('Building Slide 08 · Research Approach ...')
    slide_approach(prs)

    out_path = r'C:\Users\user\Desktop\大數據與商業分析專案\簡報設計\preview\preview_v1.pptx'
    prs.save(out_path)
    print(f'\nSaved: {out_path}')
    print(f'Slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
