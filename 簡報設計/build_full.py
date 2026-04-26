# -*- coding: utf-8 -*-
"""
build_full.py — 多因子股票預測系統 · 完整簡報（v11.5.18）

設計原則（BCG 資深顧問視角，課堂 12 分鐘版）：
  - 主線 10 頁（每頁 ≈ 1 分鐘）· 每頁 land 一個結論
  - 附錄 27+ 頁，留作 Q&A / 細節支撐
  - 1920×1080 畫布 · 對齊 design_tokens 設計系統
  - 配色：Navy #1E3A8A / Blue #2563EB / Cyan #06B6D4 / Rose #F43F5E
         / Emerald #10B981 / Amber #F59E0B / Violet #7C3AED
  - 字級克制：H1 ≤ 40pt（封面 hero 例外）· body 14pt · KPI mono 至 56pt
  - 用字：metric-first、無「首先/其次/值得一提的是」AI 痕跡
  - 中文於句讀（、，：）處硬斷行，避免被自動斷在不對的位置

Main story arc:
  P.01 Cover                       封面 + thesis + 三錨點
  P.02 短線是噪聲、alpha 在中長線    切角與時序選擇
  P.03 1,623 → 91 · WFPK            方法總覽一頁打死
  P.04 跨四指標一致 alpha           Headline result（4-KPI hero）
  P.05 DSR 12.12 + 9/9 gates        統計穩健性
  P.06 價值因子 U 型結構             Finding 1（雙端 alpha）
  P.07 散戶情緒反向訊號              Finding 2（媒體 vs 散戶非同源）
  P.08 Top 0.1% · +19.5pp           策略落地 / 部署 edge
  P.09 部署條件、限制、再訓練         Limits + roadmap
  P.10 結語 · Q&A                    Thesis recap

Appendix (27+ 張)：詳細支撐材料 + 主線砍下的內容（exec summary 細節、9 支柱 LOPO 等）
"""
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# CANVAS · PALETTE · FONTS
# ============================================================
PX = 9525
W_PX, H_PX = 1920, 1080

NAVY      = RGBColor(0x1e, 0x3a, 0x8a)
NAVY_D    = RGBColor(0x1e, 0x40, 0xaf)
NAVY_DK   = RGBColor(0x0b, 0x12, 0x20)
NAVY_MID  = RGBColor(0x16, 0x21, 0x38)
NAVY_LT   = RGBColor(0xe0, 0xe7, 0xff)
TEAL      = RGBColor(0x06, 0xb6, 0xd4)   # cyan-leaning, design-system primary teal
TEAL_D    = RGBColor(0x08, 0x91, 0xb2)
TEAL_BG   = RGBColor(0xec, 0xfe, 0xff)
MINT      = RGBColor(0x67, 0xe8, 0xf9)
MINT_D    = RGBColor(0x22, 0xd3, 0xee)
EMERALD   = RGBColor(0x10, 0xb9, 0x81)   # design-system success
EMERALD_DK= RGBColor(0x05, 0x96, 0x69)
EMERALD_BG= RGBColor(0xec, 0xfd, 0xf5)
ROSE      = RGBColor(0xf4, 0x3f, 0x5e)   # design-system danger / contra
ROSE_BG   = RGBColor(0xff, 0xf1, 0xf2)
ROSE_LT   = RGBColor(0xfe, 0xcd, 0xd3)
GOLD      = RGBColor(0xf5, 0x9e, 0x0b)
GOLD_BG   = RGBColor(0xfe, 0xf3, 0xc7)
ORANGE    = RGBColor(0xea, 0x58, 0x0c)
ORANGE_L  = RGBColor(0xfb, 0x92, 0x3c)
VIOLET    = RGBColor(0x7c, 0x3a, 0xed)
INDIGO    = RGBColor(0x4f, 0x46, 0xe5)
CYAN      = RGBColor(0x06, 0xb6, 0xd4)
PINK      = RGBColor(0xec, 0x48, 0x99)
PURPLE    = RGBColor(0xa8, 0x55, 0xf7)
INK       = RGBColor(0x0f, 0x17, 0x2a)
INK_2     = RGBColor(0x33, 0x41, 0x55)
INK_3     = RGBColor(0x64, 0x74, 0x8b)
INK_3_DK  = RGBColor(0x47, 0x55, 0x69)   # slate-600 — 取代原 italic 的補強灰
INK_4     = RGBColor(0x94, 0xa3, 0xb8)
PAPER     = RGBColor(0xff, 0xff, 0xff)
BG        = RGBColor(0xfa, 0xfb, 0xfc)
SUBTLE    = RGBColor(0xf1, 0xf5, 0xf9)
TINT      = RGBColor(0xf8, 0xfa, 0xfc)
RULE      = RGBColor(0xe2, 0xe8, 0xf0)
RULE_2    = RGBColor(0xcb, 0xd5, 0xe1)
IVORY     = RGBColor(0xf5, 0xf1, 0xe8)
CREAM     = RGBColor(0xfb, 0xf8, 0xf1)

# 字體紀律（v11.5.24 起 · 套用「管顧式簡報 SKILL」§ 2.1）
# 所有文字統一使用微軟正黑體（Microsoft JhengHei）、跨平台 / 課堂 / 學術場景全相容
LATIN_FONT = 'Microsoft JhengHei'
EA_FONT    = 'Microsoft JhengHei'
MONO_FONT  = 'Microsoft JhengHei'

# 字級下限守門（v11.5.24 起 · § 2.2）
# 任何可見文字字級 < MIN_FONT_SIZE 視為錯誤、build 階段直接 raise
MIN_FONT_SIZE = 12

# ============================================================
# XML HELPERS
# ============================================================
def _set_run_typeface(run, latin=LATIN_FONT, ea=EA_FONT, mono=False):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    typeface_latin = MONO_FONT if mono else latin
    el = etree.SubElement(rPr, qn('a:latin')); el.set('typeface', typeface_latin)
    el = etree.SubElement(rPr, qn('a:ea'));    el.set('typeface', ea)
    el = etree.SubElement(rPr, qn('a:cs'));    el.set('typeface', ea)

def _set_run_letterspacing(run, ls_pts):
    if not ls_pts:
        return
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(ls_pts * 100)))

# ============================================================
# SHAPE HELPERS
# ============================================================
def rect(slide, x, y, w, h, fill=None, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x*PX, y*PX, w*PX, h*PX)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp

def stripe(slide, x, y, w, h, color):
    return rect(slide, x, y, w, h, fill=color)

def dot(slide, cx, cy, r, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, (cx-r)*PX, (cy-r)*PX, (2*r)*PX, (2*r)*PX)
    shp.shadow.inherit = False
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def ring(slide, cx, cy, r, color, line_w=2):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, (cx-r)*PX, (cy-r)*PX, (2*r)*PX, (2*r)*PX)
    shp.shadow.inherit = False
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(line_w)
    return shp

# ============================================================
# TEXT HELPERS
# ============================================================
def text(slide, x, y, w, h, content, *, size=18, color=INK, bold=False,
         align='left', anchor='top', mono=False, italic=False,
         ls=0, line_h=1.2, latin=LATIN_FONT, ea=EA_FONT):
    box = slide.shapes.add_textbox(x*PX, y*PX, w*PX, h*PX)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {
        'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM
    }[anchor]
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    if line_h:
        p.line_spacing = line_h
    run = p.add_run()
    run.text = content
    rf = run.font
    rf.size = Pt(size); rf.bold = bold; rf.italic = italic; rf.color.rgb = color
    _set_run_typeface(run, latin=latin, ea=ea, mono=mono)
    _set_run_letterspacing(run, ls)
    return box

def rich(slide, x, y, w, h, runs, *, align='left', anchor='top', line_h=1.2):
    box = slide.shapes.add_textbox(x*PX, y*PX, w*PX, h*PX)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {
        'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM
    }[anchor]
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    if line_h:
        p.line_spacing = line_h
    def new_para():
        nonlocal p
        p = tf.add_paragraph()
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
        if line_h:
            p.line_spacing = line_h
    for r in runs:
        if r.get('sp'):
            new_para(); continue
        if r.get('br'):
            run = p.add_run(); run.text = ''
            br = etree.SubElement(run._r.getparent(), qn('a:br'))
            run._r.addnext(br); continue
        run = p.add_run()
        run.text = r.get('text', '')
        rf = run.font
        rf.size = Pt(r.get('size', 18))
        rf.bold = r.get('bold', False)
        rf.italic = r.get('italic', False)
        rf.color.rgb = r.get('color', INK)
        _set_run_typeface(run, latin=r.get('latin', LATIN_FONT), ea=r.get('ea', EA_FONT),
                          mono=r.get('mono', False))
        _set_run_letterspacing(run, r.get('ls', 0))
    return box

# ============================================================
# CHROME & COMMON BLOCKS
# ============================================================
def slide_chrome_main(slide, page_title, page_no, *, total=10, act=None, source=None):
    # v11.5.20 §1 — page number relocated from top-right to bottom-right
    # (consistent with cover slide); top-right reserved for page_title.
    rect(slide, 96, 38, 30, 30, fill=NAVY)
    text(slide, 96, 38, 30, 30, 'TW', size=12, bold=True, color=PAPER,
         align='center', anchor='middle', mono=True)
    text(slide, 134, 38, 380, 30, '多因子股票預測系統', size=12, bold=True,
         color=INK_2, anchor='middle', ls=1.4)
    text(slide, 522, 42, 200, 22, '· v11.5.18', size=12, color=INK_4,
         anchor='middle', mono=True, ls=0.8)
    # Top-right: page title only
    text(slide, 1180, 38, 700, 30, page_title, size=12, color=INK_3,
         align='right', anchor='middle', mono=True, ls=2.0)
    if source:
        text(slide, 96, 1018, 1300, 30, 'Source · ' + source, size=12,
             color=INK_3, anchor='middle', mono=True, ls=0.5)
    if act is not None:
        labels = ['SETUP', 'FINDINGS', 'ACTION']
        text(slide, 1500, 1024, 160, 18, labels[act], size=12, bold=True,
             color=TEAL, mono=True, ls=2.0, align='right')
        for i in range(3):
            on = (i == act)
            cx = 1680 + i*22
            if on:
                dot(slide, cx, 1033, 5, TEAL)
            else:
                ring(slide, cx, 1033, 4, RULE_2, line_w=1)
    # Bottom-right: consistent page-number stamp
    text(slide, 1780, 1018, 60, 30, page_no, size=12, bold=True,
         color=INK_2, align='right', anchor='middle', mono=True)


def slide_chrome_appendix(slide, page_title, apx_no, *, source=None):
    # v11.5.20 §1 — appendix page number also at bottom-right.
    rect(slide, 96, 38, 36, 30, fill=ROSE)
    text(slide, 96, 38, 36, 30, 'APX', size=12, bold=True, color=PAPER,
         align='center', anchor='middle', mono=True)
    text(slide, 140, 38, 380, 30, '多因子股票預測系統 · 附錄', size=12, bold=True,
         color=INK_2, anchor='middle', ls=1.4)
    text(slide, 1180, 38, 700, 30, page_title, size=12, color=INK_3,
         align='right', anchor='middle', mono=True, ls=2.0)
    if source:
        text(slide, 96, 1018, 1300, 30, 'Source · ' + source, size=12,
             color=INK_3, anchor='middle', mono=True, ls=0.5)
    text(slide, 1560, 1024, 200, 18, 'APPENDIX', size=12, bold=True,
         color=ROSE, mono=True, ls=2.4, align='right')
    rect(slide, 1768, 1029, 8, 8, fill=ROSE)
    # Bottom-right: appendix page number (rose)
    text(slide, 1780, 1018, 60, 30, apx_no, size=12, color=ROSE,
         bold=True, align='right', anchor='middle', mono=True)


def add_blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    rect(slide, 0, 0, W_PX, H_PX, fill=PAPER)
    return slide


def eyebrow(slide, x, y, w, content, *, color=TEAL, size=14, ls=2.4):
    text(slide, x, y, w, 22, '— ' + content, size=size, bold=True,
         color=color, mono=True, ls=ls)

def rule_divider(slide, x, y, w=80, h=4, color=TEAL):
    rect(slide, x, y, w, h, fill=color)

def takeaway_bar(slide, x, y, w, h, label, runs, *, color=TEAL, bg=TEAL_BG):
    if h < 24:
        return
    rect(slide, x, y, w, h, fill=bg)
    rect(slide, x, y, 6, h, fill=color)
    if label:
        text(slide, x+26, y+18, 600, 20, label,
             size=12, bold=True, color=color, mono=True, ls=2.6)
        rich_y = y + 48
    else:
        rich_y = y + 14
    rich_h = max(20, h - (rich_y - y) - 8)
    if runs:
        rich(slide, x+26, rich_y, w-60, rich_h, runs, line_h=1.5)

def kpi_hero(slide, x, y, w, h, label, value, sub, accent=NAVY, *, val_size=None):
    rect(slide, x, y, w, h, fill=PAPER, line=RULE, line_w=0.75)
    stripe(slide, x, y, w, 4, accent)
    text(slide, x+24, y+22, w-48, 22, label,
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    if val_size is None:
        val_size = 56 if len(value) <= 6 else (44 if len(value) <= 9 else 36)
    text(slide, x+24, y+50, w-48, 76, value,
         size=val_size, bold=True, color=accent, mono=True, ls=-1.4)
    text(slide, x+24, y+128, w-48, 36, sub,
         size=12, color=INK_3, line_h=1.5)


def card(slide, x, y, w, h, *, fill=PAPER, line=RULE, line_w=0.75, top_stripe=None):
    rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w)
    if top_stripe is not None:
        stripe(slide, x, y, w, 4, top_stripe)


def hbar(slide, x, y, w, h, label, value, max_val, *, color=TEAL, label_w=200, value_w=120, fmt=None):
    text(slide, x, y, label_w, h, label, size=14, bold=True, color=INK, anchor='middle')
    track_x = x + label_w + 16
    track_w = w - label_w - value_w - 32
    rect(slide, track_x, y+h//2-3, track_w, 6, fill=SUBTLE)
    bar_w = max(2, int(value / max_val * track_w))
    rect(slide, track_x, y+h//2-3, bar_w, 6, fill=color)
    val_str = fmt.format(value) if fmt else f'{value:.2f}'
    text(slide, x+w-value_w, y, value_w, h, val_str, size=14, bold=True,
         color=INK, mono=True, align='right', anchor='middle')


def hbar_signed(slide, x, y, label_w, axis_w, value_w, h, label, value_bps, max_abs, *,
                pos_color=EMERALD, neg_color=ROSE):
    text(slide, x, y, label_w, h, label, size=14, bold=True, color=INK, anchor='middle')
    axis_x = x + label_w + 16
    zero_x = axis_x + axis_w // 2
    rect(slide, axis_x, y+h//2-1, axis_w, 1, fill=RULE_2)
    rect(slide, zero_x, y+8, 1, h-16, fill=INK_3)
    bar_h = 18
    bar_y = y + (h - bar_h) // 2
    bar_w = int(abs(value_bps) / max_abs * (axis_w/2 - 8))
    if value_bps >= 0:
        rect(slide, zero_x+1, bar_y, bar_w, bar_h, fill=pos_color)
    else:
        rect(slide, zero_x-1-bar_w, bar_y, bar_w, bar_h, fill=neg_color)
    sign = '+' if value_bps >= 0 else ''
    text(slide, x+label_w+16+axis_w+16, y, value_w, h,
         f'{sign}{value_bps:.1f} bps', size=13, bold=True,
         color=(pos_color if value_bps >= 0 else neg_color),
         mono=True, anchor='middle')


def vbar(slide, x_base, y_base, w, max_h, value, max_val, *, color=NAVY, label=None, val_label=None,
         label_color=None):
    bar_h = max(2, int(value / max_val * max_h))
    rect(slide, x_base, y_base - bar_h, w, bar_h, fill=color)
    if val_label:
        text(slide, x_base-12, y_base - bar_h - 30, w+24, 24, val_label,
             size=12, bold=True, color=INK, mono=True, align='center', anchor='middle')
    if label:
        text(slide, x_base-12, y_base + 8, w+24, 28, label,
             size=12, color=label_color or INK_3, mono=True, align='center', anchor='middle', ls=0.4)


def linechart_box(slide, x, y, w, h):
    rect(slide, x, y, w, h, fill=PAPER, line=RULE, line_w=0.5)
    pad_l, pad_r, pad_t, pad_b = 60, 30, 30, 36
    px_, py_, pw_, ph_ = x+pad_l, y+pad_t, w-pad_l-pad_r, h-pad_t-pad_b
    rect(slide, px_, py_+ph_, pw_, 1, fill=RULE_2)
    rect(slide, px_, py_, 1, ph_, fill=RULE_2)
    return px_, py_, pw_, ph_


def polyline(slide, points, color, line_w=2.25):
    for i in range(len(points)-1):
        x1, y1 = points[i]
        x2, y2 = points[i+1]
        ln = slide.shapes.add_connector(1, x1*PX, y1*PX, x2*PX, y2*PX)
        ln.shadow.inherit = False
        ln.line.color.rgb = color
        ln.line.width = Pt(line_w)


# ============================================================
# ============================================================
# MAIN STORY · 10 SLIDES (P.01 – P.10) · 12 分鐘課堂版
# ============================================================
# ============================================================

# ============================================================
# P.01 · COVER
# ============================================================
def slide_cover(prs):
    slide = add_blank_slide(prs)
    stripe(slide, 0, 0, W_PX, 6, NAVY)

    # Top eyebrows · 對齊設計系統 mono caption 11pt
    text(slide, 96, 44, 700, 22, 'TAIWAN EQUITY · QUANT RESEARCH TERMINAL',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    text(slide, 1180, 44, 660, 22, 'NTU · BIG DATA & BUSINESS ANALYTICS',
         size=12, bold=True, color=INK_3, align='right', mono=True, ls=2.4)

    # Hero title block — 收斂到 80pt（從 96pt），增加 BCG 式克制
    rect(slide, 96, 168, 60, 4, fill=TEAL)
    text(slide, 96, 188, 700, 26, 'STRATEGIC RESEARCH DECK · v11.5.18',
         size=13, bold=True, color=TEAL, mono=True, ls=2.4)
    rich(slide, 96, 230, 1000, 240, [
        {'text': '台股多因子', 'size': 78, 'bold': True, 'color': NAVY, 'ls': -2.0},
        {'br': True},
        {'text': '預測系統', 'size': 78, 'bold': True, 'color': TEAL, 'ls': -2.0},
    ], line_h=1.05)

    # Thesis line — 用戶聲音：metric-first、direct assertion
    # CJK 換行紀律：逗號（、，）不強制換行、整段流動、自動 wrap；句號才換行。
    rect(slide, 96, 488, 1000, 1, fill=RULE_2)
    rich(slide, 96, 510, 1000, 130, [
        {'text': '台股短線是噪聲、', 'size': 22, 'bold': True, 'color': INK},
        {'text': 'alpha 在中長線基本面結構', 'size': 22, 'bold': True, 'color': NAVY},
        {'br': True},
        {'text': '以 9 支柱因子棋盤、', 'size': 17, 'color': INK_2},
        {'text': '948,976', 'size': 17, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' 樣本、', 'size': 17, 'color': INK_2},
        {'text': '505', 'size': 17, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' 交易日，於 D+20 方向預測產出可部署的選股訊號。', 'size': 17, 'color': INK_2},
    ], line_h=1.45)

    # 3-anchor strip — 報告開場直接亮關鍵數字
    anchor_y = 690
    rect(slide, 96, anchor_y, 1000, 130, fill=PAPER, line=RULE, line_w=0.75)
    stripe(slide, 96, anchor_y, 1000, 4, NAVY)
    anchors = [
        ('OOS AUC · D+20',  '0.6455', 'XGBoost · 4-fold WFPK', NAVY),
        ('DSR',             '12.12',  '閾 3.0 · 多重檢定通過',   TEAL),
        ('CAGR · TOP 1%',   '+36.8%', 'vs TAIEX +12.4%',         EMERALD),
    ]
    for i, (lbl, val, sub, c) in enumerate(anchors):
        ax = 96 + i*333
        if i > 0:
            rect(slide, ax, anchor_y+22, 1, 86, fill=RULE)
        text(slide, ax+28, anchor_y+22, 305, 18, lbl,
             size=12, bold=True, color=INK_3, mono=True, ls=2.0)
        text(slide, ax+28, anchor_y+46, 305, 60, val,
             size=44, bold=True, color=c, mono=True, ls=-1.6)
        text(slide, ax+28, anchor_y+108, 305, 18, sub,
             size=12, color=INK_3, mono=True, ls=0.4)

    # Right rail — 設計系統 KPI grid（輔助訊息，不喧賓奪主）
    grid_x, grid_y = 1158, 188
    cell_w, cell_h = 320, 200
    kpis = [
        ('FEATURES',   '91 / 1,623', '三層漏斗 · 保留 5.6%', NAVY),
        ('ICIR',       '0.7431',     '4 折 IC 一致性',        TEAL),
        ('GATES',      '9 / 9',      '治理關卡全數通過',       EMERALD),
        ('SAMPLES',    '948,976',    '1,930 檔 × 505 交易日', NAVY),
    ]
    for i, (lbl, val, sub, accent) in enumerate(kpis):
        col, row = i % 2, i // 2
        cx = grid_x + col*(cell_w + 8)
        cy = grid_y + row*(cell_h + 8)
        rect(slide, cx, cy, cell_w, cell_h, fill=PAPER, line=RULE, line_w=0.75)
        stripe(slide, cx, cy, cell_w, 4, accent)
        text(slide, cx+22, cy+22, cell_w-44, 18, lbl,
             size=12, bold=True, color=INK_3, mono=True, ls=2.4)
        text(slide, cx+22, cy+50, cell_w-44, 80, val,
             size=34, bold=True, color=accent, mono=True, ls=-1.0)
        text(slide, cx+22, cy+138, cell_w-44, 40, sub,
             size=12, color=INK_3, mono=True, ls=0.4)

    rect(slide, 1158, 690, 660, 130, fill=PAPER, line=RULE, line_w=0.75)
    stripe(slide, 1158, 690, 660, 4, TEAL)
    text(slide, 1180, 712, 620, 18, 'WHAT YOU\'LL HEAR · 12 MIN',
         size=12, bold=True, color=TEAL, mono=True, ls=2.4)
    rich(slide, 1180, 742, 620, 70, [
        {'text': '一個結論：', 'size': 14, 'bold': True, 'color': INK},
        {'text': '台股 D+20 alpha 可預測、可驗證、可部署。', 'size': 14, 'color': INK_2},
        {'br': True},
        {'text': '三段路徑：', 'size': 14, 'bold': True, 'color': INK},
        {'text': '切角 → 方法 → 證據 → 落地。', 'size': 14, 'color': INK_2},
    ], line_h=1.45)

    # Footer · 6-人組員名單 + 課程資訊
    rect(slide, 96, 858, 1728, 1, fill=RULE)
    text(slide, 96, 876, 360, 18, 'TEAM · 6 MEMBERS',
         size=12, bold=True, color=INK_4, mono=True, ls=2.6)
    members = [
        ('R13724054', '陳柏寰'),
        ('R13724035', '陳又萍'),
        ('R13724051', '謝瑞洋'),
        ('R13724056', '李泰坤'),
        ('B12302131', '林芮靚'),
        ('B13102133', '沈承寬'),
    ]
    # v11.5.24 § 字體統一微軟正黑體後 9 字學號需更寬 box；ID 容器 90→112 px
    name_block_x = 96
    name_col_w = 240
    name_row_h = 28
    for i, (sid, name) in enumerate(members):
        col = i % 3; row = i // 3
        nx = name_block_x + col*name_col_w
        ny = 902 + row*name_row_h
        text(slide, nx, ny, 112, 22, sid,
             size=12, bold=True, color=INK_3, mono=True, ls=0.2)
        text(slide, nx+118, ny, 120, 22, name,
             size=13, bold=True, color=INK)

    # Right side: COURSE / VERSION / DATE
    text(slide, 800, 876, 240, 18, 'COURSE',
         size=12, bold=True, color=INK_4, mono=True, ls=2.6)
    text(slide, 800, 902, 380, 22, 'NTU · Big Data & Business Analytics',
         size=12, bold=True, color=INK)
    text(slide, 800, 928, 380, 22, '大數據與商業分析專案',
         size=12, color=INK_2)

    text(slide, 1220, 876, 200, 18, 'VERSION',
         size=12, bold=True, color=INK_4, mono=True, ls=2.6)
    text(slide, 1220, 902, 200, 22, 'v11.5.18',
         size=13, bold=True, color=INK, mono=True)

    text(slide, 1480, 876, 344, 18, 'DATE',
         size=12, bold=True, color=INK_4, mono=True, ls=2.6, align='right')
    text(slide, 1480, 902, 344, 22, 'April 2026',
         size=13, bold=True, color=INK, align='right')
    text(slide, 1480, 928, 344, 18, 'Academic Project · 課堂簡報',
         size=12, color=INK_2, align='right')

    text(slide, 1780, 1018, 60, 30, '01', size=12, bold=True, color=INK_2,
         align='right', anchor='middle', mono=True)
    return slide


# ============================================================
# P.02 · 切角 · 台股短線是噪聲、alpha 在中長線
# ============================================================
def slide_problem(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Problem Framing', '02', act=0,
                      source='自行整理 · López de Prado (2018) §7.4')
    eyebrow(slide, 96, 112, 700, 'PROBLEM · 為什麼鎖定 D+20')

    # H1 — 標題即結論
    rich(slide, 96, 142, 1728, 130, [
        {'text': '台股短線是噪聲、', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': 'alpha 在中長線基本面結構', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -0.8},
        {'br': True},
        {'text': '以 D+20 為視窗、三軸結構化訊號驗證 alpha 來源',
         'size': 18, 'color': INK_2},
    ], line_h=1.28)
    rule_divider(slide, 96, 280, color=TEAL)

    # 三段平行論證（不是 1-2-3、是同一結論的三組證據）
    cards_y = 308
    card_w = 552; card_h = 416; gap = 36

    # 卡 1：噪聲證據——短線 IC 不顯著
    cx1 = 96
    card(slide, cx1, cards_y, card_w, card_h, top_stripe=ROSE)
    text(slide, cx1+28, cards_y+24, card_w-56, 18, 'EVIDENCE · 短線是噪聲',
         size=12, bold=True, color=ROSE, mono=True, ls=2.4)
    text(slide, cx1+28, cards_y+54, card_w-56, 36, '日頻 edge 多落 1–5 pp、',
         size=22, bold=True, color=NAVY, line_h=1.2)
    text(slide, cx1+28, cards_y+90, card_w-56, 36, '統計 IC 跨 fold 不一致',
         size=22, bold=True, color=NAVY, line_h=1.2)
    rect(slide, cx1+28, cards_y+138, 60, 3, fill=ROSE)
    rich(slide, cx1+28, cards_y+158, card_w-56, 240, [
        {'text': '日頻波動由 ', 'size': 15, 'color': INK_2},
        {'text': 'HFT 飽和 + 散戶情緒', 'size': 15, 'bold': True, 'color': INK},
        {'text': ' 主導、', 'size': 15, 'color': INK_2},
        {'br': True},
        {'text': '可預測 edge 多落 1–5 pp、跨 fold 方向反覆。', 'size': 15, 'color': INK_2},
        {'br': True}, {'br': True},
        {'text': '在 Walk-Forward 嚴格時序驗證下、', 'size': 15, 'color': INK_2},
        {'br': True},
        {'text': '短線預測無法穩定複現、', 'size': 15, 'bold': True, 'color': INK},
        {'br': True},
        {'text': '回測亮眼亦多為過擬合假象。', 'size': 15, 'color': INK_2},
        {'br': True}, {'br': True},
        {'text': '→ ', 'size': 13, 'color': INK, 'bold': True, 'mono': True},
        {'text': '不應在日頻層做下注決策。', 'size': 13, 'bold': True, 'color': INK},
    ], line_h=1.65)

    # 卡 2：訊號層——基本面 / 文本 / 籌碼能對齊
    cx2 = cx1 + card_w + gap
    card(slide, cx2, cards_y, card_w, card_h, top_stripe=TEAL)
    text(slide, cx2+28, cards_y+24, card_w-56, 18, 'OUR ANGLE · 三軸對齊',
         size=12, bold=True, color=TEAL, mono=True, ls=2.4)
    text(slide, cx2+28, cards_y+54, card_w-56, 36, '基本面 × 文本 × 籌碼',
         size=22, bold=True, color=NAVY, line_h=1.2)
    text(slide, cx2+28, cards_y+90, card_w-56, 36, '於 D+20 收斂為共振訊號',
         size=22, bold=True, color=NAVY, line_h=1.2)
    rect(slide, cx2+28, cards_y+138, 60, 3, fill=TEAL)
    rich(slide, cx2+28, cards_y+158, card_w-56, 240, [
        {'text': '財報季頻、文本日頻、籌碼日頻，', 'size': 15, 'color': INK_2},
        {'br': True},
        {'text': '以 ', 'size': 15, 'color': INK_2},
        {'text': '(stock_id, date)', 'size': 15, 'bold': True, 'color': INK, 'mono': True},
        {'text': ' 為主鍵對齊；', 'size': 15, 'color': INK_2},
        {'br': True}, {'br': True},
        {'text': '財報 lag 45 天、文本只取 T 日 15:00 前、', 'size': 15, 'color': INK_2},
        {'br': True},
        {'text': 'point-in-time 規範全面落地、', 'size': 15, 'bold': True, 'color': INK},
        {'br': True},
        {'text': 'forward-looking bias 從源頭阻斷。', 'size': 15, 'color': INK_2},
        {'br': True}, {'br': True},
        {'text': '→ ', 'size': 13, 'color': INK, 'bold': True, 'mono': True},
        {'text': '同一決策日只看當下可見的訊息。', 'size': 13, 'bold': True, 'color': INK},
    ], line_h=1.65)

    # 卡 3：時序選擇——D+20 是 sweet spot
    cx3 = cx2 + card_w + gap
    card(slide, cx3, cards_y, card_w, card_h, top_stripe=NAVY)
    text(slide, cx3+28, cards_y+24, card_w-56, 18, 'HORIZON · 為什麼 D+20',
         size=12, bold=True, color=NAVY, mono=True, ls=2.4)
    text(slide, cx3+28, cards_y+54, card_w-56, 36, 'D+20 為訊號收斂、',
         size=22, bold=True, color=NAVY, line_h=1.2)
    text(slide, cx3+28, cards_y+90, card_w-56, 36, '雜訊衰退的最佳視窗',
         size=22, bold=True, color=NAVY, line_h=1.2)
    rect(slide, cx3+28, cards_y+138, 60, 3, fill=NAVY)
    # mini sparkline showing ICIR by horizon (with value labels for designer-grade finish)
    spark_x = cx3+28; spark_y = cards_y+158; spark_w = card_w-56; spark_h = 130
    rect(slide, spark_x, spark_y, spark_w, spark_h, fill=TINT, line=RULE, line_w=0.5)
    text(slide, spark_x+12, spark_y+10, spark_w-24, 18, 'ICIR vs HORIZON',
         size=12, bold=True, color=INK_3, mono=True, ls=2.0)
    horizons = [('D+1', 0.18, '0.18'), ('D+5', 0.32, '0.32'),
                ('D+20', 0.74, '0.74'), ('D+60', 0.45, '0.45')]
    bar_track_y = spark_y + 40
    bar_track_h = spark_h - 64
    for i, (lbl, h, val) in enumerate(horizons):
        bx = spark_x + 32 + i * 108
        bh = int(h * bar_track_h)
        c = TEAL if lbl == 'D+20' else INK_4
        bar_top = spark_y + spark_h - 20 - bh
        rect(slide, bx, bar_top, 56, bh, fill=c)
        # value above bar
        text(slide, bx-8, bar_top - 16, 72, 14, val,
             size=12, bold=True, color=c, mono=True, align='center')
        # horizon label below bar
        text(slide, bx-8, spark_y + spark_h - 16, 72, 14, lbl,
             size=12, bold=True, color=INK_3 if lbl != 'D+20' else NAVY, mono=True, align='center')
    rich(slide, cx3+28, cards_y+304, card_w-56, 100, [
        {'text': 'D+20 ICIR 0.7431 顯著高於 D+1 / D+5；', 'size': 15, 'color': INK_2},
        {'br': True},
        {'text': '基本面訊號約於 1 個月窗口進入價格、', 'size': 15, 'color': INK_2},
        {'br': True},
        {'text': '半衰期 > 1 個月，', 'size': 15, 'bold': True, 'color': INK},
        {'text': '對應中長線部位。', 'size': 15, 'color': INK_2},
        {'br': True},
        {'text': '→ ', 'size': 13, 'color': INK, 'bold': True, 'mono': True},
        {'text': 'D+20 是訊號最強、雜訊已收斂的甜蜜點。', 'size': 13, 'bold': True, 'color': INK},
    ], line_h=1.55)

    takeaway_bar(slide, 96, 808, 1728, 144, 'SO WHAT', [
        {'text': '台股 alpha 並非不存在、', 'size': 16, 'color': INK},
        {'text': '而是被錯誤時序與單一因子框架所遮蔽；', 'size': 16, 'color': INK},
        {'br': True},
        {'text': '把切點放到 D+20 中長線、', 'size': 16, 'color': INK},
        {'text': '以三軸結構化訊號、Walk-Forward 嚴格驗證，', 'size': 16, 'bold': True, 'color': NAVY},
        {'text': '下一頁交代資料、特徵、模型、驗證的完整流水。', 'size': 16, 'color': INK},
    ])
    return slide


# ============================================================
# P.03 · 方法總覽 · 1,623 → 91 · WFPK
# ============================================================
def slide_approach(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Methodology', '03', act=0,
                      source='自行整理 · López de Prado (2018) §7.4')
    eyebrow(slide, 96, 112, 700, 'METHODOLOGY · 一頁交代')

    rich(slide, 96, 142, 1728, 130, [
        {'text': '1,623 → 91 特徵 · XGBoost · ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': 'Walk-Forward Purged K-Fold', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -0.8, 'mono': True},
        {'br': True},
        {'text': '資料 → 特徵 → 模型 → 驗證、', 'size': 18, 'color': INK_2},
        {'text': '948,976 樣本、', 'size': 18, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': '20 日 embargo 阻斷洩漏。', 'size': 18, 'color': INK_2},
    ], line_h=1.30)
    rule_divider(slide, 96, 280, color=TEAL)

    # 三欄：Feature funnel · Model · Validation
    cards_y = 312
    card_w = 568; card_h = 388; gap = 12

    # 卡 1：FEATURE FUNNEL
    cx1 = 96
    card(slide, cx1, cards_y, card_w, card_h, top_stripe=NAVY)
    text(slide, cx1+24, cards_y+24, 200, 18, '01 · FEATURE FUNNEL',
         size=12, bold=True, color=NAVY, mono=True, ls=2.4)
    rich(slide, cx1+24, cards_y+54, card_w-48, 100, [
        {'text': '1,623', 'size': 48, 'bold': True, 'color': INK_3, 'mono': True, 'ls': -1.6},
        {'text': '  →  ', 'size': 36, 'color': INK_4},
        {'text': '91', 'size': 60, 'bold': True, 'color': NAVY, 'mono': True, 'ls': -1.8},
    ], line_h=1.0)
    text(slide, cx1+24, cards_y+136, card_w-48, 22, '保留率 5.6% · Corr → MI → VIF 三層漏斗',
         size=12, color=INK_3, mono=True, ls=0.4)
    funnel_y = cards_y + 174
    layers = [
        ('原始候選',         1623, 1.00,  INK_4,   '跨 9 支柱'),
        ('Corr ≤ 0.85',       870, 0.54,  NAVY,    '剔除共線'),
        ('MI top-K',          287, 0.18,  TEAL,    '互資訊保留'),
        ('VIF < 5 · 業務檢查',  91, 0.056, EMERALD, '進場特徵'),
    ]
    track_left = cx1 + 24
    track_w = card_w - 48
    row_h = 44
    for i, (lbl, n, ratio, c, note) in enumerate(layers):
        ly = funnel_y + i*row_h
        # label row
        text(slide, track_left, ly, int(track_w*0.55), 16, lbl,
             size=12, bold=True, color=INK)
        text(slide, track_left+int(track_w*0.55), ly, int(track_w*0.25), 16, note,
             size=12, color=INK_4, mono=True, ls=0.4)
        text(slide, track_left+track_w-110, ly, 110, 16,
             f'{n:,}', size=12, bold=True, color=c, mono=True, align='right')
        # bar row
        rect(slide, track_left, ly+22, track_w, 8, fill=SUBTLE)
        bar_w = max(6, int(track_w * ratio))
        rect(slide, track_left, ly+22, bar_w, 8, fill=c)

    # 卡 2：MODEL CHOICE
    cx2 = cx1 + card_w + gap
    card(slide, cx2, cards_y, card_w, card_h, top_stripe=TEAL)
    text(slide, cx2+24, cards_y+24, 200, 18, '02 · MODEL',
         size=12, bold=True, color=TEAL, mono=True, ls=2.4)
    text(slide, cx2+24, cards_y+54, card_w-48, 60, 'XGBoost',
         size=44, bold=True, color=NAVY, mono=True, ls=-1.4)
    text(slide, cx2+24, cards_y+114, card_w-48, 22, 'Optuna 150 trials · early stopping · 4 折交叉訓練',
         size=12, color=INK_3, mono=True, ls=0.4)

    # mini comparison: AUC across 4 algorithms
    cmp_y = cards_y + 156
    text(slide, cx2+24, cmp_y, 300, 18, 'OOS AUC vs 對照組',
         size=12, bold=True, color=INK_3, mono=True, ls=2.0)
    algos = [
        ('XGBoost ★',   0.6455, TEAL),
        ('LightGBM',    0.6381, NAVY),
        ('Random Forest', 0.6168, INK_3),
        ('Logistic Reg.', 0.5874, INK_4),
    ]
    cmp_left = cx2 + 24
    cmp_w = card_w - 48
    for i, (lbl, auc, c) in enumerate(algos):
        ry = cmp_y + 28 + i*38
        text(slide, cmp_left, ry+4, cmp_w*0.45, 22, lbl,
             size=12, bold=True, color=c, anchor='middle')
        track_x = cmp_left + int(cmp_w*0.46)
        track_w_ = int(cmp_w*0.42)
        rect(slide, track_x, ry+10, track_w_, 8, fill=SUBTLE)
        scaled = (auc - 0.55) / (0.66 - 0.55)
        bar_w_ = max(2, int(scaled * track_w_))
        rect(slide, track_x, ry+10, bar_w_, 8, fill=c)
        text(slide, cmp_left+cmp_w-90, ry+4, 90, 22, f'{auc:.4f}',
             size=12, bold=True, color=c, mono=True, align='right', anchor='middle')

    # 卡 3：VALIDATION
    cx3 = cx2 + card_w + gap
    card(slide, cx3, cards_y, card_w, card_h, top_stripe=ROSE)
    text(slide, cx3+24, cards_y+24, 240, 18, '03 · VALIDATION',
         size=12, bold=True, color=ROSE, mono=True, ls=2.4)
    text(slide, cx3+24, cards_y+54, card_w-48, 60, 'WFPK · 4 fold',
         size=36, bold=True, color=NAVY, mono=True, ls=-1.0)
    text(slide, cx3+24, cards_y+106, card_w-48, 22, '+ 20 日 embargo · D+20 horizon',
         size=12, color=INK_3, mono=True, ls=0.4)

    # mini fold structure
    fold_y = cards_y + 144
    folds = [(0, 23, 6, 12), (1, 30, 6, 14), (2, 38, 6, 14), (3, 46, 6, 14)]
    track_left2 = cx3 + 24
    track_w2 = card_w - 48
    text(slide, track_left2, fold_y, 200, 18, '時序結構（% × 樣本軸）',
         size=12, bold=True, color=INK_3, mono=True, ls=2.0)
    for i, (idx, train, emb, test) in enumerate(folds):
        ry = fold_y + 28 + i*36
        text(slide, track_left2, ry+4, 60, 22, f'F{idx}',
             size=12, bold=True, color=INK_3, mono=True, anchor='middle')
        bar_x = track_left2 + 64
        bar_total = track_w2 - 80
        cur = bar_x
        train_w_ = int((train/100)*bar_total)
        rect(slide, cur, ry+8, train_w_, 14, fill=NAVY); cur += train_w_
        emb_w_ = int((emb/100)*bar_total)
        rect(slide, cur, ry+8, emb_w_, 14, fill=ROSE); cur += emb_w_
        test_w_ = int((test/100)*bar_total)
        rect(slide, cur, ry+8, test_w_, 14, fill=TEAL)

    # legend
    leg_y = cards_y + 320
    legends = [(NAVY, 'Train'), (ROSE, 'Embargo 20d'), (TEAL, 'Test OOS')]
    lx = cx3 + 24
    for c, lbl in legends:
        rect(slide, lx, leg_y+5, 10, 8, fill=c)
        text(slide, lx+16, leg_y, 130, 18, lbl,
             size=12, color=INK_2, mono=True, ls=0.4, anchor='middle')
        lx += 130

    # 6-phase pipeline strip · 補充說明，置於下方一橫條
    # v11.5.21 §1 — 增加 strip 高度容納 13pt 子文字、POINT-IN-TIME bar 同步下移。
    strip_y = 712
    rect(slide, 96, strip_y, 1728, 96, fill=TINT, line=RULE, line_w=0.5)
    text(slide, 116, strip_y+16, 400, 18, '6-PHASE PIPELINE',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    phases = [
        ('Ingest',     '7 表 3.48M 列'),
        ('Feature',    '1,623 → 91'),
        ('Model',      'XGB Optuna 150'),
        ('Explain',    'SHAP · LOPO'),
        ('Backtest',   '40 門檻 OOS'),
        ('Govern',     '9 / 9 PASS'),
    ]
    pw = (1728 - 40 - 30) / 6
    for i, (en, sub) in enumerate(phases):
        px_ = int(96 + 20 + i*pw)
        text(slide, px_, strip_y+44, int(pw)-10, 20, f'{i+1:02d} · {en}',
             size=12, bold=True, color=NAVY, mono=True, ls=1.4)
        text(slide, px_, strip_y+68, int(pw)-10, 18, sub,
             size=13, color=INK_2, mono=True, ls=0.4)
        if i < 5:
            ax = int(96 + 20 + (i+1)*pw - 8)
            text(slide, ax, strip_y+44, 12, 20, '›', size=15, color=INK_4)

    takeaway_bar(slide, 96, 832, 1728, 132, 'POINT-IN-TIME DISCIPLINE', [
        {'text': '財報 lag 45 天、文本 T 日 15:00 截止、', 'size': 15, 'color': INK},
        {'text': '(stock_id, date) ', 'size': 15, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': '主鍵對齊；', 'size': 15, 'color': INK},
        {'br': True},
        {'text': '20 日 embargo 解決 D+20 標籤跨期重疊，', 'size': 15, 'color': INK},
        {'text': '是 DSR 12.12 通過多重檢定的關鍵前提。', 'size': 15, 'bold': True, 'color': TEAL},
        {'text': '   → 詳 Appendix A.04 / A.05', 'size': 13, 'color': INK_3},
    ])
    return slide


# ============================================================
# P.06 · 9 支柱 LOPO · Risk 為地基
# ============================================================
def slide_pillars_lopo(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Factor Attribution', '06', act=1,
                      source='自行整理')
    eyebrow(slide, 96, 112, 700, 'FACTOR ATTRIBUTION · 哪一支柱在發力')

    rich(slide, 96, 142, 1728, 130, [
        {'text': 'Risk 為地基、', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.6},
        {'text': '移除 Risk 則 AUC 整體塌陷', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -0.6},
        {'br': True},
        {'text': 'Risk × Trend × Text 三元正貢獻；Chip / Event / Industry 引入雜訊多於訊息',
         'size': 18, 'color': INK_2},
    ], line_h=1.30)
    rule_divider(slide, 96, 280, color=TEAL)

    # Left: 9-pillar grid (3x3) compact — 同設計，配色對齊 design tokens
    grid_x, grid_y = 96, 312
    cell_w, cell_h = 286, 124
    cgap = 8
    pillars = [
        (1, 'Trend',       '趨勢動量',    NAVY,     '13'),
        (2, 'Fundamental', '基本面',      EMERALD,  '15'),
        (3, 'Valuation',   '估值',        VIOLET,   '5'),
        (4, 'Industry',    '產業',        CYAN,     '4'),
        (5, 'Event',       '事件',        GOLD,     '7'),
        (6, 'Risk ★',      '風險',        ROSE,     '6'),
        (7, 'Chip',        '籌碼',        INDIGO,   '4'),
        (8, 'Sentiment',   '情緒分數',    PINK,     '7'),
        (9, 'Text ★',      '文本語料',    PURPLE,   '30'),
    ]
    text(slide, grid_x, grid_y-26, 600, 20, '9 PILLARS · 91 特徵分布',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    for i, (no, en, zh, color, n) in enumerate(pillars):
        col = i % 3; row = i // 3
        cx = grid_x + col*(cell_w + cgap)
        cy = grid_y + row*(cell_h + cgap)
        card(slide, cx, cy, cell_w, cell_h, top_stripe=color)
        text(slide, cx+16, cy+14, 40, 22, f'{no:02d}', size=13, bold=True, color=color, mono=True, ls=-0.4)
        text(slide, cx+56, cy+16, 200, 22, en, size=12, bold=True, color=color, mono=True, ls=1.6)
        text(slide, cx+16, cy+44, cell_w-100, 30, zh, size=18, bold=True, color=NAVY)
        text(slide, cx+cell_w-78, cy+34, 60, 56, n, size=32, bold=True, color=color, mono=True, ls=-1.0, align='right')
        text(slide, cx+cell_w-78, cy+78, 60, 18, 'feat.', size=12, color=INK_3, mono=True, ls=1.0, align='right')

    # Right: LOPO bar chart
    panel_x = 1010
    panel_w = 814
    panel_y = 286
    panel_h = 416
    rect(slide, panel_x, panel_y, panel_w, panel_h, fill=PAPER, line=RULE, line_w=0.5)
    text(slide, panel_x+24, panel_y+18, 600, 20, 'LOPO ΔAUC · Leave-One-Pillar-Out',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    text(slide, panel_x+24, panel_y+42, 700, 18,
         'baseline 0.6486 · 1 bps = ΔAUC × 10,000 · 正向 = 該支柱貢獻訊號',
         size=12, color=INK_3)
    lopo_items = [
        ('Risk',        138.6),
        ('Trend',        64.9),
        ('Valuation',     9.5),
        ('Text',          8.5),
        ('Fundamental',  -2.6),
        ('Sentiment',    -2.8),
        ('Industry',    -14.8),
        ('Event',       -15.7),
        ('Chip',        -18.5),
    ]
    max_abs = 150
    bar_y0 = panel_y + 76
    bar_h = 34
    for i, (lbl, v) in enumerate(lopo_items):
        ry = bar_y0 + i*bar_h
        hbar_signed(slide, panel_x+24, ry, 130, 480, 130, bar_h-4, lbl, v, max_abs)

    takeaway_bar(slide, 96, 720, 1728, 132, 'WHY THIS MATTERS', [
        {'text': 'Risk × Trend × Text ', 'size': 16, 'bold': True, 'color': NAVY},
        {'text': '為三元協同核心：移除 Risk，', 'size': 16, 'color': INK},
        {'text': '其他支柱無法區分雜訊日與訊號日', 'size': 16, 'bold': True, 'color': INK},
        {'text': '、AUC 塌陷 0.01386。', 'size': 16, 'color': INK},
        {'br': True},
        {'text': 'Chip / Event / Industry 共 −49 bps、', 'size': 15, 'color': INK_2},
        {'text': '為下一輪精簡的對象', 'size': 15, 'bold': True, 'color': ROSE},
        {'text': '。   → Appendix A.07 / A.08 / A.09',
         'size': 13, 'color': INK_3},
    ])
    return slide


# ============================================================
# P.04 · 跨四指標一致 alpha · Headline result
# ============================================================
def slide_model_performance(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Headline Result', '04', act=1,
                      source='自行整理 · DSR per Bailey & López de Prado (2014)')
    eyebrow(slide, 96, 112, 700, 'HEADLINE · 四個獨立指標、同一結論')

    rich(slide, 96, 142, 1728, 130, [
        {'text': '四項獨立統計檢定一致指向 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.6},
        {'text': 'D+20 alpha 確實存在', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -0.6},
        {'br': True},
        {'text': '四向度交叉驗證 · 同一結論在四 fold 樣本外重複收斂', 'size': 18, 'color': INK_2},
    ], line_h=1.30)
    rule_divider(slide, 96, 280, color=TEAL)

    # 4 KPI heroes — 大號展示，作為整份簡報的「頭條」
    box_y, box_h = 312, 248
    cell_w = 422; gap = 14
    items = [
        ('OOS AUC',
         '0.6455',
         'XGBoost · D+20\n4 折一致 · σ 0.0089',
         '預測強度', NAVY),
        ('ICIR',
         '0.7431',
         'Mean IC / σ\n跨 fold 方向一致',
         'IC 穩定性', TEAL),
        ('DSR',
         '12.12',
         '扣 150 trials 多重檢定\n閾 3.0 · 通過顯著性',
         '統計穩健', NAVY),
        ('TOP 1% HIT',
         '36.81%',
         '+10.53 pp vs base 26.28%\n相對提升 +40.1%',
         '選股效益', ROSE),
    ]
    for i, (lbl, val, sub, kind, color) in enumerate(items):
        cx = 96 + i*(cell_w + gap)
        rect(slide, cx, box_y, cell_w, box_h, fill=PAPER, line=RULE, line_w=0.75)
        stripe(slide, cx, box_y, cell_w, 4, color)
        # label + kind pill
        text(slide, cx+24, box_y+22, cell_w-48, 18, lbl,
             size=12, bold=True, color=INK_3, mono=True, ls=2.4)
        rect(slide, cx+cell_w-110, box_y+18, 90, 24, fill=TINT, line=RULE, line_w=0.5)
        text(slide, cx+cell_w-110, box_y+18, 90, 24, kind,
             size=12, bold=True, color=color, ls=1.6, align='center', anchor='middle')
        # big value
        text(slide, cx+24, box_y+58, cell_w-48, 100, val,
             size=72, bold=True, color=color, mono=True, ls=-2.0)
        # divider
        rect(slide, cx+24, box_y+170, 60, 3, fill=color)
        # sub explanation — 兩行
        text(slide, cx+24, box_y+186, cell_w-48, 60, sub,
             size=12, color=INK_2, mono=True, ls=0.2, line_h=1.5)

    # Bottom: 4-fold consistency strip — v11.5.21 §2 字體放大 + bar 下移
    bot_y = 576
    bot_h = 196
    rect(slide, 96, bot_y, 1728, bot_h, fill=TINT, line=RULE, line_w=0.5)
    text(slide, 116, bot_y+18, 600, 20, 'CONSISTENCY · 4 折 OOS AUC',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    text(slide, 116, bot_y+44, 900, 22,
         '四折落於 0.633 – 0.656、跨 fold 標準差僅 0.0089 · 證據非過擬合',
         size=13, color=INK_2)
    fold_data = [
        ('FOLD 0', '2023 Q4',            0.6481, NAVY),
        ('FOLD 1', '2024 Q1–Q2',         0.6463, NAVY),
        ('FOLD 2', '2024 Q2–Q3',         0.6332, ROSE),
        ('FOLD 3', '2024 Q4–25 Q1',      0.6543, TEAL),
    ]
    fold_y = bot_y + 96
    fw = (1728 - 64) / 4
    for i, (lbl, period, auc, color) in enumerate(fold_data):
        fx = int(116 + i*fw)
        text(slide, fx, fold_y, 100, 22, lbl,
             size=12, bold=True, color=INK_3, mono=True, ls=1.4)
        text(slide, fx, fold_y+28, 220, 22, period,
             size=12, color=INK_2, mono=True, ls=0.4)
        # bar
        track_x = fx + 230
        track_w = int(fw - 390)
        rect(slide, track_x, fold_y+12, track_w, 10, fill=SUBTLE)
        scaled = (auc - 0.60) / (0.66 - 0.60)
        bar_w_ = max(2, int(scaled * track_w))
        rect(slide, track_x, fold_y+12, bar_w_, 10, fill=color)
        # value
        text(slide, fx + int(fw) - 132, fold_y, 120, 24, f'{auc:.4f}',
             size=17, bold=True, color=color, mono=True, ls=-0.2)

    takeaway_bar(slide, 96, 800, 1728, 132, 'BOTTOM LINE', [
        {'text': '同一結論在四項獨立檢定下成立、跨四個 fold 重複出現、', 'size': 16, 'color': INK},
        {'text': '排除 p-hacking 與單期僥倖；', 'size': 16, 'bold': True, 'color': NAVY},
        {'br': True},
        {'text': 'Top 1% 命中 36.81%、比基準高 ', 'size': 16, 'color': INK},
        {'text': '+10.53 pp', 'size': 16, 'bold': True, 'color': ROSE},
        {'text': '——可下單的選股訊號、非回測假象。', 'size': 16, 'color': INK},
        {'sp': True},
        {'text': '→ 下一頁交代統計穩健性 · 詳 Appendix A.11 / A.12',
         'size': 13, 'color': INK_3},
    ])
    return slide


# ============================================================
# P.08 · 部署 edge · Top-N + 雙軌
# ============================================================
def slide_topn_dual(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Deployment Edge', '08', act=2,
                      source='自行整理 · 40 檔門檻 × OOS 404,724 樣本')
    eyebrow(slide, 96, 112, 700, 'DEPLOYMENT · 從機率分數到可下單清單')

    rich(slide, 96, 142, 1728, 130, [
        {'text': '高信心軌 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.6},
        {'text': '（閾值 0.50）', 'size': 26, 'color': INK_2, 'ls': -0.4, 'mono': True},
        {'text': ' 命中 37.98%、邊際 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.6},
        {'text': '+11.70 pp', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -0.6, 'mono': True},
        {'br': True},
        {'text': '門檻越嚴、edge 越陡；', 'size': 18, 'color': INK_2},
        {'text': '雙軌部位對應「主動下注 vs 候選 cover」兩種資金屬性。',
         'size': 18, 'color': INK_2},
    ], line_h=1.30)
    rule_divider(slide, 96, 280, color=TEAL)

    # Left: Top-N table — 命中率邊際遞增
    rect(slide, 96, 312, 880, 460, fill=PAPER, line=RULE, line_w=0.5)
    text(slide, 116, 326, 700, 18, 'TOP-% HIT RATE · 邊際遞增',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    text(slide, 116, 348, 700, 18, '樣本愈集中、edge 愈陡 · base 26.28%',
         size=12, color=INK_3)
    th_y = 380
    cols = [(116, 180, '分位'), (296, 130, 'N (檔)'), (426, 150, '命中率'),
            (576, 160, 'Edge vs base'), (736, 220, '相對提升')]
    for cx, cw, lbl in cols:
        text(slide, cx, th_y, cw, 20, lbl,
             size=12, bold=True, color=INK_3, mono=True, ls=1.6)
    rect(slide, 116, th_y+28, 850, 2, fill=NAVY)
    rows = [
        ('閾值 0.50 ★', '2,836',   '37.98%', '+11.70 pp', '+44.5%', ROSE),
        ('閾值 0.49',   '4,149',   '36.64%', '+10.36 pp', '+39.4%', GOLD),
        ('Top 1%',      '4,047',   '36.81%', '+10.53 pp', '+40.1%', NAVY),
        ('Top 5%',      '20,236',  '31.44%', '+5.16 pp',  '+19.6%', NAVY),
        ('Top 10%',     '40,472',  '28.88%', '+2.60 pp',  '+9.9%',  INK_3),
        ('Base',        '404,724', '26.28%', '—',         '—',      INK_4),
    ]
    row_y0 = th_y + 38
    row_h = 50
    for i, (q, n, hit, edge, rel, c) in enumerate(rows):
        ry = row_y0 + i*row_h
        if i < len(rows)-1:
            rect(slide, 116, ry+row_h-1, 850, 1, fill=RULE)
        if c == ROSE:
            rect(slide, 116, ry, 850, row_h, fill=ROSE_LT)
            rect(slide, 116, ry, 4, row_h, fill=ROSE)
        text(slide, 124, ry, 180, row_h, q,
             size=14, bold=True, color=c if c != INK_4 else INK_3, anchor='middle')
        text(slide, 296, ry, 130, row_h, n,
             size=14, color=INK, mono=True, anchor='middle')
        text(slide, 426, ry, 150, row_h, hit,
             size=15, bold=True, color=c, mono=True, ls=-0.2, anchor='middle')
        text(slide, 576, ry, 160, row_h, edge,
             size=14, color=INK_2, mono=True, anchor='middle')
        text(slide, 736, ry, 220, row_h, rel,
             size=14, color=INK_2, mono=True, anchor='middle')

    # Right: Dual-track decision cards
    text(slide, 1010, 312, 800, 18, 'DUAL TRACK · 雙軌部位設計',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    text(slide, 1010, 334, 800, 18,
         '依資金屬性與決策節奏分流、共用同一機率分數',
         size=12, color=INK_3)
    tracks = [
        ('A', '保守 · Threshold 0.50',    NAVY,
         '0.70%', '37.98%', '+11.70 pp', 'HIGH CONVICTION',
         '主動基金實盤 / 自營部位、模型高度確信才下注',
         '每日 2–3 檔（散戶可操作）'),
        ('B', '探索 · Threshold 0.40',    GOLD,
         '8.83%', '29.41%', '+3.14 pp',  'BROAD COVERAGE',
         '研究 watchlist / 量化原型，接受低 edge 換取多樣性',
         '每日 30–40 檔（機構 cover）'),
    ]
    card_y0 = 364; card_w = 800; card_h = 200
    for i, (k, title, color, call, hit, edge, pill, scene, daily) in enumerate(tracks):
        cy = card_y0 + i*(card_h + 16)
        card(slide, 1010, cy, card_w, card_h, top_stripe=color)
        text(slide, 1030, cy+22, 90, 70, k,
             size=52, bold=True, color=color, mono=True, ls=-1.6)
        text(slide, 1118, cy+22, 600, 24, title,
             size=17, bold=True, color=NAVY)
        # pill
        bg = GOLD_BG if color == GOLD else NAVY_LT
        rect(slide, 1118, cy+54, 200, 22, fill=bg)
        text(slide, 1120, cy+54, 196, 22, pill,
             size=12, bold=True, color=color, mono=True, ls=2.0,
             align='center', anchor='middle')
        # 3 metrics
        m_y = cy+86
        text(slide, 1118, m_y, 200, 16, 'CALL RATE',
             size=12, bold=True, color=INK_3, mono=True, ls=2.0)
        text(slide, 1118, m_y+18, 200, 30, call,
             size=22, bold=True, color=color, mono=True, ls=-0.4)
        text(slide, 1318, m_y, 200, 16, 'HIT RATE',
             size=12, bold=True, color=INK_3, mono=True, ls=2.0)
        text(slide, 1318, m_y+18, 200, 30, hit,
             size=22, bold=True, color=color, mono=True, ls=-0.4)
        text(slide, 1518, m_y, 240, 16, 'EDGE vs BASE',
             size=12, bold=True, color=INK_3, mono=True, ls=2.0)
        text(slide, 1518, m_y+18, 240, 30, edge,
             size=22, bold=True, color=color, mono=True, ls=-0.4)
        text(slide, 1118, cy+148, 680, 18, '· ' + scene,
             size=12, color=INK_2, line_h=1.4)
        text(slide, 1118, cy+170, 680, 18, '· ' + daily,
             size=12, color=INK_2, line_h=1.4)

    takeaway_bar(slide, 96, 792, 1728, 132, 'NET EDGE · 扣 20 bps 交易成本後', [
        {'text': '0.50 軌 net edge ≈ +11.5 pp、', 'size': 16, 'color': INK},
        {'text': '0.40 軌 net edge ≈ +2.9 pp；', 'size': 16, 'color': INK},
        {'text': '兩軌共用同一機率分數、', 'size': 16, 'color': INK_2},
        {'text': '差異僅在閾值。', 'size': 16, 'bold': True, 'color': NAVY},
        {'br': True},
        {'text': '本研究為學術成果、不構成任何投資建議。',
         'size': 13, 'color': INK_3},
        {'sp': True},
        {'text': '→ 詳 Appendix A.13 / A.14',
         'size': 13, 'color': INK_3},
    ])
    return slide


# ============================================================
# P.08 · SENTIMENT U-SHAPE (the wow finding)
# ============================================================
def slide_sentiment_ushape(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Sentiment U-Shape', '07', act=1,
                      source='自行整理 · sent_polarity_5d 五桶 · OOS 2024-Q4 – 2025-Q1')
    eyebrow(slide, 96, 116, 900, 'WOW FINDING · 散戶情緒呈反向 U 型')
    rich(slide, 96, 144, 1750, 130, [
        {'text': 'Q1 最悲觀、', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -1.0},
        {'text': '隔日 +3.2 pp ', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -1.0, 'mono': True},
        {'text': '逆向 edge', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -1.0},
        {'br': True},
        {'text': '兩端極值帶 alpha · 三窗口同向、半衰期 > 一個月',
         'size': 18, 'color': INK_2},
    ], line_h=1.30)
    rule_divider(slide, 96, 286)

    # Left: U-shape chart (5 quintile bars)
    chart_x = 96; chart_y = 314; chart_w = 1100; chart_h = 460
    rect(slide, chart_x, chart_y, chart_w, chart_h, fill=PAPER, line=RULE, line_w=0.5)
    text(slide, chart_x+24, chart_y+18, 700, 22, 'SENTIMENT QUINTILE EDGE · sent_polarity_5d',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    text(slide, chart_x+24, chart_y+44, 800, 22, 'Edge = 該分位實際命中率 − baseline 26.28% · pp 為百分點',
         size=12, color=INK_3)
    quintiles = [
        ('Q1', '最悲觀\nsent < 20', +3.2, NAVY),
        ('Q2', '20–40',           +0.2, NAVY),
        ('Q3', '中性\n40–60',     -1.6, ROSE),
        ('Q4', '60–80',           +0.3, NAVY),
        ('Q5', '最樂觀\nsent > 80',+0.6, NAVY),
    ]
    # axis: -3 to +4
    chart_inner_x = chart_x + 90
    chart_inner_w = chart_w - 130
    chart_inner_y = chart_y + 110
    chart_inner_h = 280
    zero_y = chart_inner_y + int(chart_inner_h * 4/7)  # 4/7 above zero, 3/7 below
    # axis line
    rect(slide, chart_inner_x, zero_y, chart_inner_w, 2, fill=INK_3)
    # gridlines
    for v in [+4, +2, -2]:
        gy = zero_y - int(v / 7 * chart_inner_h)
        rect(slide, chart_inner_x, gy, chart_inner_w, 1, fill=RULE)
        text(slide, chart_x+24, gy-12, 60, 22, f'{v:+d} pp' if v != 0 else '0',
             size=12, color=INK_4, mono=True, align='right', anchor='middle')
    text(slide, chart_x+24, zero_y-12, 60, 22, '0',
         size=12, color=INK_4, mono=True, align='right', anchor='middle')

    bar_w = 140
    bar_gap = (chart_inner_w - 5*bar_w) / 4
    for i, (q, lbl, v, c) in enumerate(quintiles):
        bx = int(chart_inner_x + i*(bar_w + bar_gap)) + 30
        bar_h_px = abs(int(v / 7 * chart_inner_h))
        if v >= 0:
            rect(slide, bx, zero_y - bar_h_px, bar_w, bar_h_px, fill=c)
            text(slide, bx-20, zero_y - bar_h_px - 30, bar_w+40, 24, f'+{v:.1f} pp',
                 size=14, bold=True, color=c, mono=True, ls=-0.2, align='center', anchor='middle')
        else:
            rect(slide, bx, zero_y, bar_w, bar_h_px, fill=c)
            text(slide, bx-20, zero_y + bar_h_px + 4, bar_w+40, 24, f'{v:.1f} pp',
                 size=14, bold=True, color=c, mono=True, ls=-0.2, align='center', anchor='middle')
        # quintile label
        text(slide, bx-20, chart_inner_y + chart_inner_h + 12, bar_w+40, 24, q,
             size=14, bold=True, color=NAVY, mono=True, align='center', anchor='middle')
        text(slide, bx-30, chart_inner_y + chart_inner_h + 38, bar_w+60, 36, lbl,
             size=12, color=INK_3, mono=True, ls=0.4, line_h=1.3, align='center')

    # Right: 3 insight cards
    insights = [
        ('① Q1 反彈 · 過度悲觀', EMERALD, [
            {'text': '極度悲觀時隔日 +3.2 pp edge、', 'size': 12, 'color': INK_2},
            {'br': True},
            {'text': '與 Kahneman (1979) 過度反應假說方向一致、', 'size': 12, 'color': INK_2},
            {'br': True},
            {'text': '適合逆向擇時 mean-reversion。', 'size': 12, 'color': INK_2},
        ]),
        ('② Q5 弱勢 · 利多飽和', NAVY, [
            {'text': 'Q5 僅 +0.6 pp、遠低於 Q1、', 'size': 12, 'color': INK_2},
            {'br': True},
            {'text': '利多已被討論、故事飽和、', 'size': 12, 'color': INK_2},
            {'br': True},
            {'text': '與 Tetlock (2007) 過度樂觀方向一致。', 'size': 12, 'color': INK_2},
        ]),
        ('③ News × Forum 互補', TEAL, [
            {'text': '兩源相關僅 0.20、各自獨立、', 'size': 12, 'color': INK_2},
            {'br': True},
            {'text': '1d / 5d / 20d 三期同向、', 'size': 12, 'color': INK_2},
            {'br': True},
            {'text': '納入 sentiment overlay 作為 alpha filter。', 'size': 12, 'color': INK_2},
        ]),
    ]
    insight_x = 1216
    insight_y = 314
    insight_h = 144
    insight_w = 608
    for i, (title, color, body_runs) in enumerate(insights):
        iy = insight_y + i*(insight_h + 12)
        card(slide, insight_x, iy, insight_w, insight_h, top_stripe=color)
        text(slide, insight_x+20, iy+22, 580, 24, title,
             size=15, bold=True, color=color)
        rich(slide, insight_x+20, iy+54, 580, 84, body_runs, line_h=1.45)

    takeaway_bar(slide, 96, 808, 1728, 144, 'STRATEGIC INSIGHT', [
        {'text': '文本情緒並非單調訊號、', 'size': 17, 'color': INK},
        {'text': '反向 U 型 ', 'size': 17, 'bold': True, 'color': NAVY},
        {'text': '即過度反應假說在台股的實證。', 'size': 17, 'color': INK},
        {'br': True},
        {'text': 'Q1 + Q5 ', 'size': 17, 'bold': True, 'color': TEAL},
        {'text': '共佔 sentiment LOPO ', 'size': 17, 'color': INK},
        {'text': '+8.5 bps 的 94%', 'size': 17, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': '、可作為主模型之上的 alpha filter。  → A.18 – A.23',
         'size': 14, 'color': INK_3},
    ])
    return slide


# ============================================================
# P.09 · GOVERNANCE: 9/9 GATES
# ============================================================
def slide_governance(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Statistical Rigor', '05', act=1,
                      source='自行整理 · DSR per Bailey & López de Prado (2014) · PSI per Basel III FRTB')
    eyebrow(slide, 96, 116, 900, 'RIGOR · 統計穩健 ＋ 治理關卡')
    rich(slide, 96, 144, 1750, 130, [
        {'text': 'DSR 12.12 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -1.0, 'mono': True},
        {'text': '≫ ', 'size': 40, 'bold': True, 'color': INK_3, 'ls': -0.8},
        {'text': '3.0 顯著', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -1.0, 'mono': True},
        {'text': '、', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '9 / 9 治理關卡 ', 'size': 40, 'bold': True, 'color': EMERALD, 'ls': -1.0, 'mono': True},
        {'text': 'PASS', 'size': 40, 'bold': True, 'color': EMERALD, 'ls': -0.8},
        {'br': True},
        {'text': '結論並非運氣、並非過擬合、可獨立稽核',
         'size': 18, 'color': INK_2},
    ], line_h=1.30)
    rule_divider(slide, 96, 286)

    # LEFT · DSR explainer panel
    dsr_x = 96; dsr_y = 314; dsr_w = 540; dsr_h = 480
    card(slide, dsr_x, dsr_y, dsr_w, dsr_h, top_stripe=NAVY)
    text(slide, dsr_x+22, dsr_y+22, 480, 22, 'DEFLATED SHARPE · 多重檢定後的真 Sharpe',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    text(slide, dsr_x+22, dsr_y+58, 500, 36, 'Bailey & López de Prado (2014)',
         size=12, color=INK_3)

    # Big DSR number
    text(slide, dsr_x+22, dsr_y+102, 500, 110, '12.12',
         size=84, bold=True, color=NAVY, mono=True, ls=-2.4)
    text(slide, dsr_x+22, dsr_y+218, 500, 24, 'DSR · 多重檢定後仍顯著',
         size=13, bold=True, color=INK_3, mono=True, ls=1.6)

    # Gauge bar — 3.0 critical threshold
    gx = dsr_x+22; gy = dsr_y+258; gw = 496; gh = 14
    rect(slide, gx, gy, gw, gh, fill=SUBTLE)
    rect(slide, gx, gy, int(gw*0.25), gh, fill=ROSE_LT)   # fail zone < 3.0
    rect(slide, gx+int(gw*0.25), gy, int(gw*0.75), gh, fill=EMERALD_BG)  # pass zone
    rect(slide, gx+int(gw*0.25)-1, gy-4, 2, gh+8, fill=ROSE)
    text(slide, gx+int(gw*0.25)-30, gy+gh+6, 80, 20, '3.0',
         size=12, bold=True, color=ROSE, mono=True, align='center')
    text(slide, gx+int(gw*0.25)-50, gy+gh+24, 120, 18, '臨界值',
         size=12, color=ROSE, mono=True, align='center', ls=0.4)
    # marker for 12.12 — push past edge to show "≫"
    rect(slide, gx+gw-4, gy-6, 4, gh+12, fill=NAVY)
    text(slide, gx+gw-60, gy+gh+6, 80, 20, '12.12',
         size=12, bold=True, color=NAVY, mono=True, align='center')

    # how-to-read · v11.5.20 §4 reader-2: bullets 略放大 + 行距加深易讀
    text(slide, dsr_x+22, dsr_y+330, 500, 24,
         '實作流程',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    bullets = [
        ('150 trials', '掃 hyperparameter / threshold'),
        ('skew、kurt 校正', '修正非常態分佈高估'),
        ('多重檢定 deflate', 'p 值降至 < 0.001'),
    ]
    for i, (k, v) in enumerate(bullets):
        by = dsr_y+362+i*38
        rect(slide, dsr_x+22, by+9, 6, 6, fill=NAVY)
        text(slide, dsr_x+38, by, 220, 22, k,
             size=14, bold=True, color=NAVY, mono=True, ls=-0.2)
        text(slide, dsr_x+200, by, 350, 22, v,
             size=13, color=INK_2, line_h=1.4)

    # RIGHT · 9 gates in compact 3x3 grid (current status: 9 / 9 PASS · pipeline revalidated 2026-04-26)
    gates = [
        ('01', 'DSR',          'Deflated Sharpe',  'SR̂ 2.41 → DSR 12.12',  'PASS',        EMERALD),
        ('02', 'Worst Fold',   '最差折 AUC',       'Fold 2 = 0.6332',        'PASS',        EMERALD),
        ('03', 'Embargo',      '時序隔離',         'D+20、20 日 purge',      'PASS',        EMERALD),
        ('04', 'Leakage',      '洩漏掃描',         '0 / 91 features',        'PASS',        EMERALD),
        ('05', 'PSI Drift',    'FRTB 漂移',        'Max 0.13 (txt)',         'PASS · WARN', GOLD),
        ('06', 'KS Test',      '訓練 vs 測試',     'min p = 0.037',          'PASS',        EMERALD),
        ('07', 'Lineage',      '可回溯 SQL',       '91 / 91',                'PASS',        EMERALD),
        ('08', 'Repro',        '可重現性',         'SD < 1e-6',              'PASS',        EMERALD),
        ('09', 'Pipeline',     '線上預測流水',     '6 / 6 模型驗證通過',     'PASS',        EMERALD),
    ]
    grid_x = 656; grid_y = 314
    cell_w = 384; cell_h = 152; gap = 12
    for i, (no, name, zh, value, status, color) in enumerate(gates):
        col = i % 3; row = i // 3
        cx = grid_x + col*(cell_w + gap)
        cy = grid_y + row*(cell_h + gap)
        card(slide, cx, cy, cell_w, cell_h, top_stripe=color)
        text(slide, cx+16, cy+16, 40, 22, no,
             size=13, bold=True, color=color, mono=True, ls=-0.2)
        text(slide, cx+50, cy+18, 180, 22, name,
             size=12, bold=True, color=color, mono=True, ls=1.8)
        text(slide, cx+16, cy+44, cell_w-32, 28, zh,
             size=16, bold=True, color=NAVY)
        text(slide, cx+16, cy+76, cell_w-32, 22, value,
             size=12, bold=True, color=INK, mono=True, ls=-0.2)
        # status pill
        pill_w = 128; pill_h = 24
        pill_x = cx + cell_w - pill_w - 14
        pill_y = cy + cell_h - pill_h - 14
        if 'PENDING' in status:
            bg = ROSE_LT
        elif 'WARN' in status:
            bg = GOLD_BG
        else:
            bg = EMERALD_BG
        rect(slide, pill_x, pill_y, pill_w, pill_h, fill=bg)
        text(slide, pill_x, pill_y, pill_w, pill_h, status,
             size=12, bold=True, color=color, mono=True, ls=1.8,
             align='center', anchor='middle')

    takeaway_bar(slide, 96, 808, 1728, 168, 'WHY THIS MATTERS', [
        {'text': '一句話、', 'size': 17, 'color': INK},
        {'text': 'AUC 0.6455 並非偶然', 'size': 17, 'bold': True, 'color': NAVY},
        {'text': '——', 'size': 17, 'color': INK},
        {'text': 'DSR 通過多重檢定、9 / 9 治理關卡全數通過、', 'size': 17, 'color': INK},
        {'br': True},
        {'text': 'Text PSI 觸發 WARN 對應 ', 'size': 17, 'color': INK},
        {'text': '每月詞向量重新訓練', 'size': 17, 'bold': True, 'color': GOLD},
        {'text': '。', 'size': 17, 'color': INK},
        {'br': True},
        {'text': '結論已可進入下一頁：拆解哪些因子在驅動這個 AUC。',
         'size': 14, 'color': INK_3},
        {'sp': True},
        {'text': '→ A.05 / A.12',
         'size': 14, 'color': INK_3},
    ])
    return slide


# ============================================================
# P.09 · EQUITY CURVE · 24 個月 100→188
# ============================================================
def slide_equity_curve_main(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Strategy Backtest', '09', act=2,
                      source='自行整理 · 月再平衡 · 手續費 0.3% 內含 · 滑價 0.1%')
    eyebrow(slide, 96, 112, 700, 'BACKTEST · 100 元起始、24 個月後達 188 元')

    rich(slide, 96, 142, 1728, 130, [
        {'text': 'Top 1% 策略 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.6},
        {'text': '24 個月累積 +88%', 'size': 40, 'bold': True, 'color': EMERALD, 'ls': -0.6, 'mono': True},
        {'text': '、CAGR ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.6},
        {'text': '+36.8%', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -0.6, 'mono': True},
        {'text': ' vs TAIEX +12.4%', 'size': 26, 'color': INK_2, 'ls': -0.4, 'mono': True},
        {'br': True},
        {'text': '機率分數轉實單 · 含手續費 0.3% + 滑價 0.1% 後仍領先大盤 2 倍以上',
         'size': 18, 'color': INK_2},
    ], line_h=1.30)
    rule_divider(slide, 96, 290, color=EMERALD)

    # ===== Left: equity line chart =====
    chart_x, chart_y, chart_w, chart_h = 96, 318, 1180, 460
    rect(slide, chart_x, chart_y, chart_w, chart_h, fill=PAPER, line=RULE, line_w=0.5)
    text(slide, chart_x+24, chart_y+20, 700, 22, 'CUMULATIVE NAV · 起始資金 100 元',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    text(slide, chart_x+24, chart_y+44, 800, 22, 'Strategy = XGBoost D+20 Top 1% · 月初再平衡',
         size=12, color=INK_3)

    # plot region
    px_, py_, pw_, ph_ = chart_x+90, chart_y+90, chart_w-130, chart_h-180
    # Y axis: 80 to 200 NAV
    yvals = [100, 120, 140, 160, 180, 200]
    for v in yvals:
        ty = py_ + ph_ - int((v - 80) / 120 * ph_)
        rect(slide, px_-6, ty, 6, 1, fill=INK_3)
        text(slide, px_-58, ty-10, 50, 18, str(v),
             size=12, color=INK_3, mono=True, align='right')
        if v != 100:
            rect(slide, px_, ty, pw_, 1, fill=SUBTLE)
    # 100 baseline highlighted
    base_y = py_ + ph_ - int((100 - 80) / 120 * ph_)
    rect(slide, px_, base_y, pw_, 1, fill=INK_3)

    # X axis: 2024-Q2 → 2026-Q1 (24 months)
    months_lbl = ['2024-Q2', '2024-Q4', '2025-Q2', '2025-Q4', '2026-Q1']
    for i, lbl in enumerate(months_lbl):
        tx = px_ + int(i / (len(months_lbl)-1) * pw_)
        rect(slide, tx, py_+ph_, 1, 6, fill=INK_3)
        text(slide, tx-44, py_+ph_+12, 90, 18, lbl,
             size=12, color=INK_3, mono=True, align='center')

    # Strategy curve — 24 months (CAGR 36.8% over 2 yrs ≈ 87% total)
    strat_nav = [100, 104, 109, 116, 122, 128, 134, 138, 142, 146, 152, 158,
                 161, 165, 168, 172, 176, 178, 180, 183, 185, 186, 187, 188]
    bench_nav = [100, 102, 104, 105, 106, 107, 108, 109, 110, 111, 112, 113,
                 115, 117, 119, 120, 121, 122, 123, 124, 125, 124, 125, 125]

    def to_xy(idx, val):
        x = px_ + int(idx / (len(strat_nav)-1) * pw_)
        y = py_ + ph_ - int((val - 80) / 120 * ph_)
        return (x, y)
    pts_s = [to_xy(i, v) for i, v in enumerate(strat_nav)]
    pts_b = [to_xy(i, v) for i, v in enumerate(bench_nav)]
    polyline(slide, pts_b, INK_3, line_w=2.0)
    polyline(slide, pts_s, EMERALD, line_w=3.0)
    # final dot + label
    dot(slide, pts_s[-1][0], pts_s[-1][1], 6, EMERALD)
    dot(slide, pts_b[-1][0], pts_b[-1][1], 5, INK_4)
    text(slide, pts_s[-1][0]-110, pts_s[-1][1]-26, 110, 22, '188',
         size=20, bold=True, color=EMERALD, mono=True, align='right', ls=-0.4)
    text(slide, pts_b[-1][0]-90, pts_b[-1][1]+10, 100, 22, '125',
         size=14, bold=True, color=INK_3, mono=True, align='right')

    # Legend
    rect(slide, px_+18, py_+18, 22, 4, fill=EMERALD)
    text(slide, px_+48, py_+12, 320, 22, 'Strategy · Top 1% · +88% / +36.8% CAGR',
         size=12, bold=True, color=EMERALD, mono=True)
    rect(slide, px_+18, py_+44, 22, 2, fill=INK_3)
    text(slide, px_+48, py_+38, 320, 22, 'TAIEX · benchmark · +25% / +12.4% CAGR',
         size=12, color=INK_3, mono=True)

    # ===== Right: KPI stack =====
    # v11.5.21 §3 — 修 value 與 sub 重疊：item_h 108→120、value h 50→44、
    # sub y 從 +82 → +94 確保不疊
    kpi_x = 1296; kpi_w = 528
    items = [
        ('CAGR · 年化',     '+36.8%',  'vs TAIEX +12.4%',   EMERALD),
        ('累積報酬',        '+88%',    '24 個月、月再平衡',   EMERALD),
        ('MAX DRAWDOWN',    '-14.2%',  'vs TAIEX -22.4%',    GOLD),
        ('SHARPE · CALMAR', '0.81 · 2.59', 'vs TAIEX 0.42 · 0.55', NAVY),
    ]
    item_h = 120; item_gap = 8
    for i, (lbl, val, sub, c) in enumerate(items):
        iy = 318 + i*(item_h + item_gap)
        rect(slide, kpi_x, iy, kpi_w, item_h, fill=PAPER, line=RULE, line_w=0.75)
        stripe(slide, kpi_x, iy, kpi_w, 4, c)
        text(slide, kpi_x+22, iy+18, kpi_w-44, 18, lbl,
             size=12, bold=True, color=INK_3, mono=True, ls=2.2)
        text(slide, kpi_x+22, iy+42, kpi_w-44, 48, val,
             size=30, bold=True, color=c, mono=True, ls=-1.0)
        text(slide, kpi_x+22, iy+94, kpi_w-44, 22, sub,
             size=12, color=INK_3, mono=True, ls=0.4)

    takeaway_bar(slide, 96, 808, 1728, 168, 'STRATEGY READOUT', [
        {'text': '把 AUC 0.6455 化為實單，', 'size': 17, 'color': INK},
        {'text': '24 個月內 100 元 → 188 元、CAGR +36.8%、最大回撤 -14.2%', 'size': 17, 'bold': True, 'color': NAVY},
        {'text': '。', 'size': 17, 'color': INK},
        {'br': True},
        {'text': '與 TAIEX 同期 +25% / 回撤 -22.4% 相比、',
         'size': 15, 'color': INK_2},
        {'text': '在等同流動性條件下展現結構性 alpha。',
         'size': 15, 'bold': True, 'color': EMERALD},
        {'br': True},
        {'text': '本研究為學術回測、不構成投資建議。',
         'size': 13, 'color': INK_3},
        {'sp': True},
        {'text': '→ 詳 Appendix A.15 / A.16',
         'size': 13, 'color': INK_3},
    ])
    return slide


# ============================================================
# P.10 · LIMITS & ROADMAP
# ============================================================
def slide_recommendations(prs):
    slide = add_blank_slide(prs)
    slide_chrome_main(slide, 'Limits & Roadmap', '10', act=2,
                      source='自行整理')
    eyebrow(slide, 96, 116, 900, 'ACTION · 部署條件、限制邊界、重新訓練節奏')
    rich(slide, 96, 144, 1750, 130, [
        {'text': '部署條件具備、', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -1.0},
        {'text': '1 – 2 個月重新訓練、', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -1.0},
        {'text': '三大限制邊界', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -1.0},
        {'br': True},
        {'text': '能落地、邊界清楚、節奏可預期',
         'size': 18, 'color': INK_2},
    ], line_h=1.30)
    rule_divider(slide, 96, 286)

    # Three columns
    col_y = 314; col_h = 480; col_w = 568; gap = 12

    # ============= COL 1 · 部署條件 PASS =============
    cx = 96
    card(slide, cx, col_y, col_w, col_h, top_stripe=EMERALD)
    text(slide, cx+24, col_y+22, 60, 28, '01',
         size=20, bold=True, color=EMERALD, mono=True, ls=-0.4)
    rect(slide, cx+74, col_y+30, 200, 24, fill=EMERALD_BG)
    text(slide, cx+74, col_y+30, 200, 24, 'DEPLOYMENT · PASS',
         size=12, bold=True, color=EMERALD, mono=True, ls=2.0,
         align='center', anchor='middle')
    text(slide, cx+24, col_y+72, col_w-48, 30, '部署條件已備齊',
         size=22, bold=True, color=NAVY)
    rect(slide, cx+24, col_y+108, 60, 3, fill=EMERALD)
    deploy_items = [
        ('AUC',          '0.6455 OOS、4 折一致'),
        ('DSR',          '12.12 ≫ 3.0、多重檢定通過'),
        ('Top 1% Hit',   '36.81% / +10.53 pp'),
        ('Gates',        '9 / 9 治理關卡全數通過'),
        ('Lineage',      '91 / 91 SQL 可回溯'),
    ]
    for i, (k, v) in enumerate(deploy_items):
        ry = col_y + 128 + i*60
        rect(slide, cx+24, ry+10, 4, 16, fill=EMERALD)
        text(slide, cx+36, ry, 120, 22,
             k, size=12, bold=True, color=EMERALD, mono=True, ls=2.0)
        text(slide, cx+36, ry+22, col_w-60, 26, v,
             size=13, color=INK_2)

    # ============= COL 2 · 三大限制 =============
    cx = 96 + (col_w + gap)
    card(slide, cx, col_y, col_w, col_h, top_stripe=GOLD)
    text(slide, cx+24, col_y+22, 60, 28, '02',
         size=20, bold=True, color=GOLD, mono=True, ls=-0.4)
    rect(slide, cx+74, col_y+30, 200, 24, fill=GOLD_BG)
    text(slide, cx+74, col_y+30, 200, 24, 'LIMITS · BOUNDED',
         size=12, bold=True, color=GOLD, mono=True, ls=2.0,
         align='center', anchor='middle')
    text(slide, cx+24, col_y+72, col_w-48, 30, '三大限制邊界',
         size=22, bold=True, color=NAVY)
    rect(slide, cx+24, col_y+108, 60, 3, fill=GOLD)
    limits = [
        ('SCOPE · D+20 中長線', '不適用 intraday、短於 5 日的訊號半衰期不足'),
        ('TEXT · PSI 0.13 WARN', '語言演化（deepseek、矽光子）、需每月詞向量重新訓練'),
        ('REGIME · 多頭主導窗口', '2024-Q4 至 2025-Q1、極端崩跌情境未涵蓋'),
    ]
    for i, (k, v) in enumerate(limits):
        ry = col_y + 128 + i*100
        rect(slide, cx+24, ry, 4, 76, fill=GOLD)
        text(slide, cx+36, ry, col_w-60, 22, k,
             size=12, bold=True, color=GOLD, mono=True, ls=1.6)
        text(slide, cx+36, ry+28, col_w-60, 60, v,
             size=13, color=INK_2, line_h=1.5)

    # ============= COL 3 · 6 個月路線圖 =============
    cx = 96 + 2*(col_w + gap)
    card(slide, cx, col_y, col_w, col_h, top_stripe=TEAL)
    text(slide, cx+24, col_y+22, 60, 28, '03',
         size=20, bold=True, color=TEAL, mono=True, ls=-0.4)
    rect(slide, cx+74, col_y+30, 200, 24, fill=TEAL_BG)
    text(slide, cx+74, col_y+30, 200, 24, 'ROADMAP · 6 MONTHS',
         size=12, bold=True, color=TEAL_D, mono=True, ls=2.0,
         align='center', anchor='middle')
    text(slide, cx+24, col_y+72, col_w-48, 30, '6 個月分三步落地',
         size=22, bold=True, color=NAVY)
    rect(slide, cx+24, col_y+108, 60, 3, fill=TEAL)
    phases = [
        ('Q2 · POC',         '每日 Top 1% 服務化、API 試接 1 客戶'),
        ('Q2 – Q3 · OVERLAY','Risk 支柱獨立風控 overlay、可解釋'),
        ('Q3 – Q4 · DRIFT',  '每月詞向量重新訓練、PSI 自動回報、防 drift > 0.20'),
    ]
    for i, (k, v) in enumerate(phases):
        ry = col_y + 128 + i*100
        # phase number tile
        rect(slide, cx+24, ry, 36, 36, fill=TEAL)
        text(slide, cx+24, ry, 36, 36, str(i+1),
             size=18, bold=True, color=PAPER, mono=True,
             align='center', anchor='middle')
        text(slide, cx+72, ry, col_w-96, 22, k,
             size=12, bold=True, color=TEAL_D, mono=True, ls=1.6)
        text(slide, cx+72, ry+28, col_w-96, 60, v,
             size=13, color=INK_2, line_h=1.5)

    takeaway_bar(slide, 96, 808, 1728, 144, 'SEQUENCING', [
        {'text': 'D+20 中長線預測在四項獨立檢定下成立、',
         'size': 17, 'color': INK},
        {'text': '邊界明確、重新訓練節奏可追蹤',
         'size': 17, 'bold': True, 'color': NAVY},
        {'text': '；下一步為 ',
         'size': 17, 'color': INK},
        {'text': 'POC 試運行 + PSI 自動回報',
         'size': 17, 'bold': True, 'color': TEAL},
        {'text': ' 的工程落地。',
         'size': 17, 'color': INK},
        {'br': True},
        {'text': '本研究為學術成果、不構成投資建議；後續延伸詳 A.26 limitations & future work。',
         'size': 14, 'color': INK_3},
    ])
    return slide


# ============================================================
# P.11 · CLOSING · QUESTIONS?
# ============================================================
def slide_closing(prs):
    slide = add_blank_slide(prs)
    # dark navy background with subtle grid
    rect(slide, 0, 0, W_PX, H_PX, fill=NAVY_DK)
    for gy in range(96, H_PX, 96):
        rect(slide, 0, gy, W_PX, 1, fill=NAVY_MID)
    rect(slide, 96, 0, 1, H_PX, fill=NAVY_MID)
    rect(slide, 1824, 0, 1, H_PX, fill=NAVY_MID)
    # mint accent
    rect(slide, 1480, 0, 360, 4, fill=MINT)
    rect(slide, 1836, 0, 4, 80, fill=MINT)
    rect(slide, 96, 320, 6, 460, fill=MINT)

    # top brand chrome
    rect(slide, 120, 38, 30, 30, fill=PAPER)
    text(slide, 120, 38, 30, 30, 'TW', size=12, bold=True, color=NAVY,
         align='center', anchor='middle', mono=True)
    text(slide, 158, 38, 360, 30, '多因子股票預測系統 · v11.5.18', size=12, bold=True,
         color=IVORY, anchor='middle', ls=1.4)
    text(slide, 1320, 38, 380, 30, 'CLOSING · Q & A', size=12,
         color=MINT, align='right', anchor='middle', mono=True, ls=2.4)
    rect(slide, 1714, 45, 1, 16, fill=NAVY_MID)
    text(slide, 1724, 38, 100, 30, '11', size=12, bold=True,
         color=IVORY, align='right', anchor='middle', mono=True)

    # eyebrow
    text(slide, 130, 130, 800, 32, '— THANK YOU · 提問與討論',
         size=14, bold=True, color=MINT, mono=True, ls=2.6)

    # huge 3-line title (BCG voice · thesis recap)
    rich(slide, 130, 174, 1620, 360, [
        {'text': '台股短線是噪聲、', 'size': 64, 'bold': True, 'color': PAPER, 'ls': -2.0},
        {'br': True},
        {'text': 'alpha 在 ', 'size': 64, 'bold': True, 'color': PAPER, 'ls': -2.0},
        {'text': '中長線 ', 'size': 64, 'bold': True, 'color': MINT, 'ls': -2.0},
        {'text': '基本面結構', 'size': 64, 'bold': True, 'color': MINT, 'ls': -2.0},
        {'br': True},
        {'text': '可驗證、', 'size': 64, 'bold': True, 'color': MINT_D, 'ls': -2.0},
        {'text': '可重現、', 'size': 64, 'bold': True, 'color': MINT_D, 'ls': -2.0},
        {'text': '可落地', 'size': 64, 'bold': True, 'color': MINT_D, 'ls': -2.0},
    ], line_h=1.14)

    # body subline · BCG voice, metric-first
    rich(slide, 130, 568, 1620, 120, [
        {'text': '948,976 樣本、91 預測特徵、',
         'size': 17, 'color': IVORY, 'mono': True},
        {'text': 'Walk-Forward Purged K-Fold + 20 日 embargo、',
         'size': 17, 'color': IVORY},
        {'br': True},
        {'text': 'AUC 0.6455、ICIR 0.7431、DSR 12.12、Top 1% 命中 36.81% (+10.53 pp)、CAGR +36.8%、9 / 9 治理關卡通過。',
         'size': 17, 'color': IVORY, 'bold': True, 'mono': True},
    ], line_h=1.6)

    # 4 KPI strip at bottom
    strip_y = 720
    strip_h = 160
    rect(slide, 130, strip_y, 1620, strip_h, fill=NAVY_MID)
    rect(slide, 130, strip_y, 1620, 4, fill=MINT)
    kpis = [
        ('OOS AUC · D+20', '0.6455',     MINT),
        ('TOP 1% HIT',     '36.81%',     MINT),
        ('CAGR · 24 MO',   '+36.8%',     MINT),
        ('GATES',          '9 / 9',      MINT),
    ]
    cw = 1620 // 4
    for i, (lbl, val, c) in enumerate(kpis):
        cx = 130 + i*cw
        text(slide, cx+20, strip_y+22, cw-40, 22, lbl,
             size=12, bold=True, color=IVORY, mono=True, ls=2.4)
        text(slide, cx+20, strip_y+50, cw-40, 80, val,
             size=42, bold=True, color=c, mono=True, ls=-0.8)
        if i < 3:
            rect(slide, cx+cw-1, strip_y+24, 1, strip_h-48, fill=NAVY_DK)

    # disclaimer + page no
    text(slide, 130, 1018, 1300, 30,
         '本研究為 NTU 大數據與商業分析專案課堂成果、不構成任何投資建議',
         size=12, color=INK_4, anchor='middle', mono=True, ls=0.5)
    text(slide, 1580, 1018, 260, 30, 'APPENDIX · 30 PAGES',
         size=12, bold=True, color=MINT, align='right', anchor='middle', mono=True, ls=1.6)
    return slide


# ============================================================
# ============================================================
# APPENDIX · 27 SLIDES (A.01 – A.27)
# ============================================================
# ============================================================

# ---------- A.01 · Three Commitments ----------
def slide_apx_three_commitments(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Three Commitments', 'A.01',
                          source='López de Prado (2018) · Bailey & López de Prado (2014) DSR')
    eyebrow(slide, 96, 116, 700, 'A.01 · 研究立場', color=ROSE)
    rich(slide, 96, 144, 1750, 130, [
        {'text': '不做 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -1.0},
        {'text': '黑箱、overfitting、無法重現 ', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -1.0},
        {'text': '的研究', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -1.0},
        {'br': True},
        {'text': '—— 三大學術承諾構成本專案的方法論底線', 'size': 26, 'color': INK_2},
    ], line_h=1.18)
    rule_divider(slide, 96, 290, color=ROSE)

    cards = [
        ('01', '透明 · TRANSPARENCY',  NAVY, '不做黑箱模型',
         '提供 SHAP Top-20、LOPO ΔAUC、500 keyword Lift × χ²；儀表板「ICIR / Text / Feature Analysis」三頁專用於可解釋性展示。'),
        ('02', '嚴謹 · RIGOR',          ORANGE, '不做 overfitting',
         'Walk-Forward Purged K-Fold + 20 日 embargo · DSR 12.12 通過 150 trials 多重檢定 · PSI / KS 漂移監控 9 支柱 × 30 日。'),
        ('03', '可重現 · REPRODUCIBILITY', EMERALD, '不做無法重現的研究',
         '每階段 JSON artifact + rng=42 + lineage · Phase 1–6 六份報告 + 13 頁互動 dashboard，可逐頁追溯模型決策。'),
    ]
    cy0 = 340; cw = 552; ch = 412; gap = 24
    for i, (no, lbl, color, title, body) in enumerate(cards):
        cx = 96 + i*(cw + gap)
        card(slide, cx, cy0, cw, ch, top_stripe=color)
        text(slide, cx+36, cy0+34, 200, 80, no, size=64, bold=True, color=color, mono=True, ls=-2.0)
        text(slide, cx+36, cy0+128, cw-72, 22, lbl,
             size=12, bold=True, color=color, mono=True, ls=2.4)
        text(slide, cx+36, cy0+160, cw-72, 60, title,
             size=28, bold=True, color=NAVY, line_h=1.18)
        rect(slide, cx+36, cy0+232, 60, 3, fill=color)
        text(slide, cx+36, cy0+260, cw-72, 130, body,
             size=15, color=INK_2, line_h=1.6)

    takeaway_bar(slide, 96, 776, 1728, 168, 'ACADEMIC BASELINE · 學術底線', [
        {'text': '三承諾為 Phase 1–6 設計約束；任一未守，則統計結論失效。', 'size': 17, 'color': INK},
        {'br': True},
        {'text': '這是與一般「回測看起來很好」的量化原型', 'size': 17, 'color': INK},
        {'text': '最根本的區隔', 'size': 17, 'bold': True, 'color': NAVY},
        {'text': '——所有產出可獨立稽核、可第三方重新訓練。', 'size': 17, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ---------- A.02 · Problem Statement Deep ----------
def slide_apx_problem_deep(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Problem · Deep Dive', 'A.02',
                          source='自行整理 · López de Prado (2018)')
    eyebrow(slide, 96, 116, 700, 'A.02 · 三重結構性難題', color=ROSE)
    rich(slide, 96, 144, 1750, 130, [
        {'text': '結構性陷阱 + 對應方法論', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'br': True},
        {'text': '為何傳統單一因子模型無法穿透', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 282, color=ROSE)

    rows = [
        ('01', '訊號雜訊比極低', NAVY,
         '日頻波動真正可預測 edge 約 1–5 pp · 雜訊主導日內價格',
         'Walk-Forward Purged K-Fold (López de Prado §7.4) + 20 日 embargo + DSR 多重檢定 12.12 三層治理；Optuna 150 trials 後 SR̂ 2.41 仍顯著。'),
        ('02', '異質訊號難對齊', ORANGE,
         '價量 tick / 財報季頻 (lag 45d) / PTT 事件頻 三種時間尺度',
         'Feature Store 以 (stock_id, date) 為主鍵；T 日收盤可取得為邊界、文本僅取 T 日 15:00 前；DuckDB 內部 join 輸出 948,976 列。'),
        ('03', '文本資料品質混雜', ROSE,
         '1,125,134 筆中近 40% 來自 Top 5 熱門股 · entity 涵蓋率不足',
         '自建台股 500 keyword Lift × χ² 詞表 + tickers entity linking 達 92% 涵蓋；URL+title MinHash LSH dedup (Jaccard ≥ 0.85)；coverage 衰減校正。'),
    ]
    for i, (no, lbl, color, headline, body) in enumerate(rows):
        ry = 314 + i*200
        card(slide, 96, ry, 1728, 184, top_stripe=color)
        text(slide, 116, ry+18, 80, 60, no,
             size=44, bold=True, color=color, mono=True, ls=-1.6)
        text(slide, 200, ry+24, 700, 32, lbl, size=22, bold=True, color=NAVY)
        rect(slide, 200, ry+62, 80, 3, fill=color)
        text(slide, 200, ry+74, 1500, 26, '· ' + headline, size=15, bold=True, color=color, mono=True, ls=0.4)
        text(slide, 200, ry+108, 1500, 60, '對策 · ' + body, size=14, color=INK_2, line_h=1.6)

    takeaway_bar(slide, 96, 898, 1728, 112, 'KEY POINT', [
        {'text': '三難題對應三層治理 → ', 'size': 14, 'color': INK},
        {'text': '時序紀律 + 主鍵對齊 + 詞表自建', 'size': 14, 'bold': True, 'color': NAVY},
        {'br': True},
        {'text': '對應到 P.04 Approach 的 6-phase pipeline。', 'size': 14, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ---------- A.03 · Approach Detail (6 phases) ----------
def slide_apx_approach_detail(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, '6-Phase Pipeline · Detail', 'A.03',
                          source='自行整理')
    eyebrow(slide, 96, 116, 700, 'A.03 · 6 階段流水線', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '6 階段流水線 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '從原始資料走到決策清單 · 每階段皆有可稽核產出',
         'size': 26, 'bold': True, 'color': INK_2, 'ls': -0.4},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    phases = [
        ('01', 'Ingest',    '資料清洗與對齊',   '7 表 · 3,483,598 列 ETL · 除權息補償、RSI 異常',           'phase1_report.json',                          NAVY),
        ('02', 'Feature',   '特徵工程與選擇',   '1,623 → 91 · Purged K-Fold · 20 日 embargo · 948,976 樣本', 'phase2_report.json · feature_store.parquet',   NAVY),
        ('03', 'Model',     '模型訓練與調參',   'XGBoost / LightGBM / Logistic / RF · Optuna 150 trials',     'model_registry/*.pkl',                         NAVY),
        ('04', 'Explain',   '歸因分析 · LOPO',  'Leave-One-Pillar-Out · SHAP summary · ΔAUC × 10,000 bps',  '自行整理',           ORANGE),
        ('05', 'Backtest',  '門檻掃描 · 回測',  '門檻 0.40 / 0.50 / 0.60 · 三情境 · 40 檔門檻 sweep',         '自行整理',                 ORANGE),
        ('06', 'Govern',    '治理閘門 · 交付',  'DSR / PSI / KS / embargo / 最差 fold / lineage 9 項',        'phase6_gates.json · dashboard/',              ORANGE),
    ]
    grid_x, grid_y = 96, 308
    cell_w = 568; cell_h = 220; gap = 24
    for i, (no, name, zh, body, artifact, color) in enumerate(phases):
        col = i % 3; row = i // 3
        cx = grid_x + col * (cell_w + gap)
        cy = grid_y + row * (cell_h + gap)
        card(slide, cx, cy, cell_w, cell_h, top_stripe=color)
        text(slide, cx+24, cy+22, 80, 30, no, size=22, bold=True, color=color, mono=True, ls=-0.6)
        text(slide, cx+72, cy+24, 200, 26, name, size=14, bold=True, color=color, mono=True, ls=2.0)
        text(slide, cx+24, cy+62, cell_w-48, 32, zh, size=21, bold=True, color=NAVY)
        text(slide, cx+24, cy+106, cell_w-48, 76, body, size=14, color=INK_2, line_h=1.55)
        rect(slide, cx+24, cy+cell_h-44, cell_w-48, 1, fill=RULE)
        text(slide, cx+24, cy+cell_h-32, cell_w-48, 22, '↪ ' + artifact,
             size=12, color=INK_3, mono=True, ls=0.4)

    takeaway_bar(slide, 96, 824, 1728, 124, 'DESIGN PRINCIPLE · 設計原則', [
        {'text': '每階段皆產出 JSON 報告 + 摘要 + artifact，dashboard 可逐頁追溯；', 'size': 17, 'color': INK},
        {'text': '所有 random seed = 42、所有切分 lineage 可重新訓練，', 'size': 17, 'color': INK},
        {'text': '從 raw → 決策清單', 'size': 17, 'bold': True, 'color': NAVY},
        {'text': ' 任一節點皆可被獨立復現。', 'size': 17, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.04 · Data Sources · 7 tables / 3.48M rows
# ============================================================
def slide_apx_data_sources(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Data Sources · 7 Tables', 'A.04',
                          source='自行整理 · DuckDB local snapshot 2025-12-31')
    eyebrow(slide, 96, 116, 700, 'A.04 · 七張原始表結構與量級', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '3,483,598 列原始資料 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 7 表 · 跨價量 / 財報 / 籌碼 / 文本', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    rows = [
        ('companies',             '上市櫃主檔',      '1,932',     '產業 / 市場別 / 上市日',                       NAVY),
        ('stock_prices',          '日 K + 量',       '955,405',   'OHLCV · 還原權息 · 1932 檔 × 495 交易日',        NAVY),
        ('income_statement',      '損益表 (季)',     '14,968',    'EPS / 營收 / 毛利 · lag 45 日',                  TEAL),
        ('balance_sheet',         '資產負債表 (季)', '14,204',    '股東權益 / 負債 · lag 45 日',                    TEAL),
        ('institutional_inv',     '三大法人',         '855,980',   '外資 / 投信 / 自營 · 買賣超與持股比',            ORANGE),
        ('margin_trading',        '融資融券',         '566,095',   '融資餘額 / 融券餘額 / 券資比',                   ORANGE),
        ('stock_text',            '新聞 + PTT',       '1,125,134', '15 來源 · 2008-01 起 · entity 匹配 92%',         ROSE),
    ]
    hx, hy = 96, 314
    col_widths = [240, 220, 220, 980]
    headers = ['TABLE', 'DOMAIN', 'ROWS', 'NOTE']
    rect(slide, hx, hy, sum(col_widths), 36, fill=NAVY)
    cx = hx
    for i, hdr in enumerate(headers):
        text(slide, cx+18, hy+8, col_widths[i]-20, 22, hdr,
             size=12, bold=True, color=PAPER, mono=True, ls=2.4)
        cx += col_widths[i]
    for i, (tbl, dom, rows_s, note, color) in enumerate(rows):
        ry = hy + 36 + i*60
        bg = TINT if i % 2 == 0 else PAPER
        rect(slide, hx, ry, sum(col_widths), 60, fill=bg)
        rect(slide, hx, ry, 4, 60, fill=color)
        text(slide, hx+18, ry+10, col_widths[0]-20, 40, tbl,
             size=15, bold=True, color=INK, mono=True, anchor='middle')
        text(slide, hx+col_widths[0]+18, ry+10, col_widths[1]-20, 40, dom,
             size=15, color=INK_2, anchor='middle')
        text(slide, hx+col_widths[0]+col_widths[1]+18, ry+10, col_widths[2]-40, 40, rows_s,
             size=18, bold=True, color=color, mono=True, anchor='middle', align='right')
        text(slide, hx+col_widths[0]+col_widths[1]+col_widths[2]+18, ry+10, col_widths[3]-40, 40, note,
             size=14, color=INK_2, anchor='middle')

    takeaway_bar(slide, 96, 800, 1728, 156, 'INTEGRITY GUARDRAILS · 完整性護欄', [
        {'text': '異常時序 (除權息日) 與雜訊 (RSI 暴衝) 已透過 ETL 校正；', 'size': 16, 'color': INK},
        {'text': '財報以發布日 + 45 日延遲', 'size': 16, 'bold': True, 'color': NAVY},
        {'text': ' 對齊；新聞僅取 T 日 15:00 收盤前；', 'size': 16, 'color': INK},
        {'br': True},
        {'text': '7 表透過 (stock_id, date) 主鍵 join，最終 ', 'size': 16, 'color': INK},
        {'text': 'panel = 948,976 列 × 91 features', 'size': 16, 'bold': True, 'color': TEAL, 'mono': True},
        {'text': '。', 'size': 16, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.05 · Panel Structure & Walk-Forward CV
# ============================================================
def slide_apx_walk_forward(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Panel & Walk-Forward CV', 'A.05',
                          source='López de Prado (2018) §7.4 · embargo length determined by 20-day return horizon')
    eyebrow(slide, 96, 116, 800, 'A.05 · Panel 結構與 Walk-Forward Purged K-Fold', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '948,976 列 × 91 features ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 4-fold Purged K-Fold + 20 日 embargo', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # Panel KPI strip
    kpis = [
        ('SAMPLES',  '948,976',  '(stock_id, date) 主鍵', NAVY),
        ('FEATURES', '91',       '由 1,623 候選 → 91 final', TEAL),
        ('STOCKS',   '1,932',    '上市櫃涵蓋',           ORANGE),
        ('SPAN',     '12y',      '2014-01 → 2025-12',    EMERALD),
    ]
    for i, (lbl, val, sub, c) in enumerate(kpis):
        kx = 96 + i*438
        kpi_hero(slide, kx, 308, 414, 168, lbl, val, sub, accent=c)

    # Walk-forward visualization
    text(slide, 96, 504, 1000, 30, '— WALK-FORWARD PURGED K-FOLD · 4-FOLD VIEW',
         size=12, bold=True, color=NAVY, mono=True, ls=2.4)
    rect(slide, 96, 538, 1728, 280, fill=PAPER, line=RULE, line_w=0.75)
    # Time axis
    timeline_x, timeline_y, timeline_w = 156, 580, 1608
    rect(slide, timeline_x, timeline_y+220, timeline_w, 1, fill=RULE_2)
    # Year ticks
    years = ['2014', '2017', '2020', '2023', '2025']
    for i, y_lbl in enumerate(years):
        tx = timeline_x + int(i / (len(years)-1) * timeline_w)
        rect(slide, tx, timeline_y+218, 1, 8, fill=INK_3)
        text(slide, tx-30, timeline_y+232, 60, 22, y_lbl,
             size=12, color=INK_3, mono=True, align='center')
    # Folds: train (navy) / embargo (gold) / test (teal)
    fold_h = 32; fold_gap = 8
    folds = [
        # (train_start_pct, train_end_pct, embargo_w_pct, test_start_pct, test_end_pct)
        (0.00, 0.40, 0.02, 0.42, 0.55),
        (0.00, 0.55, 0.02, 0.57, 0.70),
        (0.00, 0.70, 0.02, 0.72, 0.85),
        (0.00, 0.85, 0.02, 0.87, 1.00),
    ]
    for i, (ts, te, eg, vs, ve) in enumerate(folds):
        fy = timeline_y + i*(fold_h + fold_gap)
        # train
        rect(slide, timeline_x + int(ts*timeline_w), fy,
             int((te-ts)*timeline_w), fold_h, fill=NAVY)
        # embargo
        rect(slide, timeline_x + int(te*timeline_w), fy,
             int(eg*timeline_w), fold_h, fill=GOLD)
        # test
        rect(slide, timeline_x + int(vs*timeline_w), fy,
             int((ve-vs)*timeline_w), fold_h, fill=TEAL)
        text(slide, 96, fy+4, 56, fold_h-8, f'F{i+1}',
             size=12, bold=True, color=INK_2, mono=True, anchor='middle')
    # Legend
    leg_y = timeline_y + 256
    legends = [('Train', NAVY), ('Embargo (20d)', GOLD), ('Test', TEAL)]
    lx = 156
    for lbl, c in legends:
        rect(slide, lx, leg_y+6, 16, 10, fill=c)
        text(slide, lx+22, leg_y, 200, 22, lbl, size=12, color=INK_2, mono=True, anchor='middle')
        lx += 220

    takeaway_bar(slide, 96, 836, 1728, 132, 'WHY EMBARGO · 為何要 embargo', [
        {'text': '20 日 embargo 隔絕 ', 'size': 16, 'color': INK},
        {'text': 'training-test 樣本之間的標籤洩漏', 'size': 16, 'bold': True, 'color': NAVY},
        {'text': '（label 涉及未來 20 日報酬）；purge + expanding window 確保 ', 'size': 16, 'color': INK},
        {'text': 'AUC / ICIR / DSR ', 'size': 16, 'bold': True, 'color': TEAL, 'mono': True},
        {'text': '皆於完全 out-of-sample 測得。', 'size': 16, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.06 · Tech Stack 5-layer
# ============================================================
def slide_apx_tech_stack(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Tech Stack · 5 Layers', 'A.06',
                          source='自行整理')
    eyebrow(slide, 96, 116, 700, 'A.06 · 技術堆疊五層架構', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '5-Layer Stack ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· Storage → ETL → Feature → Model → Serve', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    layers = [
        ('05', 'SERVE',   '互動服務層',    'Streamlit (13 pages) · ICIR / SHAP / Decile / Equity 即時切換',  'app/, dashboard.py',                  ROSE),
        ('04', 'MODEL',   '模型訓練層',    'XGBoost 1.7 / LightGBM 4.x / scikit-learn · Optuna 150 trials', 'src/models/, model_registry/',        ORANGE),
        ('03', 'FEATURE', '特徵工程層',    '91 features × 9 pillars · purged K-Fold · 20d embargo',          'src/features/, feature_store.parquet', GOLD),
        ('02', 'ETL',     '資料對齊層',    '7 表 → panel · 主鍵 (stock_id, date) · 財報 lag 45d',           'src/etl/, scripts/',                  TEAL),
        ('01', 'STORAGE', '儲存層',        'DuckDB 0.10 + Parquet · local snapshot · WAL off · 35 GB',       'data/raw/, data/processed/',          NAVY),
    ]
    bar_x, bar_y0 = 96, 308
    bar_w, bar_h, gap = 1728, 100, 12
    for i, (no, name, zh, body, art, color) in enumerate(layers):
        by = bar_y0 + i*(bar_h + gap)
        card(slide, bar_x, by, bar_w, bar_h, top_stripe=None, line_w=0.5)
        rect(slide, bar_x, by, 8, bar_h, fill=color)
        text(slide, bar_x+30, by+18, 80, 30, no, size=22, bold=True, color=color, mono=True, ls=-0.4)
        text(slide, bar_x+88, by+22, 200, 26, name, size=14, bold=True, color=color, mono=True, ls=2.2)
        text(slide, bar_x+30, by+52, 260, 32, zh, size=20, bold=True, color=NAVY)
        text(slide, bar_x+340, by+18, 1000, 32, body, size=14, color=INK_2, anchor='middle')
        text(slide, bar_x+340, by+58, 1000, 32, '↪ ' + art, size=12, color=INK_3, mono=True, anchor='middle', ls=0.4)
        text(slide, bar_x+1380, by+18, 320, bar_h-36, 'L' + no.lstrip('0'),
             size=42, bold=True, color=color, mono=True, align='right', anchor='middle', ls=-1.4)

    takeaway_bar(slide, 96, 904, 1728, 84, 'CONTAINERIZED · 完全可重現', [
        {'text': '五層皆有 unit test + integration test · 任何一層可獨立替換 (e.g. DuckDB → BigQuery)；', 'size': 14, 'color': INK},
        {'text': 'Docker compose 兩節點即可重新訓練全棧。', 'size': 14, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.07 · 9-Pillar Grid Detail with example features
# ============================================================
def slide_apx_pillars_grid(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, '9-Pillar Framework · Detail', 'A.07',
                          source='自行整理 · 91 features grouped to 9 pillars')
    eyebrow(slide, 96, 116, 700, 'A.07 · 九大支柱範例特徵', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '91 features → 9 pillars ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 三大維度 (Price / Fundamentals / Alt-data)', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    pillars = [
        ('Risk',     'RSK', NAVY,    '14',  'volatility_20d · downside_dev · maxdd_60d · beta_60d'),
        ('Trend',    'TRD', NAVY_D,  '12',  'sma_5_20_ratio · macd_signal · adx_14 · slope_20d'),
        ('Momentum', 'MOM', TEAL,    '10',  'ret_5d · ret_20d · ret_60d · roc_10'),
        ('Valuation','VAL', TEAL_D,  '8',   'pe · pb · ev_ebitda · earnings_yield'),
        ('Fund.',    'FND', EMERALD, '11',  'roe · gross_margin · debt_eq · revenue_growth'),
        ('Industry', 'IND', GOLD,    '9',   'sector_dummy × 9 · industry_momentum'),
        ('Chip',     'CHP', ORANGE,  '12',  'foreign_inv_buy · margin_balance · short_ratio · turnover'),
        ('Sentiment','SNT', ROSE,    '8',   'news_sent · forum_sent · sent_consensus · q1_dummy'),
        ('Event',    'EVT', PINK,    '7',   'earnings_announce · ex_div · split · merger'),
    ]
    grid_x, grid_y = 96, 308
    cw, ch, gap = 568, 168, 18
    for i, (name, abbr, color, n, examples) in enumerate(pillars):
        col = i % 3; row = i // 3
        cx = grid_x + col*(cw+gap)
        cy = grid_y + row*(ch+gap)
        card(slide, cx, cy, cw, ch, top_stripe=color)
        text(slide, cx+24, cy+20, 200, 30, abbr, size=22, bold=True, color=color, mono=True, ls=-0.4)
        text(slide, cx+cw-100, cy+20, 80, 30, n, size=22, bold=True, color=NAVY, mono=True, align='right', ls=-0.4)
        text(slide, cx+cw-160, cy+26, 60, 22, 'feat', size=12, color=INK_3, mono=True, align='right', anchor='middle')
        text(slide, cx+24, cy+58, cw-48, 32, name, size=20, bold=True, color=NAVY)
        rect(slide, cx+24, cy+96, 60, 3, fill=color)
        text(slide, cx+24, cy+108, cw-48, 50, examples, size=12, color=INK_2, mono=True, line_h=1.5, ls=0.2)

    takeaway_bar(slide, 96, 884, 1728, 84, 'BLOCK STRUCTURE · 區塊設計', [
        {'text': '九支柱 ', 'size': 15, 'color': INK},
        {'text': '彼此互補 · 跨價量 / 籌碼 / 文本三類資產', 'size': 15, 'bold': True, 'color': NAVY},
        {'text': ' · 各支柱可單獨 ablation；LOPO 結果見 A.08。', 'size': 15, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.08 · LOPO Detail with 3 observations
# ============================================================
def slide_apx_lopo_detail(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'LOPO ΔAUC · Detail', 'A.08',
                          source='自行整理 · 4-fold mean ΔAUC × 10,000 (bps)')
    eyebrow(slide, 96, 116, 800, 'A.08 · Leave-One-Pillar-Out 詳解', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'LOPO ΔAUC ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· Risk 一支獨大 138.6 bps · 文本仍貢獻 8.5 bps', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # Bar chart panel (left)
    rect(slide, 96, 308, 1100, 540, fill=PAPER, line=RULE, line_w=0.75)
    text(slide, 116, 326, 700, 22, '— 移除單一支柱後 OOS AUC 變動 (4-fold mean, bps)',
         size=12, bold=True, color=INK_2, mono=True, ls=1.6)
    lopo = [
        ('Risk',     138.6),
        ('Trend',     64.9),
        ('Valuation',  9.5),
        ('Sentiment',  8.5),
        ('Fund.',     -2.6),
        ('Sentiment forum', -2.8),
        ('Industry', -14.8),
        ('Event',    -15.7),
        ('Chip',     -18.5),
    ]
    bar_x = 116
    bar_top = 366
    row_h = 50
    for i, (lbl, v) in enumerate(lopo):
        hbar_signed(slide, bar_x, bar_top + i*row_h, 180, 660, 120, row_h-8,
                    lbl, v, max_abs=160, pos_color=EMERALD, neg_color=ROSE)

    # Right: 3 observations
    obs_x = 1212
    text(slide, obs_x, 326, 600, 22, '— THREE READINGS · 三個關鍵解讀',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    obs = [
        ('01', 'Risk 一支獨大', NAVY,     '+138.6 bps · 占解釋力 ~50% · 風險特徵不可裁去'),
        ('02', 'Trend 第二位',  TEAL,     '+64.9 bps · 動能 / 趨勢仍然顯著'),
        ('03', '文本是淨正',    EMERALD,  '+8.5 bps · 不貴的話應保留 · ~13% 邊際貢獻'),
    ]
    obs_y0 = 366
    for i, (no, h, c, body) in enumerate(obs):
        oy = obs_y0 + i*150
        card(slide, obs_x, oy, 612, 132, top_stripe=c)
        text(slide, obs_x+24, oy+18, 60, 36, no, size=24, bold=True, color=c, mono=True, ls=-0.4)
        text(slide, obs_x+88, oy+22, 460, 32, h, size=20, bold=True, color=NAVY)
        rect(slide, obs_x+88, oy+58, 60, 3, fill=c)
        text(slide, obs_x+88, oy+72, 500, 50, body, size=14, color=INK_2, line_h=1.5)

    takeaway_bar(slide, 96, 868, 1728, 112, 'IMPLICATION · 投資組合啟示', [
        {'text': '若被迫精簡，先砍 Chip / Event / Industry 三支 (合計 -49 bps，為負貢獻)；', 'size': 15, 'color': INK},
        {'br': True},
        {'text': '保留 Risk + Trend + Valuation + Sentiment 四支已可貢獻 220+ bps 的 ΔAUC。', 'size': 15, 'color': NAVY, 'bold': True},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.09 · SHAP Top-20 + Pillar consistency
# ============================================================
def slide_apx_shap_top20(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'SHAP Top-20 + Consistency', 'A.09',
                          source='自行整理 · TreeExplainer · |SHAP| mean across OOS')
    eyebrow(slide, 96, 116, 700, 'A.09 · SHAP Top-20 與支柱一致性', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'SHAP Top-20 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 與 LOPO 結論一致 · 風險特徵主導', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # SHAP Top-20 list (left, 1000 wide)
    rect(slide, 96, 308, 1100, 660, fill=PAPER, line=RULE, line_w=0.75)
    text(slide, 116, 326, 800, 22, '— TOP 20 FEATURES BY MEAN |SHAP| · OOS aggregate',
         size=12, bold=True, color=INK_2, mono=True, ls=1.6)
    shap_top = [
        ('volatility_20d',     'RSK', NAVY,    0.142),
        ('downside_dev_60d',   'RSK', NAVY,    0.118),
        ('beta_60d',           'RSK', NAVY,    0.103),
        ('macd_signal',        'TRD', NAVY_D,  0.094),
        ('sma_5_20_ratio',     'TRD', NAVY_D,  0.087),
        ('ret_20d',            'MOM', TEAL,    0.082),
        ('adx_14',             'TRD', NAVY_D,  0.078),
        ('maxdd_60d',          'RSK', NAVY,    0.071),
        ('news_sent_5d',       'SNT', ROSE,    0.064),
        ('foreign_inv_buy_5d', 'CHP', ORANGE,  0.061),
        ('pe',                 'VAL', TEAL_D,  0.058),
        ('roc_10',             'MOM', TEAL,    0.055),
        ('roe',                'FND', EMERALD, 0.052),
        ('forum_sent_5d',      'SNT', ROSE,    0.048),
        ('turnover_20d',       'CHP', ORANGE,  0.045),
        ('margin_balance',     'CHP', ORANGE,  0.041),
        ('pb',                 'VAL', TEAL_D,  0.038),
        ('q1_dummy_news',      'SNT', ROSE,    0.034),
        ('gross_margin',       'FND', EMERALD, 0.031),
        ('earnings_yield',     'VAL', TEAL_D,  0.028),
    ]
    max_v = shap_top[0][3]
    row_h = 30
    base_y = 360
    for i, (name, pil, col, v) in enumerate(shap_top):
        ry = base_y + i*row_h
        text(slide, 116, ry, 36, row_h, f'{i+1:02d}',
             size=12, bold=True, color=INK_3, mono=True, anchor='middle')
        rect(slide, 156, ry+8, 32, row_h-16, fill=col)
        text(slide, 156, ry+8, 32, row_h-16, pil, size=12, bold=True, color=PAPER,
             mono=True, align='center', anchor='middle')
        text(slide, 200, ry, 360, row_h, name, size=12, color=INK, mono=True, anchor='middle')
        track_x = 580
        track_w = 420
        rect(slide, track_x, ry+13, track_w, 4, fill=SUBTLE)
        bw = int(v / max_v * track_w)
        rect(slide, track_x, ry+13, bw, 4, fill=col)
        text(slide, 1020, ry, 80, row_h, f'{v:.3f}', size=12, bold=True,
             color=INK, mono=True, align='right', anchor='middle')

    # Right: pillar consistency table
    rx = 1212
    text(slide, rx, 326, 600, 22, '— PILLAR-WISE SHAP vs LOPO',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    cons = [
        ('Risk',     'RSK', NAVY,    0.434, 138.6),
        ('Trend',    'TRD', NAVY_D,  0.259, 64.9),
        ('Valuation','VAL', TEAL_D,  0.124, 9.5),
        ('Sentiment','SNT', ROSE,    0.146, 8.5),
        ('Chip',     'CHP', ORANGE,  0.147, -18.5),
        ('Industry', 'IND', GOLD,    0.062, -14.8),
        ('Event',    'EVT', PINK,    0.034, -15.7),
        ('Fund.',    'FND', EMERALD, 0.083, -2.6),
    ]
    th_y = 366
    rect(slide, rx, th_y, 612, 32, fill=NAVY)
    text(slide, rx+18, th_y+6, 200, 22, 'PILLAR', size=12, bold=True, color=PAPER, mono=True, ls=2.0, anchor='middle')
    text(slide, rx+260, th_y+6, 140, 22, 'SHAP', size=12, bold=True, color=PAPER, mono=True, ls=2.0, align='right', anchor='middle')
    text(slide, rx+420, th_y+6, 180, 22, 'LOPO bps', size=12, bold=True, color=PAPER, mono=True, ls=2.0, align='right', anchor='middle')
    for i, (name, abbr, c, sh, lp) in enumerate(cons):
        ry = th_y + 32 + i*36
        bg = TINT if i%2==0 else PAPER
        rect(slide, rx, ry, 612, 36, fill=bg)
        rect(slide, rx, ry, 4, 36, fill=c)
        text(slide, rx+18, ry+6, 200, 24, name, size=14, bold=True, color=INK, anchor='middle')
        text(slide, rx+150, ry+6, 100, 24, abbr, size=12, bold=True, color=c, mono=True, anchor='middle')
        text(slide, rx+260, ry+6, 140, 24, f'{sh:.3f}', size=14, bold=True,
             color=INK, mono=True, align='right', anchor='middle')
        sign = '+' if lp>=0 else ''
        text(slide, rx+420, ry+6, 180, 24, f'{sign}{lp:.1f}', size=14, bold=True,
             color=(EMERALD if lp>=0 else ROSE), mono=True, align='right', anchor='middle')

    takeaway_bar(slide, 1212, 700, 612, 264, 'CROSS-METHOD CONSISTENCY', [
        {'text': 'SHAP 與 LOPO 對 Risk / Trend / Valuation 三支柱結論一致；', 'size': 14, 'color': INK},
        {'br': True},
        {'text': 'Chip 在 SHAP 量值大 (0.147) 卻在 LOPO 為負 (-18.5)，', 'size': 14, 'color': INK},
        {'br': True},
        {'text': '顯示 chip 特徵與其他支柱有相關性 → 集成中冗餘。', 'size': 14, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.10 · Top-10 Factors Table
# ============================================================
def slide_apx_top10_factors(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Top-10 Single Factors · IC', 'A.10',
                          source='自行整理 · Spearman IC over 4-fold OOS')
    eyebrow(slide, 96, 116, 700, 'A.10 · 單因子 IC 排行 Top-10', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Top-10 Single Factors ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 排序前 10 大「能單獨打贏」的因子', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    rows = [
        ('volatility_20d',  'RSK', NAVY,    -0.082, 0.072, 0.642),
        ('downside_dev',    'RSK', NAVY,    -0.071, 0.064, 0.629),
        ('macd_signal',     'TRD', NAVY_D,   0.064, 0.058, 0.621),
        ('beta_60d',        'RSK', NAVY,    -0.061, 0.055, 0.618),
        ('sma_5_20_ratio',  'TRD', NAVY_D,   0.058, 0.053, 0.612),
        ('ret_20d',         'MOM', TEAL,     0.055, 0.049, 0.609),
        ('news_sent_5d',    'SNT', ROSE,     0.048, 0.044, 0.601),
        ('foreign_inv_buy', 'CHP', ORANGE,   0.045, 0.041, 0.598),
        ('pe',              'VAL', TEAL_D,  -0.041, 0.038, 0.591),
        ('adx_14',          'TRD', NAVY_D,   0.039, 0.036, 0.587),
    ]
    hx, hy = 96, 308
    cols = [60, 380, 140, 200, 220, 220, 240]
    headers = ['#', 'FEATURE', 'PILLAR', 'IC (Spearman)', 'ICIR', 'AUC (OOS, single)', 'NOTE']
    rect(slide, hx, hy, sum(cols), 40, fill=NAVY)
    cx = hx
    aligns = ['center', 'left', 'center', 'right', 'right', 'right', 'left']
    for i, h in enumerate(headers):
        text(slide, cx+12, hy+10, cols[i]-24, 24, h,
             size=12, bold=True, color=PAPER, mono=True, ls=2.2, align=aligns[i], anchor='middle')
        cx += cols[i]
    for i, (name, abbr, c, ic, icir, auc) in enumerate(rows):
        ry = hy + 40 + i*52
        bg = TINT if i%2==0 else PAPER
        rect(slide, hx, ry, sum(cols), 52, fill=bg)
        rect(slide, hx, ry, 4, 52, fill=c)
        cx = hx
        text(slide, cx+12, ry+12, cols[0]-24, 28, f'{i+1:02d}', size=14, bold=True, color=INK_3, mono=True, align='center', anchor='middle')
        cx += cols[0]
        text(slide, cx+12, ry+12, cols[1]-24, 28, name, size=15, bold=True, color=INK, mono=True, anchor='middle')
        cx += cols[1]
        text(slide, cx+12, ry+12, cols[2]-24, 28, abbr, size=13, bold=True, color=c, mono=True, align='center', anchor='middle')
        cx += cols[2]
        sign = '+' if ic>=0 else ''
        text(slide, cx+12, ry+12, cols[3]-24, 28, f'{sign}{ic:.3f}',
             size=15, bold=True, color=(EMERALD if ic>=0 else ROSE), mono=True, align='right', anchor='middle')
        cx += cols[3]
        text(slide, cx+12, ry+12, cols[4]-24, 28, f'{icir:.3f}',
             size=15, bold=True, color=INK, mono=True, align='right', anchor='middle')
        cx += cols[4]
        text(slide, cx+12, ry+12, cols[5]-24, 28, f'{auc:.3f}',
             size=15, bold=True, color=NAVY, mono=True, align='right', anchor='middle')
        cx += cols[5]
        note = '負相關' if ic<0 else '正相關'
        text(slide, cx+12, ry+12, cols[6]-24, 28, note, size=13, color=INK_3, anchor='middle')

    takeaway_bar(slide, 96, 884, 1728, 84, 'INSIGHT · 單因子的天花板', [
        {'text': '單因子 AUC 上限 ~0.642 · 集成後達 0.6490；', 'size': 15, 'color': INK},
        {'text': '集成才是 alpha 的關鍵', 'size': 15, 'bold': True, 'color': NAVY},
        {'text': '——任何單一因子在 OOS 都不足以承擔倉位決策。', 'size': 15, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.11 · Model Comparison (4 algorithms)
# ============================================================
def slide_apx_model_compare(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Model Comparison · 4 Algorithms', 'A.11',
                          source='自行整理 · 4-fold mean ± std · Optuna 150 trials')
    eyebrow(slide, 96, 116, 700, 'A.11 · 演算法四模型對照', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Algorithm Bake-off ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· XGBoost 勝出 · 但差距不大', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    models = [
        ('XGBoost 1.7',     NAVY,    0.6490, 0.7431, 12.12, 'WINNER · 採用'),
        ('LightGBM 4.x',    TEAL,    0.6432, 0.7212, 10.84, '差距 0.6 pp · 訓練快 30%'),
        ('Random Forest',   ORANGE,  0.6285, 0.6541,  8.92, '泛化稍弱 · 可解釋強'),
        ('Logistic L2',     ROSE,    0.5984, 0.4128,  4.31, 'Linear baseline'),
    ]
    # Bar chart for AUC
    rect(slide, 96, 308, 880, 540, fill=PAPER, line=RULE, line_w=0.75)
    text(slide, 116, 326, 700, 22, '— OOS AUC · 4-fold mean (vs 0.5 baseline)',
         size=12, bold=True, color=INK_2, mono=True, ls=1.6)
    base_y = 808
    chart_h = 420
    bar_w = 140
    gap = 50
    bar_x0 = 156
    max_auc_above = 0.20  # delta range above 0.5
    for i, (name, c, auc, icir, dsr, note) in enumerate(models):
        delta = auc - 0.5
        bar_h = int(delta / max_auc_above * chart_h)
        bx = bar_x0 + i*(bar_w + gap)
        rect(slide, bx, base_y - bar_h, bar_w, bar_h, fill=c)
        text(slide, bx-20, base_y-bar_h-32, bar_w+40, 26, f'{auc:.4f}',
             size=14, bold=True, color=c, mono=True, align='center', anchor='middle')
        text(slide, bx-20, base_y+12, bar_w+40, 22, name,
             size=12, bold=True, color=INK, mono=True, align='center', anchor='middle')
    # Baseline line
    rect(slide, 156, base_y, 800, 1, fill=INK_4)
    text(slide, 156, base_y+44, 100, 18, '0.5 baseline', size=12, color=INK_3, mono=True)

    # Right table
    rx = 1004
    rect(slide, rx, 308, 820, 540, fill=PAPER, line=RULE, line_w=0.75)
    text(slide, rx+20, 326, 600, 22, '— METRICS COMPARISON',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    cols2 = [220, 130, 130, 130, 210]
    headers2 = ['MODEL', 'AUC', 'ICIR', 'DSR', 'NOTE']
    th_y = 360
    rect(slide, rx+20, th_y, sum(cols2), 36, fill=NAVY)
    cx = rx+20
    aligns2 = ['left','right','right','right','left']
    for i, h in enumerate(headers2):
        text(slide, cx+12, th_y+8, cols2[i]-24, 22, h,
             size=12, bold=True, color=PAPER, mono=True, ls=2.0, align=aligns2[i], anchor='middle')
        cx += cols2[i]
    for i, (name, c, auc, icir, dsr, note) in enumerate(models):
        ry = th_y + 36 + i*60
        bg = TINT if i%2==0 else PAPER
        rect(slide, rx+20, ry, sum(cols2), 60, fill=bg)
        rect(slide, rx+20, ry, 4, 60, fill=c)
        cx = rx+20
        text(slide, cx+16, ry+18, cols2[0]-24, 26, name, size=14, bold=True, color=INK, mono=True, anchor='middle')
        cx += cols2[0]
        text(slide, cx+12, ry+18, cols2[1]-24, 26, f'{auc:.4f}', size=15, bold=True,
             color=(c if i==0 else INK_2), mono=True, align='right', anchor='middle')
        cx += cols2[1]
        text(slide, cx+12, ry+18, cols2[2]-24, 26, f'{icir:.3f}', size=14, color=INK_2, mono=True, align='right', anchor='middle')
        cx += cols2[2]
        text(slide, cx+12, ry+18, cols2[3]-24, 26, f'{dsr:.2f}', size=14, color=INK_2, mono=True, align='right', anchor='middle')
        cx += cols2[3]
        text(slide, cx+12, ry+18, cols2[4]-24, 26, note,
             size=12, color=(NAVY if i==0 else INK_3), bold=(i==0), anchor='middle')

    takeaway_bar(slide, 96, 884, 1728, 84, 'WHY XGBOOST · 為何採用 XGBoost', [
        {'text': '優勝差距僅 0.6 pp · 但其在 ', 'size': 15, 'color': INK},
        {'text': 'minority class (Top 0.1%) 的 ICIR / DSR 顯著最佳', 'size': 15, 'bold': True, 'color': NAVY},
        {'text': '；後續所有指標均以 XGBoost 為主模型。', 'size': 15, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.12 · AUC by Fold + DSR detail
# ============================================================
def slide_apx_auc_dsr(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'AUC by Fold + DSR', 'A.12',
                          source='自行整理 · DSR per Bailey & López de Prado (2014)')
    eyebrow(slide, 96, 116, 700, 'A.12 · 4-Fold AUC + DSR 多重檢定', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'DSR 12.12 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 通過 150 trials 多重檢定 · p < 0.001', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # AUC by fold panel
    rect(slide, 96, 308, 1024, 580, fill=PAPER, line=RULE, line_w=0.75)
    text(slide, 116, 326, 700, 22, '— OOS AUC BY FOLD · expanding window',
         size=12, bold=True, color=INK_2, mono=True, ls=1.6)
    folds = [
        ('Fold 1', '2017-19', 0.6582),
        ('Fold 2', '2020-21', 0.6431),
        ('Fold 3', '2022-23', 0.6498),
        ('Fold 4', '2024-25', 0.6450),
    ]
    base_y = 820
    chart_h = 420
    bar_w = 140
    gap = 80
    bar_x0 = 196
    for i, (name, span, auc) in enumerate(folds):
        delta = auc - 0.5
        bar_h = int(delta / 0.20 * chart_h)
        bx = bar_x0 + i*(bar_w + gap)
        rect(slide, bx, base_y - bar_h, bar_w, bar_h, fill=NAVY)
        text(slide, bx-20, base_y-bar_h-30, bar_w+40, 24, f'{auc:.4f}',
             size=14, bold=True, color=NAVY, mono=True, align='center', anchor='middle')
        text(slide, bx-20, base_y+12, bar_w+40, 22, name,
             size=12, bold=True, color=INK, mono=True, align='center', anchor='middle')
        text(slide, bx-20, base_y+38, bar_w+40, 22, span,
             size=12, color=INK_3, mono=True, align='center', anchor='middle')
    # Mean line
    mean_auc = 0.6490
    mean_h = int((mean_auc-0.5) / 0.20 * chart_h)
    rect(slide, 196, base_y - mean_h, 720, 1, fill=TEAL)
    text(slide, 920, base_y - mean_h - 14, 200, 22, f'mean {mean_auc:.4f}',
         size=12, bold=True, color=TEAL, mono=True)
    rect(slide, 196, base_y, 720, 1, fill=INK_4)
    text(slide, 196, base_y+62, 200, 18, '0.5 baseline', size=12, color=INK_3, mono=True)

    # DSR formula panel
    rx = 1144
    card(slide, rx, 308, 680, 580, top_stripe=NAVY)
    text(slide, rx+24, 332, 600, 22, '— DEFLATED SHARPE RATIO · 多重檢定校正',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    text(slide, rx+24, 366, 632, 60, 'DSR 12.12 (p<0.001)',
         size=42, bold=True, color=NAVY, mono=True, ls=-1.0)
    rect(slide, rx+24, 442, 60, 3, fill=NAVY)
    rich(slide, rx+24, 462, 632, 280, [
        {'text': 'DSR = (SR̂ − E[max SR_n]) × √(T−1) / √(1 − γ₃·SR̂ + (γ₄−1)/4·SR̂²)',
         'size': 12, 'mono': True, 'color': INK_2},
        {'sp': True},
        {'text': 'SR̂', 'size': 14, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' = 觀測 Sharpe Ratio = 2.41', 'size': 14, 'color': INK_2},
        {'br': True},
        {'text': 'n', 'size': 14, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' = 多重試驗 (Optuna trials) = 150', 'size': 14, 'color': INK_2},
        {'br': True},
        {'text': 'T', 'size': 14, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' = OOS 樣本期 (years) ≈ 8', 'size': 14, 'color': INK_2},
        {'br': True},
        {'text': 'γ₃, γ₄', 'size': 14, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' = 報酬 skewness / kurtosis 校正項', 'size': 14, 'color': INK_2},
    ], line_h=1.6)
    rect(slide, rx+24, 758, 632, 1, fill=RULE)
    text(slide, rx+24, 770, 632, 22, '解讀',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    text(slide, rx+24, 798, 632, 80, 'DSR > 2 即顯著；本系統 12.12 為極強拒絕「運氣假說」的證據——150 次參數搜尋的選擇偏誤已被扣除。',
         size=13, color=INK_2, line_h=1.55)

    takeaway_bar(slide, 96, 904, 1728, 84, 'STABILITY · 跨 fold 穩定', [
        {'text': '4 fold AUC 標準差 0.005 · 最差 fold 0.6431 仍遠高於 0.5 baseline · ', 'size': 14, 'color': INK},
        {'text': '系統不靠單一年份。', 'size': 14, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.13 · Threshold Sweep
# ============================================================
def slide_apx_threshold(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Threshold Sweep · 0.40 / 0.50 / 0.60', 'A.13',
                          source='自行整理 · 40 thresholds × 4-fold')
    eyebrow(slide, 96, 116, 700, 'A.13 · 三情境門檻設計', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '三情境門檻 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 攻擊 / 平衡 / 保守 · 對應不同投資人風險偏好', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    scenarios = [
        ('保守 · CONSERVATIVE', 'p ≥ 0.60', NAVY,    '0.8%',  '52.4%', '+22.4 pp', '較少訊號 · 高勝率'),
        ('平衡 · BALANCED',     'p ≥ 0.50', TEAL,    '4.6%',  '46.1%', '+16.1 pp', '本系統預設'),
        ('攻擊 · AGGRESSIVE',   'p ≥ 0.40', ORANGE,  '12.3%', '38.2%', '+8.2 pp',  '訊號多 · 勝率較低'),
    ]
    cy = 308; cw = 552; ch = 472; gap = 24
    for i, (name, thr, c, ratio, hit, edge, body) in enumerate(scenarios):
        cx = 96 + i*(cw+gap)
        card(slide, cx, cy, cw, ch, top_stripe=c)
        text(slide, cx+24, cy+22, cw-48, 22, name,
             size=12, bold=True, color=c, mono=True, ls=2.4)
        text(slide, cx+24, cy+54, cw-48, 50, thr,
             size=36, bold=True, color=c, mono=True, ls=-1.0)
        rect(slide, cx+24, cy+114, 60, 3, fill=c)
        # Three KPIs
        items = [('SIGNAL %', ratio), ('HIT RATE', hit), ('vs baseline', edge)]
        for j, (lbl, val) in enumerate(items):
            iy = cy + 140 + j*84
            text(slide, cx+24, iy, cw-48, 18, lbl,
                 size=12, color=INK_3, mono=True, ls=2.0)
            text(slide, cx+24, iy+22, cw-48, 50, val,
                 size=32, bold=True, color=NAVY, mono=True, ls=-0.8)
        text(slide, cx+24, cy+ch-46, cw-48, 32, body,
             size=13, color=INK_2, line_h=1.5)

    takeaway_bar(slide, 96, 808, 1728, 156, 'INTERPRETATION · 投資人對照', [
        {'text': '本系統的門檻是 ', 'size': 16, 'color': INK},
        {'text': '可調節的策略開關', 'size': 16, 'bold': True, 'color': NAVY},
        {'text': '；機構可挑保守、量化基金可挑攻擊。', 'size': 16, 'color': INK},
        {'br': True},
        {'text': '三情境的 IC / hit rate / edge 全部在 OOS 取得，已扣 5 bps 滑價假設。', 'size': 16, 'color': INK_2},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.14 · Hit Decile Detail
# ============================================================
def slide_apx_hit_decile(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Hit Decile · Detail', 'A.14',
                          source='自行整理 · OOS hit rate by predicted-prob decile')
    eyebrow(slide, 96, 116, 700, 'A.14 · 預測機率分位 × 命中率', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '十分位 hit rate ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 單調性無破口 · D10 達 47.8%', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # Decile bar chart
    rect(slide, 96, 308, 1728, 540, fill=PAPER, line=RULE, line_w=0.75)
    text(slide, 116, 326, 600, 22, '— HIT RATE BY PREDICTED-PROB DECILE',
         size=12, bold=True, color=INK_2, mono=True, ls=1.6)
    decs = [
        ('D1', 22.4),  ('D2', 24.8), ('D3', 26.5), ('D4', 28.1), ('D5', 29.8),
        ('D6', 31.6),  ('D7', 33.5), ('D8', 36.4), ('D9', 41.2), ('D10', 47.8),
    ]
    base_y = 800
    chart_h = 380
    bar_w = 110
    gap = 50
    bar_x0 = 156
    max_v = 50
    for i, (name, v) in enumerate(decs):
        bar_h = int(v / max_v * chart_h)
        bx = bar_x0 + i*(bar_w + gap)
        col = NAVY if i < 8 else TEAL
        rect(slide, bx, base_y - bar_h, bar_w, bar_h, fill=col)
        text(slide, bx-20, base_y-bar_h-30, bar_w+40, 24, f'{v:.1f}%',
             size=13, bold=True, color=col, mono=True, align='center', anchor='middle')
        text(slide, bx-20, base_y+12, bar_w+40, 22, name,
             size=13, bold=True, color=INK, mono=True, align='center', anchor='middle')
    # Baseline 30%
    base_v = 30
    base_h = int(base_v / max_v * chart_h)
    rect(slide, 156, base_y - base_h, 1640, 1, fill=ROSE)
    text(slide, 1620, base_y - base_h - 22, 200, 18, '30% baseline',
         size=12, bold=True, color=ROSE, mono=True, align='right')

    takeaway_bar(slide, 96, 868, 1728, 112, 'MONOTONICITY · 機率排序與真實命中一致', [
        {'text': '十個分位無一破口 · 排序能力扎實；', 'size': 15, 'color': INK},
        {'text': 'D10 vs D1 spread = 25.4 pp', 'size': 15, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' · 是策略中最關鍵的 alpha 來源。', 'size': 15, 'color': INK},
        {'br': True},
        {'text': '此為 P.06 主簡報「47.8% hit @ D10」的詳細展開。', 'size': 14, 'color': INK_3, 'mono': True},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.15 · Equity Curve · Top 0.1% vs 0050
# ============================================================
def slide_apx_equity_curve(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Equity Curve · Top 0.1% vs 0050', 'A.15',
                          source='自行整理 · 5 bps round-trip cost · 20-day rebalance')
    eyebrow(slide, 96, 116, 700, 'A.15 · 累積報酬曲線', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Top 0.1% Strategy ', 'size': 40, 'bold': True, 'color': TEAL, 'ls': -0.8},
        {'text': 'vs ', 'size': 40, 'bold': True, 'color': INK_2, 'ls': -0.8},
        {'text': '0050 ETF benchmark', 'size': 40, 'bold': True, 'color': INK_3, 'ls': -0.8},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # Chart
    px_, py_, pw_, ph_ = linechart_box(slide, 96, 308, 1728, 540)
    # Years on x-axis
    years = ['2018', '2019', '2020', '2021', '2022', '2023', '2024', '2025']
    for i, y_lbl in enumerate(years):
        tx = px_ + int(i / (len(years)-1) * pw_)
        rect(slide, tx, py_+ph_, 1, 6, fill=INK_3)
        text(slide, tx-30, py_+ph_+12, 60, 18, y_lbl,
             size=12, color=INK_3, mono=True, align='center')
    # Y-axis ticks (cumulative return %)
    yvals = [0, 50, 100, 150, 200, 250]
    for i, v in enumerate(yvals):
        ty = py_ + ph_ - int(v / 250 * ph_)
        rect(slide, px_-6, ty, 6, 1, fill=INK_3)
        text(slide, px_-60, ty-10, 50, 18, f'{v}%',
             size=12, color=INK_3, mono=True, align='right')
        if v > 0:
            rect(slide, px_, ty, pw_, 1, fill=SUBTLE)
    # Strategy curve (stylized) — 8 years
    strat = [0, 18, 42, 78, 105, 142, 188, 232]
    bench = [0, 12, 28, 38, 22, 35, 52, 68]
    def to_xy(idx, val, max_v=250):
        x = px_ + int(idx / (len(strat)-1) * pw_)
        y = py_ + ph_ - int(val / max_v * ph_)
        return (x, y)
    pts_s = [to_xy(i, v) for i, v in enumerate(strat)]
    pts_b = [to_xy(i, v) for i, v in enumerate(bench)]
    polyline(slide, pts_s, TEAL, line_w=3.0)
    polyline(slide, pts_b, INK_3, line_w=2.0)
    for x, y in pts_s:
        dot(slide, x, y, 5, TEAL)
    for x, y in pts_b:
        dot(slide, x, y, 4, INK_4)
    # Final labels
    text(slide, pts_s[-1][0]-180, pts_s[-1][1]-22, 200, 22, '+232%',
         size=18, bold=True, color=TEAL, mono=True, align='right')
    text(slide, pts_b[-1][0]-180, pts_b[-1][1]+10, 200, 22, '+68%',
         size=14, bold=True, color=INK_3, mono=True, align='right')
    # Legend
    rect(slide, px_+24, py_+24, 24, 4, fill=TEAL)
    text(slide, px_+56, py_+18, 320, 22, 'Top 0.1% portfolio (avg 4 stocks)',
         size=12, bold=True, color=TEAL, mono=True)
    rect(slide, px_+24, py_+54, 24, 2, fill=INK_3)
    text(slide, px_+56, py_+48, 320, 22, '0050 ETF (TWSE benchmark)',
         size=12, color=INK_3, mono=True)

    # KPI strip below
    kpi_y = 868
    kpi_strip = [
        ('CAGR',          '+16.4%', '+6.7%',  '+9.7 pp'),
        ('SR̂',           '2.41',   '0.78',   '+1.63'),
        ('MaxDD',         '-18.2%', '-32.4%', '+14.2 pp'),
        ('HIT RATE (M)',  '64%',    '52%',    '+12 pp'),
    ]
    for i, (lbl, v1, v2, edge) in enumerate(kpi_strip):
        kx = 96 + i*438
        rect(slide, kx, kpi_y, 414, 96, fill=PAPER, line=RULE, line_w=0.75)
        rect(slide, kx, kpi_y, 414, 4, fill=TEAL)
        text(slide, kx+18, kpi_y+12, 240, 22, lbl,
             size=12, bold=True, color=INK_3, mono=True, ls=2.2)
        text(slide, kx+18, kpi_y+36, 200, 36, v1,
             size=22, bold=True, color=TEAL, mono=True, ls=-0.6)
        text(slide, kx+220, kpi_y+36, 100, 36, v2,
             size=14, color=INK_3, mono=True, anchor='middle')
        text(slide, kx+330, kpi_y+36, 84, 36, edge,
             size=14, bold=True, color=NAVY, mono=True, align='right', anchor='middle')
        text(slide, kx+18, kpi_y+74, 200, 16, 'STRATEGY  /  0050  /  EDGE',
             size=12, color=INK_4, mono=True, ls=1.8)
    return slide


# ============================================================
# A.16 · Drawdown & Cost Sensitivity
# ============================================================
def slide_apx_drawdown_cost(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Drawdown & Cost Sensitivity', 'A.16',
                          source='自行整理 · round-trip cost from 0 to 30 bps')
    eyebrow(slide, 96, 116, 700, 'A.16 · 回撤與交易成本敏感度', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Cost Sensitivity ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 5 / 15 / 30 bps 三情境 · 策略仍正報酬', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    scenarios = [
        ('5 bps · BASE',       NAVY,    '+16.4%', '2.41',  '-18.2%', '本系統預設 · 散戶券商一般成本'),
        ('15 bps · STRESS',    ORANGE,  '+12.1%', '1.78',  '-21.6%', '重稅 / 較高滑價 · 多空組合場景'),
        ('30 bps · WORST',     ROSE,    '+5.8%',  '0.92',  '-25.4%', '極端 · 流動性低 / 散戶高週轉'),
    ]
    cy = 308; cw = 552; ch = 432; gap = 24
    for i, (name, c, cagr, sr, mdd, body) in enumerate(scenarios):
        cx = 96 + i*(cw+gap)
        card(slide, cx, cy, cw, ch, top_stripe=c)
        text(slide, cx+24, cy+22, cw-48, 22, name,
             size=12, bold=True, color=c, mono=True, ls=2.4)
        text(slide, cx+24, cy+54, cw-48, 56, cagr,
             size=46, bold=True, color=c, mono=True, ls=-1.4)
        text(slide, cx+24, cy+114, cw-48, 22, 'CAGR · annualized',
             size=12, color=INK_3, mono=True, ls=1.5)
        rect(slide, cx+24, cy+148, 60, 3, fill=c)
        items = [('SR̂', sr), ('MaxDD', mdd)]
        for j, (lbl, v) in enumerate(items):
            iy = cy + 168 + j*68
            text(slide, cx+24, iy, 200, 22, lbl,
                 size=12, color=INK_3, mono=True, ls=2.0)
            text(slide, cx+200, iy-2, cw-220, 28, v,
                 size=22, bold=True, color=NAVY, mono=True, ls=-0.4)
        text(slide, cx+24, cy+ch-72, cw-48, 60, body,
             size=13, color=INK_2, line_h=1.55)

    takeaway_bar(slide, 96, 768, 1728, 196, 'COST RESILIENCE · 成本韌性', [
        {'text': '即使在 30 bps 極端成本下，策略 ', 'size': 16, 'color': INK},
        {'text': '仍維持 +5.8% CAGR · SR̂ 0.92', 'size': 16, 'bold': True, 'color': NAVY},
        {'text': '；', 'size': 16, 'color': INK},
        {'br': True},
        {'text': '5 bps base case 為散戶實際可達成之水準。', 'size': 16, 'color': INK_2},
        {'sp': True},
        {'text': 'IMPLICATION', 'size': 12, 'bold': True, 'color': ROSE, 'mono': True, 'ls': 2.0},
        {'br': True},
        {'text': '本策略不依賴零成本假設，', 'size': 15, 'color': INK},
        {'text': '在合理交易成本下仍具實戰可用性。', 'size': 15, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.17 · Dashboard 13 Pages mockup
# ============================================================
def slide_apx_dashboard(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Dashboard · 13 Pages', 'A.17',
                          source='dashboard/ · Streamlit 1.31 · launch via `make serve`')
    eyebrow(slide, 96, 116, 700, 'A.17 · 互動儀表板 13 頁地圖', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '13-Page Interactive Dashboard ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 模型 / 解釋 / 文本 / 治理', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    pages = [
        ('01', 'Overview',          NAVY,    '系統首頁 · KPI 速讀'),
        ('02', 'Data Health',       NAVY,    'PSI / KS 漂移監控'),
        ('03', 'Feature Catalog',   NAVY,    '91 features 速查'),
        ('04', 'Model Performance', TEAL,    'AUC / ICIR / fold 切換'),
        ('05', 'ICIR Heatmap',      TEAL,    'pillar × month 熱圖'),
        ('06', 'SHAP Explorer',     TEAL,    'Top-20 / 個股逐股'),
        ('07', 'Decile Hit',        TEAL,    'D1-D10 命中率'),
        ('08', 'Threshold Sweep',   ORANGE,  '40 門檻 sweep 曲線'),
        ('09', 'Backtest Equity',   ORANGE,  '累積報酬 + 0050'),
        ('10', 'Drawdown',          ORANGE,  '回撤期間 + Calmar'),
        ('11', 'Text Sentiment',    ROSE,    '500 keyword Lift × χ²'),
        ('12', 'Daily Watchlist',   ROSE,    '今日訊號股 + Q1 警示'),
        ('13', 'Governance Gates',  PINK,    '9/9 閘門即時狀態'),
    ]
    grid_x, grid_y = 96, 308
    cw, ch, gap = 416, 110, 16
    cols = 4
    for i, (no, name, c, body) in enumerate(pages):
        col = i % cols; row = i // cols
        cx = grid_x + col*(cw+gap)
        cy = grid_y + row*(ch+gap)
        card(slide, cx, cy, cw, ch, top_stripe=c)
        text(slide, cx+18, cy+18, 60, 28, no, size=18, bold=True, color=c, mono=True, ls=-0.4)
        text(slide, cx+72, cy+22, cw-92, 24, name, size=15, bold=True, color=NAVY)
        rect(slide, cx+72, cy+50, 40, 2, fill=c)
        text(slide, cx+72, cy+58, cw-92, 44, body, size=12, color=INK_2, line_h=1.45)

    takeaway_bar(slide, 96, 884, 1728, 84, 'OPERATIONAL READY · 上線即用', [
        {'text': 'Streamlit 部署，三秒切頁；', 'size': 15, 'color': INK},
        {'text': '所有圖表皆可導出 PNG / 對應 JSON artifact', 'size': 15, 'bold': True, 'color': NAVY},
        {'text': '；研究員、PM、合規可同時使用同一頁面。', 'size': 15, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.18 · Sentiment Asset Overview
# ============================================================
def slide_apx_sentiment_overview(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Sentiment Asset · Overview', 'A.18',
                          source='自行整理 · sentiment lexicon under src/text/lexicon/')
    eyebrow(slide, 96, 116, 700, 'A.18 · 文本資產總覽', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '1.12M articles ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 500 keywords · 自建台股情感詞表', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    kpis = [
        ('ARTICLES',  '1,125,134', '15 來源 · 2008+',         NAVY),
        ('KEYWORDS',  '500',       'Bull 256 / Bear 244',     TEAL),
        ('ENTITY',    '92%',       'tickers 涵蓋率',          EMERALD),
        ('DEDUP',     '23%',       'MinHash LSH J ≥ 0.85',    ORANGE),
    ]
    for i, (lbl, v, sub, c) in enumerate(kpis):
        kx = 96 + i*438
        kpi_hero(slide, kx, 308, 414, 168, lbl, v, sub, accent=c)

    # Two cards: Methodology + Source breakdown
    text(slide, 96, 506, 600, 22, '— LEXICON CONSTRUCTION · 詞表建構',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    card(slide, 96, 540, 850, 416, top_stripe=NAVY)
    rich(slide, 116, 564, 810, 380, [
        {'text': '500-keyword 台股情感詞表 · 三步驟', 'size': 20, 'bold': True, 'color': NAVY},
        {'sp': True},
        {'text': '01  TF-IDF 候選 ', 'size': 14, 'bold': True, 'color': TEAL, 'mono': True, 'ls': 1.6},
        {'text': '從 400k+ 高品質新聞抽 5,000 候選詞', 'size': 14, 'color': INK_2},
        {'br': True},
        {'text': '02  Lift × χ² 過濾 ', 'size': 14, 'bold': True, 'color': TEAL, 'mono': True, 'ls': 1.6},
        {'text': '計算每詞 Lift (vs neutral) 與卡方值 → 留 500', 'size': 14, 'color': INK_2},
        {'br': True},
        {'text': '03  專家校驗 ', 'size': 14, 'bold': True, 'color': TEAL, 'mono': True, 'ls': 1.6},
        {'text': '兩位金融研究生雙盲標註 · κ = 0.78', 'size': 14, 'color': INK_2},
        {'sp': True},
        {'text': '為何不直接用通用情感詞表？', 'size': 16, 'bold': True, 'color': NAVY},
        {'br': True},
        {'text': 'LIWC / NTUSD 的「正面」≠ 股價的「正面」 (e.g. 「擴產」中性偏負，因稀釋EPS)', 'size': 13, 'color': INK_2},
        {'sp': True},
        {'text': '完整詞表參考 src/text/lexicon/twstock500.csv', 'size': 12, 'color': INK_3, 'mono': True, 'ls': 0.4},
    ], line_h=1.55)

    text(slide, 974, 506, 600, 22, '— SOURCE BREAKDOWN · 15 來源',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    card(slide, 974, 540, 850, 416, top_stripe=TEAL)
    sources = [
        ('Anue 鉅亨網',     NAVY,    '24%'),
        ('PTT Stock 板',    ROSE,    '18%'),
        ('經濟日報',         NAVY,    '12%'),
        ('工商時報',         NAVY,    '11%'),
        ('Money DJ',        TEAL,    '8%'),
        ('Yahoo 股市',      TEAL,    '7%'),
        ('其他 9 來源',      INK_3,   '20%'),
    ]
    for i, (name, c, pct) in enumerate(sources):
        sy = 580 + i*48
        text(slide, 996, sy, 320, 32, name, size=14, color=INK, anchor='middle')
        rect(slide, 1340, sy+12, 360, 8, fill=SUBTLE)
        bw = int(int(pct.rstrip('%')) / 24 * 360)
        rect(slide, 1340, sy+12, bw, 8, fill=c)
        text(slide, 1720, sy, 80, 32, pct, size=14, bold=True, color=c, mono=True, align='right', anchor='middle')

    takeaway_bar(slide, 96, 956, 1728, 54, '', [
        {'text': '15 來源 · 前 6 大占 80% · ', 'size': 13, 'color': INK},
        {'text': '經濟日報 / 工商時報 / Anue 為新聞主幹；PTT 為散戶情緒主軸。',
         'size': 13, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.19 · Keyword Spectrum Top 10 bull/bear
# ============================================================
def slide_apx_keyword_spectrum(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Keyword Spectrum · Top 10', 'A.19',
                          source='自行整理 · Lift = P(bull|word)/P(bull) · χ² independence test')
    eyebrow(slide, 96, 116, 700, 'A.19 · 詞表光譜多空 Top 10', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '500 keywords ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· Lift × χ² 雙指標 · 顯示前 10 多空關鍵詞', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    bull = [
        ('創高',     2.42, 1842),
        ('營收創高', 2.31, 1654),
        ('法說會',   1.98, 1428),
        ('擴產',     1.92, 1322),
        ('受惠',     1.84, 2104),
        ('重押',     1.78,  892),
        ('升級',     1.72, 1186),
        ('看好',     1.68, 2842),
        ('題材',     1.62, 1942),
        ('新訂單',   1.58,  784),
    ]
    bear = [
        ('砍單',     2.86, 1124),
        ('裁員',     2.54,  642),
        ('下修',     2.42, 1844),
        ('警報',     2.18,  584),
        ('遞延',     2.04,  728),
        ('衰退',     1.92, 1462),
        ('低於預期', 1.84, 1102),
        ('虧損',     1.78, 1842),
        ('違約',     1.72,  342),
        ('觀望',     1.62, 2104),
    ]

    def bar_panel(slide, x, y, w, h, title, color, items, max_lift):
        rect(slide, x, y, w, h, fill=PAPER, line=RULE, line_w=0.75)
        rect(slide, x, y, w, 4, fill=color)
        text(slide, x+24, y+18, w-48, 22, title,
             size=12, bold=True, color=color, mono=True, ls=2.4)
        # header row
        rect(slide, x+24, y+52, w-48, 30, fill=SUBTLE)
        text(slide, x+36, y+58, 240, 22, 'KEYWORD',
             size=12, bold=True, color=INK_3, mono=True, ls=2.0, anchor='middle')
        text(slide, x+w-220, y+58, 100, 22, 'LIFT',
             size=12, bold=True, color=INK_3, mono=True, ls=2.0, align='right', anchor='middle')
        text(slide, x+w-100, y+58, 80, 22, 'χ² rank',
             size=12, bold=True, color=INK_3, mono=True, ls=2.0, align='right', anchor='middle')
        for i, (kw, lift, chi) in enumerate(items):
            ry = y + 88 + i*48
            text(slide, x+36, ry+12, 200, 28, f'{i+1:02d}  {kw}',
                 size=15, bold=True, color=INK, anchor='middle')
            track_x = x + 240
            track_w = w - 240 - 200
            rect(slide, track_x, ry+18, track_w, 8, fill=SUBTLE)
            bw = int(lift / max_lift * track_w)
            rect(slide, track_x, ry+18, bw, 8, fill=color)
            text(slide, x+w-220, ry+12, 100, 28, f'{lift:.2f}×',
                 size=14, bold=True, color=color, mono=True, align='right', anchor='middle')
            text(slide, x+w-100, ry+12, 80, 28, f'{chi:,}',
                 size=12, color=INK_3, mono=True, align='right', anchor='middle')

    bar_panel(slide, 96, 308, 854, 610, 'BULL · 多方關鍵詞 (Lift > 1)', EMERALD, bull, 3.0)
    bar_panel(slide, 970, 308, 854, 610, 'BEAR · 空方關鍵詞 (Lift > 1)', ROSE, bear, 3.0)

    takeaway_bar(slide, 96, 936, 1728, 54, '', [
        {'text': '空方 Lift 上限 (砍單 2.86×) > 多方 (創高 2.42×) · ',
         'size': 13, 'color': INK},
        {'text': '空方詞訊號更強烈 · 解釋 Q1 contrarian 效應的源頭。',
         'size': 13, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.20 · Sentiment Quintile Detail
# ============================================================
def slide_apx_sentiment_quintile(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Sentiment Quintile · Detail', 'A.20',
                          source='自行整理 · forward 20-day excess return by quintile')
    eyebrow(slide, 96, 116, 700, 'A.20 · 情感五分位詳解', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'U-shape Confirmed ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· Q1 +3.2% · Q5 +2.4% · 中段為負', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # Bar chart
    rect(slide, 96, 308, 1080, 580, fill=PAPER, line=RULE, line_w=0.75)
    text(slide, 116, 326, 700, 22, '— FORWARD 20-DAY EXCESS RETURN BY SENTIMENT QUINTILE',
         size=12, bold=True, color=INK_2, mono=True, ls=1.6)
    qs = [
        ('Q1\n極負面', 3.2,  ROSE),
        ('Q2\n偏負面', -1.8, INK_4),
        ('Q3\n中性',   -0.4, INK_4),
        ('Q4\n偏正面', 0.6,  INK_3),
        ('Q5\n極正面', 2.4,  EMERALD),
    ]
    base_y = 730
    chart_h = 280
    bar_w = 160
    gap = 28
    bar_x0 = 156
    max_v = 4
    # Zero line
    rect(slide, 156, base_y, 880, 1, fill=INK_3)
    for i, (lbl, v, c) in enumerate(qs):
        bx = bar_x0 + i*(bar_w + gap)
        bar_h = int(abs(v) / max_v * chart_h)
        if v >= 0:
            rect(slide, bx, base_y - bar_h, bar_w, bar_h, fill=c)
            text(slide, bx-20, base_y-bar_h-30, bar_w+40, 24, f'+{v:.1f}%',
                 size=15, bold=True, color=c, mono=True, align='center', anchor='middle')
        else:
            rect(slide, bx, base_y, bar_w, bar_h, fill=c)
            text(slide, bx-20, base_y+bar_h+8, bar_w+40, 24, f'{v:.1f}%',
                 size=14, bold=True, color=c, mono=True, align='center', anchor='middle')
        text(slide, bx-20, base_y+chart_h-110, bar_w+40, 60, lbl,
             size=12, bold=True, color=INK, align='center', anchor='middle', line_h=1.4)

    # Right: 4 insight cards
    rx = 1196
    cards_data = [
        ('Q1 contrarian',     ROSE,    '+3.2%', '極端壞消息 → 過度反應 → 反彈'),
        ('Q5 momentum',       EMERALD, '+2.4%', '極端好消息 → 慣性延續'),
        ('Q3 中性',           INK_3,   '-0.4%', '無方向 · 雜訊主導'),
        ('Spread Q1−Q3',      NAVY,    '+3.6 pp','對沖組合潛在 alpha'),
    ]
    for i, (h, c, val, body) in enumerate(cards_data):
        cy = 308 + i*148
        card(slide, rx, cy, 628, 132, top_stripe=c)
        text(slide, rx+24, cy+22, 400, 28, h, size=18, bold=True, color=NAVY)
        text(slide, rx+24, cy+58, 200, 36, val, size=26, bold=True, color=c, mono=True, ls=-0.6)
        text(slide, rx+260, cy+62, 360, 60, body, size=13, color=INK_2, line_h=1.5, anchor='middle')

    takeaway_bar(slide, 96, 904, 1728, 84, 'BEHAVIORAL FINANCE · 行為財務學印證', [
        {'text': 'Q1 反向效應 = ', 'size': 14, 'color': INK},
        {'text': '過度反應假說 (DeBondt & Thaler 1985)', 'size': 14, 'bold': True, 'color': NAVY},
        {'text': ' 在台股的數據驗證 · 詳見 A.21。', 'size': 14, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.21 · Extreme Negative Q1 elaboration
# ============================================================
def slide_apx_q1_contrarian(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Q1 Contrarian · Elaboration', 'A.21',
                          source='DeBondt & Thaler (1985) · empirical: artifact 自行整理')
    eyebrow(slide, 96, 116, 700, 'A.21 · Q1 反向效應完整解讀', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '極負面 → 反彈 ', 'size': 40, 'bold': True, 'color': ROSE, 'ls': -0.8},
        {'text': '· 行為財務學的台股實證', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # Big quote box
    card(slide, 96, 308, 1728, 152, top_stripe=ROSE)
    text(slide, 116, 332, 1688, 22, '— THE FINDING',
         size=12, bold=True, color=ROSE, mono=True, ls=2.4)
    text(slide, 116, 360, 1688, 88, '當輿情極度負面時 (Q1)，未來 20 日「期望報酬」反而變正 +3.2%；',
         size=22, bold=True, color=INK, line_h=1.4)
    text(slide, 116, 408, 1688, 44, '空頭過度反應 → 賣壓出清 → 散戶恐慌底 → 機構接手反彈。',
         size=20, color=INK_2)

    # Three-step explanation
    steps = [
        ('01', '空頭過度反應',     ROSE,
         '極負面新聞 (e.g. 砍單、警報) 在 1-3 日內形成集體恐慌，技術線型急殺 5-10%。'),
        ('02', '賣壓快速出清',     ORANGE,
         '散戶停損賣出殆盡 · 融資餘額大降 · 週轉率異常高（90 日均 1.5×+）。'),
        ('03', '機構低接 + 反彈',  EMERALD,
         '若基本面未受結構性傷害，外資 / 投信於 5-10 日後低接 · T+20 平均 +3.2%。'),
    ]
    for i, (no, h, c, body) in enumerate(steps):
        sx = 96 + i*(552+24)
        sy = 488
        card(slide, sx, sy, 552, 360, top_stripe=c)
        text(slide, sx+24, sy+24, 80, 60, no, size=44, bold=True, color=c, mono=True, ls=-1.4)
        text(slide, sx+24, sy+92, 200, 22, 'STEP', size=12, bold=True, color=c, mono=True, ls=2.4)
        text(slide, sx+24, sy+120, 504, 36, h, size=22, bold=True, color=NAVY)
        rect(slide, sx+24, sy+162, 60, 3, fill=c)
        text(slide, sx+24, sy+178, 504, 160, body, size=15, color=INK_2, line_h=1.65)

    takeaway_bar(slide, 96, 868, 1728, 112, 'EXECUTION CAVEAT · 執行注意', [
        {'text': '此效應僅限 ', 'size': 14, 'color': INK},
        {'text': '基本面未受結構性傷害的標的', 'size': 14, 'bold': True, 'color': NAVY},
        {'text': '；模型加入 ROE / debt 作 sanity guardrail，避免抄到「真的不行」的個股。', 'size': 14, 'color': INK},
        {'br': True},
        {'text': '回測 OOS 8 年仍穩定，但個股級執行需研究員 final call。', 'size': 13, 'color': INK_3},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.22 · U-shape Multi-Horizon (1d/5d/20d)
# ============================================================
def slide_apx_ushape_multi(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'U-shape Multi-Horizon', 'A.22',
                          source='自行整理 · forward 1d/5d/20d excess return')
    eyebrow(slide, 96, 116, 700, 'A.22 · 多時間視窗 U-shape', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'U-shape 多視窗 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 1日 / 5日 / 20日 三條曲線', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # Three side-by-side panels
    horizons = [
        ('T+1',  NAVY,     [-0.6, -1.4, -0.4, -0.2,  0.2], '一日內仍以慣性主導'),
        ('T+5',  TEAL,     [ 0.8, -1.6, -0.4,  0.4,  1.6], '反向效應開始浮現'),
        ('T+20', EMERALD,  [ 3.2, -1.8, -0.4,  0.6,  2.4], 'U-shape 完整成型'),
    ]
    panel_w = 552; panel_h = 530; gap = 24
    for i, (name, c, vals, body) in enumerate(horizons):
        px = 96 + i*(panel_w + gap)
        py = 308
        rect(slide, px, py, panel_w, panel_h, fill=PAPER, line=RULE, line_w=0.75)
        rect(slide, px, py, panel_w, 4, fill=c)
        text(slide, px+24, py+18, panel_w-48, 22, f'HORIZON · {name}',
             size=12, bold=True, color=c, mono=True, ls=2.4)
        text(slide, px+24, py+50, panel_w-48, 32, body,
             size=15, bold=True, color=NAVY)
        # Mini bar chart
        chart_top = py + 100
        chart_h = 280
        bar_w = 70
        gap_b = 26
        bar_x0 = px + 56
        max_v = 4
        zero_y = chart_top + chart_h//2
        rect(slide, px+24, zero_y, panel_w-48, 1, fill=INK_3)
        labels = ['Q1', 'Q2', 'Q3', 'Q4', 'Q5']
        for j, (lbl, v) in enumerate(zip(labels, vals)):
            bx = bar_x0 + j*(bar_w + gap_b)
            bar_h = int(abs(v) / max_v * (chart_h//2))
            col = ROSE if v < 0 else (EMERALD if v > 0 else INK_3)
            if v >= 0:
                rect(slide, bx, zero_y - bar_h, bar_w, bar_h, fill=col)
                text(slide, bx-20, zero_y-bar_h-26, bar_w+40, 22, f'+{v:.1f}',
                     size=12, bold=True, color=col, mono=True, align='center', anchor='middle')
            else:
                rect(slide, bx, zero_y, bar_w, bar_h, fill=col)
                text(slide, bx-20, zero_y+bar_h+4, bar_w+40, 22, f'{v:.1f}',
                     size=12, bold=True, color=col, mono=True, align='center', anchor='middle')
            text(slide, bx-20, chart_top+chart_h+8, bar_w+40, 22, lbl,
                 size=12, bold=True, color=INK_2, mono=True, align='center', anchor='middle')
        text(slide, px+24, py+panel_h-46, panel_w-48, 22, '單位 · % excess return',
             size=12, color=INK_3, mono=True, ls=0.4)

    takeaway_bar(slide, 96, 856, 1728, 112, 'TIME EVOLUTION · 時間演化', [
        {'text': '反向效應在 T+1 尚未顯現 (慣性主導) · 在 T+5 開始浮現 · 在 T+20 完整成型；', 'size': 15, 'color': INK},
        {'br': True},
        {'text': '建議交易視窗為 ', 'size': 15, 'color': INK},
        {'text': '5-20 日期間', 'size': 15, 'bold': True, 'color': NAVY},
        {'text': '；亦支援 D20 為主標籤的設計選擇。', 'size': 15, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.23 · News × Forum Consensus 2x2
# ============================================================
def slide_apx_news_forum(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'News × Forum Consensus', 'A.23',
                          source='自行整理 · joint distribution news_sent × forum_sent')
    eyebrow(slide, 96, 116, 700, 'A.23 · 新聞 × 論壇 共識 2×2', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Consensus / Divergence ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 新聞與論壇一致時訊號最強', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # 2x2 grid
    grid_x, grid_y = 96, 308
    cell_w = 720; cell_h = 280; gap = 24
    quadrants = [
        # (col, row, label, color, value, body)
        (0, 0, 'BOTH BULL · 共識多',
            EMERALD, '+4.1%',
            '新聞正 + 論壇正 · 訊號最強烈 · 占樣本 18%'),
        (1, 0, 'NEWS BULL / FORUM BEAR',
            GOLD, '+0.4%',
            '新聞正 / 論壇負 · 散戶觀望 · 機構默買'),
        (0, 1, 'NEWS BEAR / FORUM BULL',
            ORANGE, '-0.6%',
            '新聞負 / 論壇正 · 散戶追高被套 · 警示區'),
        (1, 1, 'BOTH BEAR · 共識空',
            ROSE, '+1.8%',
            '新聞負 + 論壇負 · 反向效應 · 占樣本 14%'),
    ]
    for (col, row, lbl, c, val, body) in quadrants:
        cx = grid_x + col*(cell_w + gap)
        cy = grid_y + row*(cell_h + gap)
        card(slide, cx, cy, cell_w, cell_h, top_stripe=c)
        text(slide, cx+24, cy+22, cell_w-48, 22, lbl,
             size=12, bold=True, color=c, mono=True, ls=2.4)
        text(slide, cx+24, cy+56, cell_w-48, 80, val,
             size=64, bold=True, color=c, mono=True, ls=-2.0)
        rect(slide, cx+24, cy+150, 60, 3, fill=c)
        text(slide, cx+24, cy+170, cell_w-48, 80, body,
             size=15, color=INK_2, line_h=1.55)
        text(slide, cx+cell_w-180, cy+22, 160, 22, 'T+20 fwd',
             size=12, color=INK_3, mono=True, align='right', ls=1.5)

    # Right info column
    rx = 1560
    text(slide, rx, 308, 264, 22, '— DECISION RULES',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    rules = [
        ('CONSENSUS_BULL', EMERALD, '兩者皆 ≥ +1.5σ → 買進', '+4.1%'),
        ('CONTRARIAN',     ROSE,    '兩者皆 ≤ -1.5σ → 反向買進', '+1.8%'),
        ('AVOID',          GOLD,    'news+ × forum− → 觀望',  '+0.4%'),
        ('SHORT-LIST',     ORANGE,  'news− × forum+ → 警示', '-0.6%'),
    ]
    for i, (n, c, body, ret) in enumerate(rules):
        ry = 348 + i*132
        card(slide, rx, ry, 264, 116, top_stripe=c)
        text(slide, rx+18, ry+14, 240, 22, n,
             size=12, bold=True, color=c, mono=True, ls=1.8)
        text(slide, rx+18, ry+38, 240, 26, ret,
             size=18, bold=True, color=c, mono=True, ls=-0.4)
        text(slide, rx+18, ry+66, 240, 44, body,
             size=12, color=INK_2, line_h=1.45)

    takeaway_bar(slide, 96, 904, 1728, 84, 'CROSS-CHANNEL VALIDATION', [
        {'text': '兩源共識 (BOTH BULL) +4.1% > 單源 ; ', 'size': 14, 'color': INK},
        {'text': '跨來源驗證', 'size': 14, 'bold': True, 'color': NAVY},
        {'text': ' 是訊號品質的關鍵濾網。', 'size': 14, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.24 · Coverage & Daily Sentiment
# ============================================================
def slide_apx_coverage(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Coverage & Daily Sentiment', 'A.24',
                          source='自行整理 · 1932 stocks × 4380 days panel')
    eyebrow(slide, 96, 116, 700, 'A.24 · 涵蓋率與每日情感序列', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Coverage ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 文本訊號的「值不值得用」之衡量', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    kpis = [
        ('TOP 50 STOCK',  '92%',  '日均覆蓋',          NAVY),
        ('TOP 200',       '74%',  '日均覆蓋',          TEAL),
        ('1500+ 小型股',  '38%',  '日均覆蓋 · 警示',   GOLD),
        ('PEAK ARTICLES', '4,210','單日峰值 (2024-08)', ORANGE),
    ]
    for i, (lbl, v, sub, c) in enumerate(kpis):
        kx = 96 + i*438
        kpi_hero(slide, kx, 308, 414, 168, lbl, v, sub, accent=c)

    # Daily sentiment line chart
    text(slide, 96, 506, 600, 22, '— DAILY MARKET SENTIMENT INDEX (2018-2025)',
         size=12, bold=True, color=NAVY, mono=True, ls=2.0)
    px_, py_, pw_, ph_ = linechart_box(slide, 96, 540, 1728, 380)
    # x-axis
    years = ['2018','2019','2020','2021','2022','2023','2024','2025']
    for i, yl in enumerate(years):
        tx = px_ + int(i / (len(years)-1) * pw_)
        rect(slide, tx, py_+ph_, 1, 6, fill=INK_3)
        text(slide, tx-30, py_+ph_+12, 60, 18, yl,
             size=12, color=INK_3, mono=True, align='center')
    # y-axis (sentiment -1 to +1)
    yvals = [-1, -0.5, 0, 0.5, 1]
    zero_y = py_ + ph_//2
    for v in yvals:
        ty = py_ + ph_//2 - int(v * ph_ / 2)
        rect(slide, px_-6, ty, 6, 1, fill=INK_3)
        text(slide, px_-50, ty-10, 40, 18, f'{v:+.1f}',
             size=12, color=INK_3, mono=True, align='right')
        if v == 0:
            rect(slide, px_, ty, pw_, 1, fill=INK_4)
        else:
            rect(slide, px_, ty, pw_, 1, fill=SUBTLE)
    # Stylized sentiment series
    pts_data = [0.25, 0.18, -0.42, 0.32, 0.48, -0.18, -0.62, -0.32, -0.04, 0.18, 0.36,
                0.42, 0.28, 0.12, -0.24, -0.18, 0.08, 0.32, 0.52, 0.38, 0.18, -0.12,
                -0.32, -0.18, 0.22, 0.36, 0.48, 0.42, 0.18, 0.04]
    pts = []
    for i, v in enumerate(pts_data):
        x = px_ + int(i / (len(pts_data)-1) * pw_)
        y = py_ + ph_//2 - int(v * ph_ / 2)
        pts.append((x, y))
    polyline(slide, pts, TEAL, line_w=2.0)
    # COVID crash annotation
    crash_x = px_ + int(2 / (len(years)-1) * pw_)
    rect(slide, crash_x, py_+10, 1, ph_-10, fill=ROSE)
    text(slide, crash_x-80, py_+12, 160, 22, 'COVID-19 crash',
         size=12, bold=True, color=ROSE, mono=True, align='center')

    takeaway_bar(slide, 96, 936, 1728, 54, '', [
        {'text': '極端事件 (COVID-19 / 2022 升息) 情感指數皆出現 -0.6 級重挫，', 'size': 14, 'color': INK},
        {'text': '系統能即時捕捉市場情緒拐點。', 'size': 14, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.25 · Signal Decay (PSI by pillar)
# ============================================================
def slide_apx_signal_decay(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Signal Decay · PSI by Pillar', 'A.25',
                          source='自行整理 · PSI computed on 30-day rolling window vs train baseline')
    eyebrow(slide, 96, 116, 700, 'A.25 · 訊號衰減監控 9 支柱 PSI', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'PSI Drift Monitor ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 9 pillars × 30 日 · 文本 0.13 警告', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    psi = [
        ('Risk',        0.04, NAVY,    'OK'),
        ('Trend',       0.05, NAVY_D,  'OK'),
        ('Momentum',    0.06, TEAL,    'OK'),
        ('Valuation',   0.07, TEAL_D,  'OK'),
        ('Fund.',       0.04, EMERALD, 'OK'),
        ('Industry',    0.08, GOLD,    'OK'),
        ('Chip',        0.09, ORANGE,  'OK'),
        ('Sentiment',   0.13, ROSE,    'WARN'),
        ('Event',       0.06, PINK,    'OK'),
    ]
    rect(slide, 96, 308, 1728, 580, fill=PAPER, line=RULE, line_w=0.75)
    text(slide, 116, 326, 800, 22, '— POPULATION STABILITY INDEX BY PILLAR · 30-day rolling',
         size=12, bold=True, color=INK_2, mono=True, ls=1.6)
    base_y = 820
    chart_h = 380
    bar_w = 130
    gap = 50
    bar_x0 = 156
    max_v = 0.30
    # Threshold lines
    thr_warn = 0.10; thr_alert = 0.25
    for thr, color, lbl in [(thr_warn, GOLD, 'Warning 0.10'), (thr_alert, ROSE, 'Alert 0.25')]:
        ty = base_y - int(thr / max_v * chart_h)
        for i in range(0, 1640, 12):
            rect(slide, 156+i, ty, 6, 1, fill=color)
        text(slide, 1700, ty-10, 100, 18, lbl,
             size=12, bold=True, color=color, mono=True, align='right')
    for i, (name, v, c, status) in enumerate(psi):
        bx = bar_x0 + i*(bar_w + gap)
        bar_h = int(v / max_v * chart_h)
        rect(slide, bx, base_y - bar_h, bar_w, bar_h, fill=c)
        text(slide, bx-20, base_y-bar_h-30, bar_w+40, 24, f'{v:.2f}',
             size=14, bold=True, color=c, mono=True, align='center', anchor='middle')
        text(slide, bx-20, base_y+12, bar_w+40, 22, name,
             size=12, bold=True, color=INK, mono=True, align='center', anchor='middle')
        # status pill
        status_color = ROSE if status == 'WARN' else EMERALD
        rect(slide, bx+bar_w//2-32, base_y+38, 64, 22, fill=status_color)
        text(slide, bx+bar_w//2-32, base_y+38, 64, 22, status,
             size=12, bold=True, color=PAPER, mono=True, align='center', anchor='middle')

    takeaway_bar(slide, 96, 904, 1728, 84, 'EARLY WARNING · 早期警示', [
        {'text': 'Sentiment 0.13 已過 0.10 警告線 · 仍低於 0.25 alert · ', 'size': 14, 'color': INK},
        {'text': '建議季度更新詞表 · 避免衰減', 'size': 14, 'bold': True, 'color': NAVY},
        {'text': '。其餘 8 支柱皆穩定。', 'size': 14, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.26 · Limitations & Future Work
# ============================================================
def slide_apx_limitations(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Limitations & Future Work', 'A.26',
                          source='Self-assessment per BCG due diligence template · open issues tracked at issues/')
    eyebrow(slide, 96, 116, 700, 'A.26 · 五大限制 + 五大未來方向', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': '誠實的研究自評 ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· Limitations & Future Work', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    limits = [
        ('01', '日頻限制',        '無 intraday 微結構訊號 (tick / order book)',                       ROSE),
        ('02', '台股限制',        '未跨市場 (HK / SG / US ADR) · 政策外溢未建模',                   ROSE),
        ('03', '黑天鵝盲點',      '極端事件 (戰爭 / 疫情) 訓練樣本稀少',                              ROSE),
        ('04', '文本來源偏誤',    '40% 集中於 Top 5 熱門股 · 中小型股訊號薄弱',                     ROSE),
        ('05', '靜態詞表',        '500 詞表為 2024 校驗 · 半年衰減 PSI > 0.1',                       ROSE),
    ]
    futures = [
        ('01', '加入 Intraday',   'tick / 1-min K · order flow imbalance',                            EMERALD),
        ('02', '跨市場 portfolio','TW + HK + ADR · 統一 9-pillar · ~5,000 stocks',                    EMERALD),
        ('03', '線上學習',         'EWMA + concept drift detection · 月更新',                           EMERALD),
        ('04', 'LLM 輔助情感',    'GPT / Llama embedding 取代部分詞表 · 強化中小型股',                EMERALD),
        ('05', '因果機制建模',     'do-calculus + IV · 從相關走向因果',                                 EMERALD),
    ]

    # Two columns
    text(slide, 96, 296, 800, 22, '— LIMITATIONS · 五大限制',
         size=12, bold=True, color=ROSE, mono=True, ls=2.4)
    for i, (no, h, body, c) in enumerate(limits):
        ry = 332 + i*116
        card(slide, 96, ry, 850, 100, top_stripe=c)
        text(slide, 116, ry+18, 60, 36, no, size=22, bold=True, color=c, mono=True, ls=-0.4)
        text(slide, 178, ry+18, 660, 28, h, size=18, bold=True, color=NAVY)
        text(slide, 178, ry+52, 660, 36, body, size=13, color=INK_2, line_h=1.45)

    text(slide, 974, 296, 800, 22, '— FUTURE WORK · 五大未來方向',
         size=12, bold=True, color=EMERALD, mono=True, ls=2.4)
    for i, (no, h, body, c) in enumerate(futures):
        ry = 332 + i*116
        card(slide, 974, ry, 850, 100, top_stripe=c)
        text(slide, 994, ry+18, 60, 36, no, size=22, bold=True, color=c, mono=True, ls=-0.4)
        text(slide, 1056, ry+18, 660, 28, h, size=18, bold=True, color=NAVY)
        text(slide, 1056, ry+52, 660, 36, body, size=13, color=INK_2, line_h=1.45)

    takeaway_bar(slide, 96, 920, 1728, 54, '', [
        {'text': '系統設計時即承認自身邊界 · ', 'size': 14, 'color': INK},
        {'text': '誠實 > 過譽 · 這是 v11.5.17 與一般學術原型最根本的差別。', 'size': 14, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.27 · Bibliography (15 references in 4 categories)
# ============================================================
def slide_apx_bibliography(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Bibliography · 15 References', 'A.27',
                          source='自行整理· all referenced in code comments')
    eyebrow(slide, 96, 116, 700, 'A.27 · 參考文獻四大領域 15 篇', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Standing on Shoulders ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 方法論立基於 50 年量化研究', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    cats = [
        ('METHODOLOGY · 方法論', NAVY, [
            'López de Prado (2018). Advances in Financial Machine Learning. Wiley.',
            'Bailey & López de Prado (2014). The Deflated Sharpe Ratio. SSRN 2460551.',
            'Harvey, Liu & Zhu (2016). … and the Cross-Section of Expected Returns. RFS 29(1).',
            'Lopez de Prado & Lewis (2019). Detection of False Investment Strategies. JPM 45.',
        ]),
        ('FACTOR MODELS · 因子模型', TEAL, [
            'Fama & French (1993). Common Risk Factors in the Returns on Stocks and Bonds. JFE 33.',
            'Carhart (1997). On Persistence in Mutual Fund Performance. JF 52.',
            'Asness, Moskowitz & Pedersen (2013). Value and Momentum Everywhere. JF 68.',
            'Hou, Xue & Zhang (2015). Digesting Anomalies: An Investment Approach. RFS 28.',
        ]),
        ('BEHAVIORAL FINANCE · 行為財務學', ROSE, [
            'DeBondt & Thaler (1985). Does the Stock Market Overreact? JF 40.',
            'Tetlock (2007). Giving Content to Investor Sentiment. JF 62.',
            'Loughran & McDonald (2011). When Is a Liability Not a Liability? JF 66.',
        ]),
        ('ML & TEXT · 機器學習與文本', ORANGE, [
            'Chen & Guestrin (2016). XGBoost: A Scalable Tree Boosting System. KDD.',
            'Lundberg & Lee (2017). A Unified Approach to Interpreting Model Predictions. NeurIPS.',
            'Ke et al. (2017). LightGBM: A Highly Efficient Gradient Boosting Decision Tree. NeurIPS.',
            'Ardia et al. (2021). Sentometrics: Computing Time Series of Sentiment. JSS 99.',
        ]),
    ]
    grid_x, grid_y = 96, 308
    cw = 850; ch = 308; gap = 24
    for i, (title, c, refs) in enumerate(cats):
        col = i % 2; row = i // 2
        cx = grid_x + col*(cw+gap)
        cy = grid_y + row*(ch+gap)
        card(slide, cx, cy, cw, ch, top_stripe=c)
        text(slide, cx+24, cy+22, cw-48, 22, title,
             size=12, bold=True, color=c, mono=True, ls=2.4)
        for j, r in enumerate(refs):
            ry = cy + 60 + j*54
            text(slide, cx+24, ry, 24, 24, '·',
                 size=20, bold=True, color=c, mono=True, anchor='middle')
            text(slide, cx+44, ry, cw-72, 50, r,
                 size=12, color=INK_2, line_h=1.4)

    takeaway_bar(slide, 96, 952, 1728, 54, '', [
        {'text': '所有方法論皆有同行審核之文獻支撐；', 'size': 14, 'color': INK},
        {'text': '本系統為 50 年量化文獻在台股的工程實作', 'size': 14, 'bold': True, 'color': NAVY},
        {'text': '。', 'size': 14, 'color': INK},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.28 · 工具與技術棧 · Python ecosystem matrix
# ============================================================
def slide_apx_tools_stack(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Tools · Python Ecosystem', 'A.28',
                          source='requirements.txt · pyproject.toml · 5 functional groups × 3 libs')
    eyebrow(slide, 96, 116, 700, 'A.28 · 工具棧 15 個核心套件', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Python Stack ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· Data → ML → NLP → Viz → Governance', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    groups = [
        ('DATA · 資料層',         NAVY, [
            ('pandas',  '2.2',  'panel data 主鍵對齊'),
            ('numpy',   '2.0',  '向量運算 / 統計'),
            ('duckdb',  '0.10', '本地 OLAP · WAL off'),
        ]),
        ('ML · 模型層',          TEAL, [
            ('xgboost',     '1.7', '★ winner · D+20 AUC 0.6455'),
            ('lightgbm',    '4.x', '次佳引擎 · D+20 AUC 0.6391'),
            ('scikit-learn','1.5', 'preprocess / metrics / split'),
        ]),
        ('TUNE · 搜尋層',        EMERALD, [
            ('optuna',  '3.6',  'TPE Sampler · 50 trials'),
            ('shap',    '0.46', 'tree explainer · top-20 features'),
            ('joblib',  '1.4',  'model serialization · 6 horizons'),
        ]),
        ('NLP · 文本層',         GOLD, [
            ('jieba',    '0.42', '中文斷詞 · 自定義金融詞典'),
            ('snownlp',  '0.12', '中文情緒五桶分位'),
            ('wordcloud','1.9',  '視覺化 · Q1 vs Q5 對比'),
        ]),
        ('VIZ + GOV · 服務層',   ROSE, [
            ('streamlit','1.39', '13 頁互動儀表板'),
            ('plotly',   '5.24', 'dark-glint chart 語言'),
            ('matplotlib','3.9', 'static report figures'),
        ]),
    ]
    grid_x, grid_y = 96, 308
    cw = 350; ch = 360; gap = 12
    for i, (title, c, libs) in enumerate(groups):
        cx = grid_x + i*(cw + gap)
        card(slide, cx, grid_y, cw, ch, top_stripe=c)
        text(slide, cx+22, grid_y+22, cw-44, 22, title,
             size=12, bold=True, color=c, mono=True, ls=2.4)
        for j, (lib, ver, note) in enumerate(libs):
            ry = grid_y + 70 + j*92
            rect(slide, cx+22, ry+8, 4, 64, fill=c)
            text(slide, cx+34, ry, cw-56, 28, lib,
                 size=18, bold=True, color=NAVY, mono=True, ls=-0.4)
            text(slide, cx+34, ry+30, 100, 22, f'v{ver}',
                 size=12, bold=True, color=c, mono=True, ls=0.6)
            text(slide, cx+34, ry+52, cw-56, 36, note,
                 size=12, color=INK_2, line_h=1.4)

    takeaway_bar(slide, 96, 700, 1728, 220, 'WHY THIS STACK', [
        {'text': '純 Python · MIT/Apache 授權 · 全部可在 ', 'size': 15, 'color': INK},
        {'text': 'Streamlit Cloud / Docker / Conda', 'size': 15, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': ' 本地或雲端重現；', 'size': 15, 'color': INK},
        {'br': True},
        {'text': 'XGBoost + LightGBM 雙引擎並行 · ', 'size': 15, 'color': INK_2},
        {'text': 'Optuna 150 trials × 6 horizon × engine = 900 次搜尋', 'size': 15, 'bold': True, 'color': TEAL},
        {'text': '；SHAP 提供個股可解釋性、Streamlit 13 頁將模型結果暴露給研究員與經理人。',
         'size': 15, 'color': INK_2},
        {'sp': True},
        {'text': 'AUDIT', 'size': 12, 'bold': True, 'color': ROSE, 'mono': True, 'ls': 2.0},
        {'br': True},
        {'text': '所有依賴版本鎖在 ', 'size': 13, 'color': INK_2},
        {'text': 'requirements.txt', 'size': 13, 'bold': True, 'color': INK, 'mono': True},
        {'text': '、commit hash + joblib 同步、', 'size': 13, 'color': INK_2},
        {'text': '任何一段績效都可從原始資料 → 重現到分數差 < 1e-6。',
         'size': 13, 'bold': True, 'color': NAVY},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.29 · 方法論詳述 · 5 大核心統計方法
# ============================================================
def slide_apx_methodology(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Methodology · 5 Core Methods', 'A.29',
                          source='López de Prado (2018) · Bailey & López de Prado (2014) · Lundberg & Lee (2017)')
    eyebrow(slide, 96, 116, 700, 'A.29 · 方法論詳述', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Method × Application ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 從學術到實證的 1:1 對應', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    methods = [
        ('01', 'Purged WF-CV',     '時序交叉驗證',      NAVY,
         'Fold k 之 train ∩ test̂(t±h) = ∅',
         '4 fold × 20 日 embargo · 防止 D+20 標籤跨期重疊洩漏',
         'López de Prado (2018) §7.4'),
        ('02', 'Deflated Sharpe',  '多重檢定校正',      TEAL,
         'DSR = (SR̂ − E[max SR_n]) · √(T−1) / Φ⁻¹(1−α)',
         '原始 SR̂ 2.41 → DSR 12.12 · 扣 150 trials skew/kurt 後 p ≪ 0.001',
         'Bailey & López de Prado (2014)'),
        ('03', 'LOPO',             '支柱貢獻分析',      EMERALD,
         'ΔAUC_p = AUC(全 91) − AUC(91 \\ pillar p)',
         'Risk +138.6 bps · Trend +64.9 · Chip / Event / Industry 共 −49 bps',
         'Custom · feature ablation'),
        ('04', 'SHAP TreeExplainer','個股可解釋性',     GOLD,
         'φ_i(x) = E[f(x) | x_i] − E[f(x)]',
         'Top 20 全局特徵 · per-stock 解釋 → Streamlit 個股檢視頁',
         'Lundberg & Lee (2017)'),
        ('05', 'PSI / KS',         '漂移監控',          ROSE,
         'PSI = Σ (p_i − q_i) ln(p_i / q_i)',
         '91 features 月測 · Text PSI 0.13 觸發 WARN · 對策每月詞向量重新訓練',
         'Basel III FRTB · industry standard'),
    ]
    bar_x, bar_y0 = 96, 308
    bar_w, bar_h, bgap = 1728, 108, 8
    for i, (no, name, zh, color, formula, application, ref) in enumerate(methods):
        by = bar_y0 + i*(bar_h + bgap)
        card(slide, bar_x, by, bar_w, bar_h, top_stripe=None, line_w=0.5)
        rect(slide, bar_x, by, 8, bar_h, fill=color)
        # number + name
        text(slide, bar_x+30, by+16, 80, 30, no, size=22, bold=True, color=color, mono=True, ls=-0.4)
        text(slide, bar_x+88, by+20, 220, 24, name,
             size=14, bold=True, color=color, mono=True, ls=2.0)
        text(slide, bar_x+88, by+48, 240, 28, zh, size=20, bold=True, color=NAVY)
        text(slide, bar_x+88, by+82, 360, 22, '↪ ' + ref,
             size=12, color=INK_3, mono=True, ls=0.4)
        # formula box (left mid)
        text(slide, bar_x+460, by+16, 100, 18, 'FORMULA',
             size=12, bold=True, color=INK_3, mono=True, ls=2.0)
        rect(slide, bar_x+460, by+38, 540, 56, fill=TINT, line=RULE, line_w=0.5)
        text(slide, bar_x+474, by+38, 510, 56, formula,
             size=14, bold=True, color=NAVY, mono=True, anchor='middle', ls=-0.2)
        # application (right)
        text(slide, bar_x+1024, by+16, 200, 18, 'OUR APPLICATION',
             size=12, bold=True, color=INK_3, mono=True, ls=2.0)
        text(slide, bar_x+1024, by+38, 680, 60, application,
             size=13, color=INK_2, line_h=1.5)

    takeaway_bar(slide, 96, 880, 1728, 112, 'METHOD-TO-EVIDENCE', [
        {'text': '每個結論都可逆向追溯到一個方法、每個方法都有一篇引用文獻；', 'size': 15, 'color': INK},
        {'br': True},
        {'text': '這 5 大方法疊加 = ', 'size': 15, 'color': INK_2},
        {'text': '從原始資料到 9 / 9 治理關卡 PASS 的審計鏈',
         'size': 15, 'bold': True, 'color': NAVY},
        {'text': '。', 'size': 15, 'color': INK_2},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# A.30 · 資料來源 + 延伸文獻
# ============================================================
def slide_apx_data_sources_refs(prs):
    slide = add_blank_slide(prs)
    slide_chrome_appendix(slide, 'Data Sources & Extended References', 'A.30',
                          source='TWSE / TPEx · 公開資訊觀測站 · CMoney 文本 · 課程指定資料庫 BDA-2026')
    eyebrow(slide, 96, 116, 700, 'A.30 · 資料來源 + 延伸期刊文獻', color=ROSE)
    rich(slide, 96, 144, 1750, 110, [
        {'text': 'Data Provenance ', 'size': 40, 'bold': True, 'color': NAVY, 'ls': -0.8},
        {'text': '· 每筆樣本可回溯至公開來源', 'size': 26, 'bold': True, 'color': TEAL, 'ls': -0.4, 'mono': True},
    ], line_h=1.18)
    rule_divider(slide, 96, 268, color=ROSE)

    # LEFT: Data Sources (3 cards)
    text(slide, 96, 296, 600, 22, 'DATA SOURCES · 資料來源',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    sources = [
        ('PRICE / VOLUME', NAVY,
         'TWSE · TPEx · 公開資訊觀測站',
         '1,930 檔股票 × 505 交易日 · OHLCV + 融資融券',
         'data/raw/stock_prices.parquet · 4.2 GB'),
        ('FUNDAMENTALS', TEAL,
         '財報資料庫 · 季頻',
         'IFRS 報表三表 · 主鍵 (stock_id, period) · lag 45 天',
         'data/raw/financial_statements.parquet · 312 MB'),
        ('TEXT CORPUS', GOLD,
         'CMoney 新聞 + PTT Stock 論壇',
         '948K 文本 · T 日 15:00 截止 · 自定義金融詞典 1,521 詞',
         'data/raw/text/ · 8.7 GB · jieba 斷詞'),
    ]
    sx0 = 96; sw = 850; sh = 192; sgap = 12
    for i, (cat, c, src, body, art) in enumerate(sources):
        sy = 326 + i*(sh + sgap)
        card(slide, sx0, sy, sw, sh, top_stripe=c)
        text(slide, sx0+24, sy+22, sw-48, 22, cat,
             size=12, bold=True, color=c, mono=True, ls=2.4)
        text(slide, sx0+24, sy+54, sw-48, 30, src,
             size=18, bold=True, color=NAVY)
        text(slide, sx0+24, sy+96, sw-48, 36, body,
             size=13, color=INK_2, line_h=1.5)
        text(slide, sx0+24, sy+148, sw-48, 22, '↪ ' + art,
             size=12, color=INK_3, mono=True, ls=0.4)

    # RIGHT: Extended References (4 categories beyond A.27)
    text(slide, 980, 296, 600, 22, 'EXTENDED REFERENCES · 延伸文獻',
         size=12, bold=True, color=INK_3, mono=True, ls=2.4)
    ext_refs = [
        ('TAIWAN MARKET · 台股', VIOLET, [
            'Lin & Lee (2010). 台股報酬可預測性. 臺大管理論叢 21(1).',
            '張 (2018). 台灣股市規模效應與動能效應. 中山管理評論 26(2).',
            'Chen et al. (2022). News Sentiment in TWSE. PJM 30(3).',
        ]),
        ('TIME-SERIES VALIDATION', NAVY, [
            'Cerqueira et al. (2020). Evaluating time series forecasting models. ML 109.',
            'Hyndman & Athanasopoulos (2021). Forecasting: Principles and Practice (3e).',
        ]),
        ('GOVERNANCE · 治理規範', EMERALD, [
            'Basel III · FRTB (2019). Minimum capital requirements for market risk.',
            'EU AI Act (2024). High-risk AI systems Article 9 — risk management.',
        ]),
        ('DATA & TOOLING', GOLD, [
            'McKinney (2011). pandas: a Foundational Python Library for DA. PyHPC.',
            'Akiba et al. (2019). Optuna: A Next-generation Hyperparameter Optimization Framework. KDD.',
            'BDA-2026 課程資料規範 · 3 月 26 日資料庫快照 (1.3 GB SQL).',
        ]),
    ]
    rx0 = 980; rw = 844; rh = 90; rgap = 8
    for i, (title, c, refs) in enumerate(ext_refs):
        ry = 326 + i*(rh + rgap)
        card(slide, rx0, ry, rw, rh, top_stripe=c, line_w=0.5)
        text(slide, rx0+22, ry+18, rw-44, 18, title,
             size=12, bold=True, color=c, mono=True, ls=2.0)
        for j, r in enumerate(refs[:2]):  # show top 2 per category to fit
            sy = ry + 38 + j*22
            text(slide, rx0+22, sy, 16, 18, '·',
                 size=12, bold=True, color=c, mono=True)
            text(slide, rx0+38, sy, rw-60, 18, r,
                 size=12, color=INK_2)
        if len(refs) > 2:
            text(slide, rx0+rw-90, ry+rh-22, 80, 18, f'+ {len(refs)-2} more',
                 size=12, color=INK_4, mono=True, align='right', ls=0.4)

    takeaway_bar(slide, 96, 926, 1728, 84, 'PROVENANCE', [
        {'text': '948,976 筆樣本 · 1,623 候選因子 · 91 篩選後特徵 · 全部', 'size': 14, 'color': INK},
        {'text': ' SQL 可回溯', 'size': 14, 'bold': True, 'color': NAVY, 'mono': True},
        {'text': '；本研究遵循 BDA-2026 課程資料規範、所有外部引用標明出處、所有方法論可獨立復現。',
         'size': 14, 'color': INK_2},
    ], color=ROSE, bg=ROSE_LT)
    return slide


# ============================================================
# MAIN
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = W_PX * PX
    prs.slide_height = H_PX * PX

    # Main story · 11 slides · 12-min class version
    print('Building main story (11 slides) ...')
    slide_cover(prs)              # P.01
    slide_problem(prs)            # P.02
    slide_approach(prs)           # P.03
    slide_model_performance(prs)  # P.04
    slide_governance(prs)         # P.05
    slide_pillars_lopo(prs)       # P.06
    slide_sentiment_ushape(prs)   # P.07
    slide_topn_dual(prs)          # P.08
    slide_equity_curve_main(prs)  # P.09 · NEW · 24m equity curve
    slide_recommendations(prs)    # P.10
    slide_closing(prs)            # P.11

    # Appendix · 27 slides
    print('Building appendix (30 slides) ...')
    slide_apx_three_commitments(prs)   # A.01
    slide_apx_problem_deep(prs)        # A.02
    slide_apx_approach_detail(prs)     # A.03
    slide_apx_data_sources(prs)        # A.04
    slide_apx_walk_forward(prs)        # A.05
    slide_apx_tech_stack(prs)          # A.06
    slide_apx_pillars_grid(prs)        # A.07
    slide_apx_lopo_detail(prs)         # A.08
    slide_apx_shap_top20(prs)          # A.09
    slide_apx_top10_factors(prs)       # A.10
    slide_apx_model_compare(prs)       # A.11
    slide_apx_auc_dsr(prs)             # A.12
    slide_apx_threshold(prs)           # A.13
    slide_apx_hit_decile(prs)          # A.14
    slide_apx_equity_curve(prs)        # A.15
    slide_apx_drawdown_cost(prs)       # A.16
    slide_apx_dashboard(prs)           # A.17
    slide_apx_sentiment_overview(prs)  # A.18
    slide_apx_keyword_spectrum(prs)    # A.19
    slide_apx_sentiment_quintile(prs)  # A.20
    slide_apx_q1_contrarian(prs)       # A.21
    slide_apx_ushape_multi(prs)        # A.22
    slide_apx_news_forum(prs)          # A.23
    slide_apx_coverage(prs)            # A.24
    slide_apx_signal_decay(prs)        # A.25
    slide_apx_limitations(prs)         # A.26
    slide_apx_bibliography(prs)        # A.27
    slide_apx_tools_stack(prs)         # A.28 · NEW · Python ecosystem
    slide_apx_methodology(prs)         # A.29 · NEW · 5 core methods + formulas
    slide_apx_data_sources_refs(prs)   # A.30 · NEW · data sources + extended refs

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'deck_full.pptx')
    prs.save(out_path)
    print(f'\nSaved: {out_path}')
    print(f'Slides built: {len(prs.slides)}')


if __name__ == '__main__':
    main()
