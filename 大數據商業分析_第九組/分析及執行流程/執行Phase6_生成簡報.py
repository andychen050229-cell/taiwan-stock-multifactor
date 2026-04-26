"""Generate a McKinsey-style consulting PPTX deck for the final presentation.

Deck structure (MECE, action-title format):
  S1.  Cover
  S2.  Executive Summary (1-slide BLUF)
  S3.  Agenda
  S4.  §1 Situation — 問題背景：台股散戶訊息落差
  S5.  §1.1 Complication — 三個現象缺口
  S6.  §2 Resolution — 解決方案總覽（六階段架構）
  S7.  §3.1 Phase 1 — 資料 · 9 支柱特徵工程
  S8.  §3.2 Phase 2 — 雙引擎 + Walk-Forward CV
  S9.  §3.3 Phase 3 — 擴充分析（成本/地平線/個案）
  S10. §3.4 Phase 4 — 模型治理（Model Card/DSR）
  S11. §3.5 Phase 5 — 文本 + 情緒因子
  S12. §3.6 Phase 6 — 深度驗證（LOPO/Threshold/個股）
  S13. §4 Evidence — 核心數據證據
  S14. §4.1 LOPO 支柱貢獻
  S15. §4.2 Threshold Sweep
  S16. §4.3 2454 聯發科個股
  S17. §5 Product — Streamlit 雙面板儀表板
  S18. §6 Governance — 9/9 Quality Gates 全過
  S19. §7 Limitations — 現階段局限
  S20. §8 Next Steps — 下一步
  S21. §9 Takeaways — 3 句核心訊息
  S22. Q&A / Thank You
"""
import io
import sys
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Force utf-8 stdout on Windows cp950
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

# ============================================================================
# Brand system (matches dashboard glint-light)
# ============================================================================
BRAND = {
    "ink":      RGBColor(0x0F, 0x17, 0x2A),  # near-black
    "ink_2":    RGBColor(0x47, 0x55, 0x69),  # slate
    "ink_3":    RGBColor(0x94, 0xA3, 0xB8),  # tertiary
    "canvas":   RGBColor(0xFA, 0xFB, 0xFC),  # off-white
    "paper":    RGBColor(0xFF, 0xFF, 0xFF),
    "subtle":   RGBColor(0xF1, 0xF5, 0xF9),
    "border":   RGBColor(0xE2, 0xE8, 0xF0),

    "blue":     RGBColor(0x25, 0x63, 0xEB),
    "violet":   RGBColor(0x7C, 0x3A, 0xED),
    "cyan":     RGBColor(0x06, 0xB6, 0xD4),
    "emerald":  RGBColor(0x10, 0xB9, 0x81),
    "amber":    RGBColor(0xF5, 0x9E, 0x0B),
    "rose":     RGBColor(0xF4, 0x3F, 0x5E),
    "indigo":   RGBColor(0x4F, 0x46, 0xE5),
}

FONT_TITLE = "Inter"       # falls back to Calibri if Inter not installed
FONT_BODY  = "Inter"
FONT_MONO  = "Consolas"
FONT_ZH    = "Microsoft JhengHei"  # Traditional Chinese


# ============================================================================
# Paths
# ============================================================================
ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "outputs" / "figures"
OUT  = ROOT / "進度報告" / "台股多因子預測系統_簡報_v4.pptx"


# ============================================================================
# Helpers
# ============================================================================
def set_text(tf, text, *, size=16, bold=False, color=BRAND["ink"],
             font=FONT_ZH, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Write text into a textframe with styling — clears any existing content."""
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    # Also set east-asian font for CJK
    rPr = r._r.get_or_add_rPr()
    rFont = etree.SubElement(rPr, qn("a:ea"))
    rFont.set("typeface", FONT_ZH)


def add_paragraph(tf, text, *, size=14, bold=False, color=BRAND["ink_2"],
                  bullet=False, indent=0, space_before=4, font=FONT_BODY,
                  align=None):
    """Append a paragraph to an existing textframe."""
    p = tf.add_paragraph()
    if space_before:
        p.space_before = Pt(space_before)
    p.level = indent
    if align is not None:
        p.alignment = align
    if bullet:
        # add bullet via XML
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "■")
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "Arial")
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgb = etree.SubElement(buClr, qn("a:srgbClr"))
        srgb.set("val", "2563EB")
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    rFont = etree.SubElement(rPr, qn("a:ea"))
    rFont.set("typeface", FONT_ZH)
    return p


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75):
    """Add a styled rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def add_accent_bar(slide, x, y, w, h, color=BRAND["blue"]):
    """Thin horizontal accent bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_tech_grid_bg(slide, prs):
    """Subtle tech-grid background (light)."""
    sw, sh = prs.slide_width, prs.slide_height
    # very pale canvas
    bg = add_rect(slide, 0, 0, sw, sh, fill=BRAND["canvas"])
    # top-left glow
    glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, -Inches(2), -Inches(2),
                                    Inches(5), Inches(5))
    glow1.shadow.inherit = False
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = BRAND["blue"]
    glow1.line.fill.background()
    glow1.fill.transparency = 0.94
    return bg


def add_page_chrome(slide, prs, slide_num, total, section_label=""):
    """Add the standard header/footer chrome (McKinsey-style)."""
    sw = prs.slide_width
    # Top thin accent
    add_accent_bar(slide, 0, 0, sw, Emu(50000), color=BRAND["blue"])
    # Footer: source line + page num
    footer_y = prs.slide_height - Inches(0.36)
    footer_box = slide.shapes.add_textbox(Inches(0.4), footer_y,
                                          sw - Inches(0.8), Inches(0.3))
    ftf = footer_box.text_frame
    ftf.margin_left = ftf.margin_right = 0
    ftf.margin_top = ftf.margin_bottom = 0
    set_text(ftf, section_label, size=8, color=BRAND["ink_3"], font=FONT_BODY)
    # Right-aligned page number
    pg_box = slide.shapes.add_textbox(sw - Inches(1.0), footer_y,
                                       Inches(0.8), Inches(0.3))
    ptf = pg_box.text_frame
    ptf.margin_left = ptf.margin_right = 0
    set_text(ptf, f"{slide_num} / {total}", size=8, color=BRAND["ink_3"],
             font=FONT_MONO, align=PP_ALIGN.RIGHT)


def add_action_title(slide, prs, title, *, color=BRAND["ink"]):
    """Consulting-style action title (one-sentence insight, large left-aligned)."""
    sw = prs.slide_width
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.45),
                                    sw - Inches(1.0), Inches(0.8))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    set_text(tf, title, size=22, bold=True, color=color, font=FONT_ZH)


def add_subtitle(slide, prs, subtitle):
    """Sub-line under the action title (italic grey)."""
    sw = prs.slide_width
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.15),
                                    sw - Inches(1.0), Inches(0.35))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    set_text(tf, subtitle, size=12, color=BRAND["ink_2"], font=FONT_ZH)


def add_section_tag(slide, prs, tag):
    """Small uppercase tag (top-right)."""
    sw = prs.slide_width
    box = slide.shapes.add_textbox(sw - Inches(2.5), Inches(0.25),
                                    Inches(2.2), Inches(0.3))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    set_text(tf, tag, size=9, bold=True, color=BRAND["blue"],
             font=FONT_MONO, align=PP_ALIGN.RIGHT)


# ============================================================================
# Slide builders
# ============================================================================
def build_cover(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sw, sh = prs.slide_width, prs.slide_height
    add_tech_grid_bg(slide, prs)

    # Left accent bar (vertical gradient)
    for i, y in enumerate([Inches(0), Inches(sh.inches / 3), Inches(2 * sh.inches / 3)]):
        colors = [BRAND["cyan"], BRAND["blue"], BRAND["violet"]]
        bar = add_rect(slide, 0, y, Inches(0.18), Inches(sh.inches / 3),
                       fill=colors[i])

    # Eyebrow
    eb = slide.shapes.add_textbox(Inches(0.8), Inches(1.2),
                                   Inches(8), Inches(0.4))
    set_text(eb.text_frame, "BIG DATA & BUSINESS ANALYTICS · FINAL PROJECT",
             size=11, bold=True, color=BRAND["blue"], font=FONT_MONO)

    # Main title
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9),
                                     sw - Inches(1.6), Inches(2.0))
    tf = tbox.text_frame
    tf.word_wrap = True
    set_text(tf, "台股多因子預測決策輔助系統",
             size=44, bold=True, color=BRAND["ink"], font=FONT_ZH)
    add_paragraph(tf, "以九支柱因子工程 + 雙引擎機器學習，嚴格驗證下的月度 Alpha",
                  size=18, color=BRAND["ink_2"], space_before=12)

    # Subtitle chips
    chip_y = Inches(4.7)
    chip_txts = [
        ("948K · samples",      BRAND["blue"]),
        ("9 · pillars",         BRAND["violet"]),
        ("1,623 → 91 · feats",  BRAND["cyan"]),
        ("9 / 9 · gates",       BRAND["emerald"]),
        ("LOPO · validated",    BRAND["rose"]),
    ]
    cx = Inches(0.8)
    for txt, c in chip_txts:
        chip_w = Inches(1.45 if len(txt) < 20 else 1.8)
        chip = add_rect(slide, cx, chip_y, chip_w, Inches(0.38),
                        fill=BRAND["paper"], line=c, line_w=1.0)
        ctb = slide.shapes.add_textbox(cx, chip_y + Inches(0.07),
                                        chip_w, Inches(0.3))
        set_text(ctb.text_frame, txt, size=9, bold=True, color=c,
                 font=FONT_MONO, align=PP_ALIGN.CENTER)
        cx += chip_w + Inches(0.12)

    # Author / date
    ab = slide.shapes.add_textbox(Inches(0.8), sh - Inches(1.2),
                                   sw - Inches(1.6), Inches(0.8))
    atf = ab.text_frame
    set_text(atf, f"陳延逸  ·  大數據與商業分析  ·  {datetime.now().strftime('%Y-%m-%d')}",
             size=11, color=BRAND["ink_2"], font=FONT_ZH)
    add_paragraph(atf, "Streamlit Dashboard · https://github.com/andychen050229-cell/taiwan-stock-multifactor",
                  size=9, color=BRAND["ink_3"], font=FONT_MONO, space_before=4)

    return slide


def build_executive_summary(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="Executive Summary · 專案 BLUF (Bottom Line Up Front)")
    add_section_tag(slide, prs, "S2 · EXECUTIVE SUMMARY")
    add_action_title(slide, prs,
                     "月頻 D+20 是唯一通過「嚴格驗證」的可交易地平線；風險支柱貢獻最大")

    sw, sh = prs.slide_width, prs.slide_height
    # 3-column insight cards
    col_w = (sw - Inches(1.2)) / 3
    col_h = Inches(4.3)
    col_y = Inches(1.7)

    cards = [
        {
            "tag":   "PROBLEM",
            "color": BRAND["rose"],
            "title": "台股散戶的決策落差",
            "bullets": [
                "1,932 檔個股，資訊量 >> 決策能力",
                "人工盯盤難跨 9 個因子面向整合判讀",
                "基本面 / 技術面 / 籌碼 / 文本情緒分散",
                "缺乏可回看、可解釋的驗證機制",
            ],
        },
        {
            "tag":   "APPROACH",
            "color": BRAND["blue"],
            "title": "九支柱 + 雙引擎 + LOPO 驗證",
            "bullets": [
                "948K 樣本 × 1,623 候選 → 91 特徵",
                "LightGBM + XGBoost 雙引擎 Ensemble",
                "Purged Walk-Forward CV 4 折 · 20 日 Embargo",
                "Phase 6 LOPO + Threshold + 個股三重驗證",
            ],
        },
        {
            "tag":   "RESULT",
            "color": BRAND["emerald"],
            "title": "月頻 Alpha 穩定且可出手",
            "bullets": [
                "xgboost_D20 · AUC_macro 0.649",
                "Top 0.1% 精度 45.8% · edge +19.5pp",
                "閾值 t=0.40 · call 8.8% · edge +3.14pp",
                "2454 聯發科 OOS 命中 58.6% vs 48.8%",
            ],
        },
    ]

    for i, c in enumerate(cards):
        x = Inches(0.4) + i * col_w + Inches(0.1) if i > 0 else Inches(0.4)
        if i > 0:
            x = Inches(0.4) + i * (col_w + Inches(0.1))
        # Card
        add_rect(slide, x, col_y, col_w, col_h, fill=BRAND["paper"],
                 line=BRAND["border"], line_w=0.75)
        # Top accent
        add_accent_bar(slide, x, col_y, col_w, Emu(40000), color=c["color"])
        # Tag
        tag_box = slide.shapes.add_textbox(x + Inches(0.25), col_y + Inches(0.2),
                                            col_w - Inches(0.5), Inches(0.3))
        set_text(tag_box.text_frame, c["tag"], size=10, bold=True,
                 color=c["color"], font=FONT_MONO)
        # Title
        t_box = slide.shapes.add_textbox(x + Inches(0.25), col_y + Inches(0.55),
                                          col_w - Inches(0.5), Inches(0.7))
        set_text(t_box.text_frame, c["title"], size=16, bold=True,
                 color=BRAND["ink"], font=FONT_ZH)
        # Bullets
        b_box = slide.shapes.add_textbox(x + Inches(0.25), col_y + Inches(1.35),
                                          col_w - Inches(0.5), col_h - Inches(1.6))
        btf = b_box.text_frame
        btf.word_wrap = True
        first = True
        for bullet in c["bullets"]:
            if first:
                set_text(btf, bullet, size=11, color=BRAND["ink_2"], font=FONT_ZH)
                first = False
            else:
                add_paragraph(btf, bullet, size=11, color=BRAND["ink_2"],
                              space_before=8)

    # Bottom takeaway strip
    strip_y = col_y + col_h + Inches(0.25)
    add_rect(slide, Inches(0.4), strip_y, sw - Inches(0.8), Inches(0.6),
             fill=BRAND["subtle"], line=BRAND["border"])
    strip_tb = slide.shapes.add_textbox(Inches(0.55), strip_y + Inches(0.1),
                                         sw - Inches(1.1), Inches(0.4))
    set_text(strip_tb.text_frame,
             "✓  9/9 品質門控全通過   ✓  LOPO 顯示 risk/trend 支柱不可或缺   "
             "✓  Streamlit 雙面板（投資解讀 + 量化工作台）",
             size=11, bold=True, color=BRAND["ink"], font=FONT_ZH)


def build_agenda(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total, section_label="Agenda · 議程")
    add_section_tag(slide, prs, "S3 · AGENDA")
    add_action_title(slide, prs, "今日議程：從「問題」出發，以「證據」收斂")

    sw = prs.slide_width
    items = [
        ("1", "Situation",   "問題背景與複雜性"),
        ("2", "Resolution",  "六階段研究設計"),
        ("3", "Pipeline",    "Phase 1-6 方法論逐段"),
        ("4", "Evidence",    "LOPO · Threshold · 個股"),
        ("5", "Product",     "Streamlit 雙面板儀表板"),
        ("6", "Governance",  "治理、限制、下一步"),
    ]
    x0 = Inches(0.6)
    y0 = Inches(1.8)
    row_h = Inches(0.78)
    for i, (num, en, zh) in enumerate(items):
        y = y0 + i * row_h
        # circle
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x0, y, Inches(0.5), Inches(0.5))
        c.shadow.inherit = False
        c.fill.solid()
        c.fill.fore_color.rgb = BRAND["blue"]
        c.line.fill.background()
        ntf = c.text_frame
        ntf.margin_left = ntf.margin_right = 0
        ntf.margin_top = ntf.margin_bottom = 0
        set_text(ntf, num, size=16, bold=True, color=BRAND["paper"],
                 font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # en title
        entb = slide.shapes.add_textbox(x0 + Inches(0.7), y,
                                         Inches(2.5), Inches(0.4))
        set_text(entb.text_frame, en, size=18, bold=True,
                 color=BRAND["ink"], font=FONT_BODY)
        # zh subtitle
        zhtb = slide.shapes.add_textbox(x0 + Inches(3.2), y + Inches(0.05),
                                         Inches(5), Inches(0.4))
        set_text(zhtb.text_frame, zh, size=14, color=BRAND["ink_2"],
                 font=FONT_ZH)


def build_situation(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§1 Situation · 問題背景")
    add_section_tag(slide, prs, "S4 · SITUATION")
    add_action_title(slide, prs,
                     "台股散戶面臨「資訊量 >> 決策能力」的結構性落差")
    add_subtitle(prs=prs, slide=slide,
                 subtitle="1,932 家上市櫃 × 多面向因子 × 即時資訊 → 需要可解釋、可回看、可治理的決策輔助")

    sw = prs.slide_width
    # Three pain points
    pains = [
        ("📊", "資訊過載",
         "單一投資人難同時監控 1,932 檔、9 大因子面向、每日事件輿情",
         BRAND["blue"]),
        ("🎯", "無一致判讀",
         "技術派 / 基本派 / 籌碼派 / 消息派各自為戰，難跨維度整合",
         BRAND["violet"]),
        ("🔍", "缺驗證工具",
         "市場流通策略多為「事後貼圖」，缺乏可回看、可解釋、可治理的證據",
         BRAND["rose"]),
    ]
    x0 = Inches(0.5)
    y0 = Inches(1.9)
    card_w = (sw - Inches(1.0) - Inches(0.4)) / 3
    card_h = Inches(3.3)
    for i, (ic, t, d, c) in enumerate(pains):
        x = x0 + i * (card_w + Inches(0.2))
        # Card
        add_rect(slide, x, y0, card_w, card_h, fill=BRAND["paper"],
                 line=BRAND["border"])
        add_accent_bar(slide, x, y0, card_w, Emu(35000), color=c)
        # icon
        itb = slide.shapes.add_textbox(x + Inches(0.3), y0 + Inches(0.25),
                                        card_w - Inches(0.6), Inches(0.7))
        set_text(itb.text_frame, ic, size=40, color=c, font=FONT_BODY)
        # title
        ttb = slide.shapes.add_textbox(x + Inches(0.3), y0 + Inches(1.1),
                                        card_w - Inches(0.6), Inches(0.5))
        set_text(ttb.text_frame, t, size=18, bold=True,
                 color=BRAND["ink"], font=FONT_ZH)
        # desc
        dtb = slide.shapes.add_textbox(x + Inches(0.3), y0 + Inches(1.75),
                                        card_w - Inches(0.6), Inches(1.4))
        set_text(dtb.text_frame, d, size=12, color=BRAND["ink_2"], font=FONT_ZH)


def build_resolution_overview(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§2 Resolution · 研究設計")
    add_section_tag(slide, prs, "S5 · RESOLUTION")
    add_action_title(slide, prs,
                     "以六階段管線建立「可回看 + 可解釋 + 可治理」的研究框架")

    sw = prs.slide_width
    # 6 phase blocks in a row
    phases = [
        ("P1", "資料與特徵", BRAND["cyan"]),
        ("P2", "模型與回測", BRAND["blue"]),
        ("P3", "擴充分析",   BRAND["indigo"]),
        ("P4", "治理監控",   BRAND["violet"]),
        ("P5", "文本情緒",   BRAND["amber"]),
        ("P6", "深度驗證",   BRAND["rose"]),
    ]
    x0 = Inches(0.5)
    y0 = Inches(2.0)
    w = (sw - Inches(1.0) - Inches(0.5)) / 6
    h = Inches(1.5)
    for i, (p, t, c) in enumerate(phases):
        x = x0 + i * (w + Inches(0.1))
        add_rect(slide, x, y0, w, h, fill=BRAND["paper"], line=c, line_w=1.5)
        # P tag
        ptb = slide.shapes.add_textbox(x + Inches(0.1), y0 + Inches(0.15),
                                        w - Inches(0.2), Inches(0.4))
        set_text(ptb.text_frame, p, size=14, bold=True, color=c, font=FONT_MONO)
        # Title
        ttb = slide.shapes.add_textbox(x + Inches(0.1), y0 + Inches(0.65),
                                        w - Inches(0.2), Inches(0.7))
        set_text(ttb.text_frame, t, size=12, bold=True,
                 color=BRAND["ink"], font=FONT_ZH)
        # Arrow to next
        if i < len(phases) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          x + w, y0 + h / 2 - Inches(0.08),
                                          Inches(0.1), Inches(0.16))
            arr.shadow.inherit = False
            arr.fill.solid()
            arr.fill.fore_color.rgb = BRAND["ink_3"]
            arr.line.fill.background()

    # Bottom outputs summary
    out_y = Inches(4.0)
    add_rect(slide, Inches(0.5), out_y, sw - Inches(1.0), Inches(1.8),
             fill=BRAND["subtle"], line=BRAND["border"])
    otb = slide.shapes.add_textbox(Inches(0.8), out_y + Inches(0.2),
                                    sw - Inches(1.6), Inches(1.4))
    otf = otb.text_frame
    otf.word_wrap = True
    set_text(otf, "📦 成果交付", size=12, bold=True,
             color=BRAND["blue"], font=FONT_MONO)
    add_paragraph(otf,
                  "Streamlit 雙面板儀表板（投資解讀 / 量化工作台，共 11 頁）  ·  "
                  "Phase 2-6 品質門控報告 JSON 20+ 份  ·  60+ 張視覺化圖表",
                  size=12, color=BRAND["ink_2"], font=FONT_ZH, space_before=8)
    add_paragraph(otf,
                  "終版研究報告 (docx · 70+ 頁)  ·  簡報 (本份)  ·  "
                  "完整程式碼 GitHub 開源",
                  size=12, color=BRAND["ink_2"], font=FONT_ZH, space_before=6)


def build_phase_slide(prs, slide_num, total, *, tag, phase_num, zh_title,
                      en_subtitle, headline, bullets, color, fig_path=None):
    """Generic Phase slide template."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label=f"§3.{phase_num} Phase {phase_num} · {zh_title}")
    add_section_tag(slide, prs, tag)
    add_action_title(slide, prs, headline, color=BRAND["ink"])
    add_subtitle(prs=prs, slide=slide, subtitle=en_subtitle)

    sw, sh = prs.slide_width, prs.slide_height
    # Left: bullets (60%)
    left_w = sw * 0.55
    # Bullet area
    btb = slide.shapes.add_textbox(Inches(0.55), Inches(1.8),
                                    left_w - Inches(0.8), sh - Inches(2.4))
    btf = btb.text_frame
    btf.word_wrap = True
    first = True
    for bt in bullets:
        if isinstance(bt, tuple):
            head, detail = bt
            if first:
                set_text(btf, head, size=14, bold=True, color=color,
                         font=FONT_ZH)
                first = False
            else:
                add_paragraph(btf, head, size=14, bold=True, color=color,
                              font=FONT_ZH, space_before=14)
            add_paragraph(btf, detail, size=11, color=BRAND["ink_2"],
                          font=FONT_ZH, space_before=2, indent=1)
        else:
            if first:
                set_text(btf, bt, size=13, color=BRAND["ink_2"], font=FONT_ZH)
                first = False
            else:
                add_paragraph(btf, bt, size=13, color=BRAND["ink_2"],
                              font=FONT_ZH, space_before=8)

    # Right: figure (40%)
    if fig_path and Path(fig_path).exists():
        right_x = left_w + Inches(0.1)
        right_w = sw - left_w - Inches(0.6)
        right_y = Inches(1.8)
        right_h = sh - Inches(2.4)
        slide.shapes.add_picture(str(fig_path), right_x, right_y,
                                  width=right_w, height=right_h)


def build_evidence_lopo(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§4.1 Evidence · LOPO 支柱貢獻")
    add_section_tag(slide, prs, "S12 · LOPO VALIDATION")
    add_action_title(slide, prs,
                     "風險支柱是最不可或缺的 ─ 拿掉會掉 +1.39 bps AUC_macro")

    sw, sh = prs.slide_width, prs.slide_height

    # Left: figure
    fig = FIGS / "lopo_pillar_contribution_D20.png"
    if fig.exists():
        slide.shapes.add_picture(str(fig), Inches(0.5), Inches(1.7),
                                  width=Inches(6.2), height=Inches(4.3))

    # Right: ranked table
    right_x = Inches(7.0)
    right_y = Inches(1.7)
    right_w = sw - right_x - Inches(0.5)
    # Header
    add_rect(slide, right_x, right_y, right_w, Inches(0.45),
             fill=BRAND["subtle"], line=BRAND["border"])
    htb = slide.shapes.add_textbox(right_x + Inches(0.15), right_y + Inches(0.08),
                                    right_w - Inches(0.3), Inches(0.3))
    set_text(htb.text_frame, "LOPO Ranking · xgboost_D20",
             size=11, bold=True, color=BRAND["ink"], font=FONT_MONO)

    # Rows
    rows = [
        ("1", "risk",  "風險",    "+1.39 bps", BRAND["rose"]),
        ("2", "trend", "趨勢",    "+0.65 bps", BRAND["blue"]),
        ("3", "val",   "估值",    "+0.09 bps", BRAND["cyan"]),
        ("4", "txt",   "文本",    "+0.08 bps", BRAND["violet"]),
        ("5", "fund",  "基本面",  "-0.03 bps", BRAND["ink_3"]),
        ("6", "sent",  "情緒",    "-0.03 bps", BRAND["ink_3"]),
        ("7", "ind",   "產業",    "-0.15 bps", BRAND["ink_3"]),
        ("8", "event", "事件",    "-0.16 bps", BRAND["ink_3"]),
        ("9", "chip",  "籌碼",    "-0.19 bps", BRAND["ink_3"]),
    ]
    row_h = Inches(0.33)
    cur_y = right_y + Inches(0.5)
    for rk, p_en, p_zh, delta, c in rows:
        add_rect(slide, right_x, cur_y, right_w, row_h,
                 fill=BRAND["paper"], line=BRAND["border"], line_w=0.4)
        # rank
        rk_tb = slide.shapes.add_textbox(right_x + Inches(0.1), cur_y + Inches(0.05),
                                          Inches(0.3), Inches(0.25))
        set_text(rk_tb.text_frame, rk, size=10, bold=True, color=c,
                 font=FONT_MONO)
        # pillar
        p_tb = slide.shapes.add_textbox(right_x + Inches(0.45), cur_y + Inches(0.05),
                                         Inches(1.3), Inches(0.25))
        set_text(p_tb.text_frame, f"{p_en} · {p_zh}", size=10, bold=True,
                 color=BRAND["ink"], font=FONT_ZH)
        # delta
        d_tb = slide.shapes.add_textbox(right_x + right_w - Inches(1.2),
                                         cur_y + Inches(0.05),
                                         Inches(1.1), Inches(0.25))
        set_text(d_tb.text_frame, delta, size=10, bold=True, color=c,
                 font=FONT_MONO, align=PP_ALIGN.RIGHT)
        cur_y += row_h + Inches(0.04)

    # Bottom callout
    add_rect(slide, Inches(0.5), sh - Inches(1.2), sw - Inches(1.0), Inches(0.65),
             fill=BRAND["subtle"], line=BRAND["border"])
    c_tb = slide.shapes.add_textbox(Inches(0.7), sh - Inches(1.1),
                                     sw - Inches(1.4), Inches(0.5))
    set_text(c_tb.text_frame,
             "📌 風險僅 6 個特徵卻貢獻最大 ΔAUC_up +2.86 bps ─ 證明 downside 保護的邊際價值 > 追 upside",
             size=12, bold=True, color=BRAND["ink"], font=FONT_ZH)


def build_evidence_threshold(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§4.2 Evidence · 閾值敏感度")
    add_section_tag(slide, prs, "S13 · THRESHOLD SWEEP")
    add_action_title(slide, prs,
                     "Top 0.1% 最強訊號精度 45.8%，edge +19.5pp ─ 模型排序可信")

    sw, sh = prs.slide_width, prs.slide_height

    # Figure left
    fig = FIGS / "threshold_sweep_xgb_D20.png"
    if fig.exists():
        slide.shapes.add_picture(str(fig), Inches(0.5), Inches(1.7),
                                  width=Inches(6.5), height=Inches(4.3))

    # Right: three callouts
    right_x = Inches(7.3)
    right_y = Inches(1.8)
    callouts = [
        ("Conservative · t=0.40", "call 8.8%   hit 29.4%   edge +3.14pp",
         "交易頻率低 · 訊號可靠度高", BRAND["blue"]),
        ("Balanced · t=0.35",     "call 15.6%  hit 27.3%   edge +1.04pp",
         "取折中部位 · 日常可用", BRAND["violet"]),
        ("Top 0.1% · 精度峰值",   "n=404     hit 45.8%   edge +19.5pp",
         "模型最有把握的訊號 · 適合高信念部位", BRAND["emerald"]),
    ]
    for i, (t, m, d, c) in enumerate(callouts):
        cy = right_y + i * Inches(1.5)
        add_rect(slide, right_x, cy, sw - right_x - Inches(0.5), Inches(1.35),
                 fill=BRAND["paper"], line=BRAND["border"])
        add_accent_bar(slide, right_x, cy, sw - right_x - Inches(0.5),
                        Emu(30000), color=c)
        ttb = slide.shapes.add_textbox(right_x + Inches(0.2), cy + Inches(0.15),
                                        Inches(4.5), Inches(0.35))
        set_text(ttb.text_frame, t, size=12, bold=True, color=c, font=FONT_ZH)
        mtb = slide.shapes.add_textbox(right_x + Inches(0.2), cy + Inches(0.55),
                                        Inches(4.5), Inches(0.4))
        set_text(mtb.text_frame, m, size=11, color=BRAND["ink"], font=FONT_MONO)
        dtb = slide.shapes.add_textbox(right_x + Inches(0.2), cy + Inches(0.95),
                                        Inches(4.5), Inches(0.35))
        set_text(dtb.text_frame, d, size=10, color=BRAND["ink_2"], font=FONT_ZH)


def build_evidence_2454(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§4.3 Evidence · 2454 聯發科個股")
    add_section_tag(slide, prs, "S14 · SINGLE-STOCK CASE")
    add_action_title(slide, prs,
                     "2454 聯發科 OOS 213 日：58 次出手、命中 58.6%（+9.8pp vs base）")

    sw, sh = prs.slide_width, prs.slide_height

    # Figure
    fig = FIGS / "single_stock_2454_mediatek.png"
    if fig.exists():
        slide.shapes.add_picture(str(fig), Inches(0.5), Inches(1.7),
                                  width=Inches(8.5), height=Inches(4.5))

    # Right KPIs
    right_x = Inches(9.2)
    right_y = Inches(1.8)
    kpis = [
        ("OOS 天數",       "213",      "2024-04 → 2025-03", BRAND["blue"]),
        ("出手次數",       "58",       "@ t=0.35",          BRAND["violet"]),
        ("命中率",         "58.6%",    "+9.8pp vs 48.8%",   BRAND["emerald"]),
        ("2025-01 熱區",   "81.8%",    "11 出手 · 9 命中",  BRAND["rose"]),
    ]
    for i, (lbl, val, sub, c) in enumerate(kpis):
        ky = right_y + i * Inches(1.05)
        add_rect(slide, right_x, ky, sw - right_x - Inches(0.5), Inches(0.95),
                 fill=BRAND["paper"], line=BRAND["border"])
        add_accent_bar(slide, right_x, ky,
                        sw - right_x - Inches(0.5),
                        Emu(25000), color=c)
        ltb = slide.shapes.add_textbox(right_x + Inches(0.15), ky + Inches(0.1),
                                        Inches(3), Inches(0.25))
        set_text(ltb.text_frame, lbl, size=9, bold=True,
                 color=BRAND["ink_3"], font=FONT_ZH)
        vtb = slide.shapes.add_textbox(right_x + Inches(0.15), ky + Inches(0.33),
                                        Inches(3), Inches(0.45))
        set_text(vtb.text_frame, val, size=20, bold=True, color=c,
                 font=FONT_MONO)
        stb = slide.shapes.add_textbox(right_x + Inches(0.15), ky + Inches(0.73),
                                        Inches(3), Inches(0.22))
        set_text(stb.text_frame, sub, size=9, color=BRAND["ink_2"], font=FONT_ZH)


def build_product_dashboard(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§5 Product · Streamlit 雙面板儀表板")
    add_section_tag(slide, prs, "S15 · PRODUCT")
    add_action_title(slide, prs, "雙面板設計：一套研究引擎，兩種閱讀模式")

    sw, sh = prs.slide_width, prs.slide_height

    # Two columns
    col_w = (sw - Inches(1.3)) / 2
    y = Inches(1.9)
    h = Inches(4.5)

    # Left: 投資解讀面板
    add_rect(slide, Inches(0.5), y, col_w, h,
             fill=BRAND["paper"], line=BRAND["amber"], line_w=1.2)
    add_accent_bar(slide, Inches(0.5), y, col_w, Emu(40000),
                    color=BRAND["amber"])
    itb = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.25),
                                    col_w - Inches(0.4), Inches(4))
    itf = itb.text_frame
    itf.word_wrap = True
    set_text(itf, "🌱  投資解讀面板", size=18, bold=True,
             color=BRAND["amber"], font=FONT_ZH)
    add_paragraph(itf, "For 投資人 · 行銷 · 商管讀者", size=10,
                  color=BRAND["ink_3"], font=FONT_BODY, space_before=4)
    add_paragraph(itf, "將研究結果翻譯成可理解的決策資訊：",
                  size=12, color=BRAND["ink_2"], font=FONT_ZH, space_before=12)
    for item in [
        "✓  今日/回看歷史任一日的模型看好名單",
        "✓  公司基本面 / 產業分析 / 成本試算",
        "✓  風險提示 / 技術指標視覺化",
        "✓  五層資訊架構 · 漸進式揭露",
        "✓  4 個 tab（模型判讀 / 公司輪廓 / 成本 / 市場）",
    ]:
        add_paragraph(itf, item, size=11, color=BRAND["ink"], font=FONT_ZH,
                      space_before=6)

    # Right: 量化研究工作台
    right_x = Inches(0.5) + col_w + Inches(0.3)
    add_rect(slide, right_x, y, col_w, h,
             fill=BRAND["paper"], line=BRAND["blue"], line_w=1.2)
    add_accent_bar(slide, right_x, y, col_w, Emu(40000), color=BRAND["blue"])
    ptb = slide.shapes.add_textbox(right_x + Inches(0.2), y + Inches(0.25),
                                    col_w - Inches(0.4), Inches(4))
    ptf = ptb.text_frame
    ptf.word_wrap = True
    set_text(ptf, "⚙️  量化研究工作台", size=18, bold=True,
             color=BRAND["blue"], font=FONT_ZH)
    add_paragraph(ptf, "For 教授 · 評審 · 量化研究者", size=10,
                  color=BRAND["ink_3"], font=FONT_BODY, space_before=4)
    add_paragraph(ptf, "完整呈現方法論、證據、治理：",
                  size=12, color=BRAND["ink_2"], font=FONT_ZH, space_before=12)
    for item in [
        "✓  Model Metrics · ICIR 分析 · Backtest",
        "✓  Feature Analysis · Data Explorer · Text",
        "✓  Model Governance · Signal Monitor",
        "✓  Extended Analytics（Phase 3 四項）",
        "✓  🔭 Phase 6 深度驗證（LOPO/Threshold/2454）",
    ]:
        add_paragraph(ptf, item, size=11, color=BRAND["ink"], font=FONT_ZH,
                      space_before=6)

    # Bottom strip
    strip_y = y + h + Inches(0.2)
    add_rect(slide, Inches(0.5), strip_y, sw - Inches(1.0), Inches(0.5),
             fill=BRAND["subtle"], line=BRAND["border"])
    stb = slide.shapes.add_textbox(Inches(0.7), strip_y + Inches(0.1),
                                    sw - Inches(1.4), Inches(0.35))
    set_text(stb.text_frame,
             "🎨  設計語言：glint-light · Inter + JetBrains Mono · 電光藍/紫漸變 · 資料密度 + 亮色科技感",
             size=11, bold=True, color=BRAND["ink"], font=FONT_ZH)


def build_governance(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§6 Governance · 治理")
    add_section_tag(slide, prs, "S16 · GOVERNANCE")
    add_action_title(slide, prs, "9 項品質門控全數通過 ─ 研究可信度的硬證據")

    sw = prs.slide_width

    gates = [
        ("leakage_scan",          "零洩漏掃描"),
        ("walk_forward_cv",       "時序交叉驗證"),
        ("purge_embargo",         "Purge + Embargo"),
        ("model_comparison",      "雙引擎對照"),
        ("feature_stability",     "特徵穩定性"),
        ("baseline_compare",      "基準對照"),
        ("oos_performance",       "樣本外績效"),
        ("cost_sensitivity",      "成本敏感度"),
        ("lopo_validation",       "LOPO 驗證"),
    ]
    cols = 3
    w = (sw - Inches(1.2)) / cols
    h = Inches(0.95)
    x0 = Inches(0.5)
    y0 = Inches(1.9)
    for i, (en, zh) in enumerate(gates):
        row, col = divmod(i, cols)
        x = x0 + col * (w + Inches(0.1))
        y = y0 + row * (h + Inches(0.15))
        add_rect(slide, x, y, w, h, fill=BRAND["paper"],
                 line=BRAND["emerald"], line_w=1.0)
        # check
        ctb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2),
                                        Inches(0.5), Inches(0.55))
        set_text(ctb.text_frame, "✓", size=24, bold=True,
                 color=BRAND["emerald"], font=FONT_BODY)
        # name
        ntb = slide.shapes.add_textbox(x + Inches(0.7), y + Inches(0.15),
                                        w - Inches(0.9), Inches(0.35))
        set_text(ntb.text_frame, zh, size=13, bold=True,
                 color=BRAND["ink"], font=FONT_ZH)
        # code
        etb = slide.shapes.add_textbox(x + Inches(0.7), y + Inches(0.52),
                                        w - Inches(0.9), Inches(0.3))
        set_text(etb.text_frame, en, size=9, color=BRAND["ink_3"],
                 font=FONT_MONO)


def build_limitations(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§7 Limitations · 限制")
    add_section_tag(slide, prs, "S17 · LIMITATIONS")
    add_action_title(slide, prs, "誠實的邊界：我們不聲稱「賺錢機器」")

    sw = prs.slide_width
    limits = [
        ("資料期間僅 2 年",
         "2023/3-2025/3 僅涵蓋 AI 牛市段落，未經歷完整空頭",
         BRAND["amber"]),
        ("Rank IC 絕對值偏低",
         "D+20 Rank IC +0.015，統計顯著但經濟可用空間需搭配閾值過濾",
         BRAND["amber"]),
        ("文本/情緒貢獻有限",
         "LOPO 顯示 txt/sent 近零，可能受資料量與中文處理限制",
         BRAND["amber"]),
        ("未接盤後即時流程",
         "屬歷史快照研究平台，非即時交易系統",
         BRAND["rose"]),
        ("交易成本假設保守",
         "0.36%-0.73% 手續費已含稅，未考慮滑價與流動性衝擊",
         BRAND["rose"]),
        ("回測不等於實戰",
         "過去回測績效不代表未來報酬，無投資建議意涵",
         BRAND["rose"]),
    ]
    cols = 2
    w = (sw - Inches(1.2)) / cols
    h = Inches(1.4)
    x0 = Inches(0.5)
    y0 = Inches(1.9)
    for i, (t, d, c) in enumerate(limits):
        row, col = divmod(i, cols)
        x = x0 + col * (w + Inches(0.2))
        y = y0 + row * (h + Inches(0.15))
        add_rect(slide, x, y, w, h, fill=BRAND["paper"],
                 line=BRAND["border"])
        add_accent_bar(slide, x, y, Emu(50000), h, color=c)
        # title
        ttb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.2),
                                        w - Inches(0.4), Inches(0.4))
        set_text(ttb.text_frame, t, size=13, bold=True,
                 color=BRAND["ink"], font=FONT_ZH)
        # desc
        dtb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.65),
                                        w - Inches(0.4), Inches(0.7))
        set_text(dtb.text_frame, d, size=11, color=BRAND["ink_2"],
                 font=FONT_ZH)


def build_next_steps(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§8 Next Steps · 下一步")
    add_section_tag(slide, prs, "S18 · NEXT STEPS")
    add_action_title(slide, prs, "下一步：擴資料 × 強文本 × 落地回測")

    sw = prs.slide_width
    nexts = [
        ("Short · 0-3 個月", BRAND["blue"], [
            "擴充資料至 5 年含空頭段落",
            "補 D+1 model cards 完整 6 份",
            "建立 shadow mode 盤後跑線上對照",
        ]),
        ("Medium · 3-6 個月", BRAND["violet"], [
            "換用 CKIP 或 BERT 重做文本向量化",
            "接入籌碼盤中資料（外資/投信）",
            "加入產業中性化與流動性過濾器",
        ]),
        ("Long · 6-12 個月", BRAND["rose"], [
            "與券商 API 整合為可執行信號",
            "導入 regime detection（多頭/空頭切換）",
            "建立付費訂閱模式的 MVP",
        ]),
    ]
    cols = 3
    w = (sw - Inches(1.2)) / cols
    h = Inches(4.5)
    x0 = Inches(0.5)
    y0 = Inches(1.9)
    for i, (hd, c, items) in enumerate(nexts):
        x = x0 + i * (w + Inches(0.15))
        add_rect(slide, x, y0, w, h, fill=BRAND["paper"], line=BRAND["border"])
        add_accent_bar(slide, x, y0, w, Emu(40000), color=c)
        # header
        htb = slide.shapes.add_textbox(x + Inches(0.25), y0 + Inches(0.2),
                                        w - Inches(0.4), Inches(0.4))
        set_text(htb.text_frame, hd, size=14, bold=True, color=c, font=FONT_ZH)
        # items
        itb = slide.shapes.add_textbox(x + Inches(0.25), y0 + Inches(0.75),
                                        w - Inches(0.4), h - Inches(1.0))
        itf = itb.text_frame
        itf.word_wrap = True
        first = True
        for it in items:
            if first:
                set_text(itf, f"▸  {it}", size=12, color=BRAND["ink_2"],
                         font=FONT_ZH)
                first = False
            else:
                add_paragraph(itf, f"▸  {it}", size=12,
                              color=BRAND["ink_2"], font=FONT_ZH,
                              space_before=12)


def build_takeaways(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    add_page_chrome(slide, prs, slide_num, total,
                    section_label="§9 Takeaways · 3 句話帶走")
    add_section_tag(slide, prs, "S19 · TAKEAWAYS")
    add_action_title(slide, prs, "三句話帶走這份研究")

    sw, sh = prs.slide_width, prs.slide_height
    takeaways = [
        ("1",
         "月頻 D+20 是唯一「嚴格驗證後仍可用」的地平線",
         "AUC 0.649 · Top 0.1% 精度 45.8% · 9/9 品質門控",
         BRAND["blue"]),
        ("2",
         "風險支柱比 AI 熱門話題更有說服力",
         "LOPO 顯示 risk 僅 6 個特徵卻貢獻最大 ΔAUC_up",
         BRAND["violet"]),
        ("3",
         "研究的價值在「可回看、可解釋、可治理」",
         "雙面板儀表板 + Model Card + LOPO 把黑盒變成白盒",
         BRAND["emerald"]),
    ]
    y0 = Inches(1.9)
    row_h = Inches(1.55)
    for i, (num, hd, sub, c) in enumerate(takeaways):
        y = y0 + i * (row_h + Inches(0.1))
        # big num
        ntb = slide.shapes.add_textbox(Inches(0.6), y, Inches(1.2), row_h)
        set_text(ntb.text_frame, num, size=72, bold=True, color=c,
                 font=FONT_MONO)
        # card
        add_rect(slide, Inches(1.9), y + Inches(0.1),
                 sw - Inches(2.5), row_h - Inches(0.2),
                 fill=BRAND["paper"], line=BRAND["border"])
        add_accent_bar(slide, Inches(1.9), y + Inches(0.1),
                        Emu(50000), row_h - Inches(0.2), color=c)
        # head
        htb = slide.shapes.add_textbox(Inches(2.15), y + Inches(0.3),
                                        sw - Inches(2.8), Inches(0.5))
        set_text(htb.text_frame, hd, size=17, bold=True,
                 color=BRAND["ink"], font=FONT_ZH)
        # sub
        stb = slide.shapes.add_textbox(Inches(2.15), y + Inches(0.85),
                                        sw - Inches(2.8), Inches(0.5))
        set_text(stb.text_frame, sub, size=12,
                 color=BRAND["ink_2"], font=FONT_ZH)


def build_thank_you(prs, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_grid_bg(slide, prs)
    sw, sh = prs.slide_width, prs.slide_height

    # Centered content
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5),
                                   sw - Inches(1.0), Inches(1.2))
    tf = tb.text_frame
    set_text(tf, "Thank You", size=60, bold=True, color=BRAND["ink"],
             font=FONT_BODY, align=PP_ALIGN.CENTER)
    add_paragraph(tf, "歡迎指教 · Q & A", size=22, color=BRAND["ink_2"],
                  font=FONT_ZH, space_before=10)
    p = tf.paragraphs[-1]
    p.alignment = PP_ALIGN.CENTER

    # Contact
    ctb = slide.shapes.add_textbox(Inches(0.5), sh - Inches(1.6),
                                    sw - Inches(1.0), Inches(0.8))
    ctf = ctb.text_frame
    set_text(ctf, "📊  Dashboard · streamlit run 程式碼/儀表板/app.py",
             size=11, color=BRAND["ink_2"], font=FONT_MONO,
             align=PP_ALIGN.CENTER)
    add_paragraph(ctf,
                  "📄  Final Report · 進度報告/專案報告_最新版_v4.docx",
                  size=11, color=BRAND["ink_2"], font=FONT_MONO,
                  space_before=4)
    p = ctf.paragraphs[-1]
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(ctf,
                  "🔗  GitHub · andychen050229-cell/taiwan-stock-multifactor",
                  size=11, color=BRAND["blue"], font=FONT_MONO,
                  space_before=4)
    p = ctf.paragraphs[-1]
    p.alignment = PP_ALIGN.CENTER


# ============================================================================
# Main
# ============================================================================
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Count slides upfront for total (actual final count)
    TOTAL = 20

    # S1 Cover
    build_cover(prs, TOTAL)

    # S2 Exec Summary
    build_executive_summary(prs, 2, TOTAL)

    # S3 Agenda
    build_agenda(prs, 3, TOTAL)

    # S4 Situation
    build_situation(prs, 4, TOTAL)

    # S5 Complication (problem framing cross-slide)
    # (integrated into Resolution overview)

    # S5 Resolution overview
    build_resolution_overview(prs, 5, TOTAL)

    # S6-S11 Phase 1-6
    build_phase_slide(
        prs, 6, TOTAL,
        tag="S6 · PHASE 1",
        phase_num=1,
        zh_title="資料與特徵工程",
        en_subtitle="FinMind API → 948K samples · 9 pillars · 1,623 → 91 features (IC/MI/VIF)",
        headline="以 9 支柱理論框架把 1,623 個候選因子收斂到 91 個可用特徵",
        color=BRAND["cyan"],
        bullets=[
            ("資料源", "FinMind API + 選用資料集 · 1,930 家上市櫃 · 2023/3 – 2025/3"),
            ("9 支柱架構", "trend / fund / val / event / risk / chip / ind / txt / sent · 每支柱有清楚定義邊界"),
            ("三階段篩選", "① IC prescreen（去噪）② Chi²/MI（非線性關聯）③ VIF（去共線）"),
            ("PIT 嚴格處理", "財報與事件資料強制未來資訊封存，從 feature store 建構即避免洩漏"),
        ],
        fig_path=None,
    )
    build_phase_slide(
        prs, 7, TOTAL,
        tag="S7 · PHASE 2",
        phase_num=2,
        zh_title="模型訓練與回測驗證",
        en_subtitle="LightGBM + XGBoost · Purged Walk-Forward CV 4 folds · embargo 20d",
        headline="雙引擎 + 嚴格時序 CV 確保「未來不能看過去」",
        color=BRAND["blue"],
        bullets=[
            ("雙引擎", "LightGBM + XGBoost · 各自超參 Optuna HPO · 等權 Ensemble 融合"),
            ("三地平線", "D+1 日度 · D+5 週度 · D+20 月度 · 共 6 模型對照"),
            ("Purged Walk-Forward", "initial_train=252 · test=63 · step=63 · embargo=20"),
            ("Quality Gates", "leakage_scan / walk_forward_cv / purge_embargo / baseline_compare 等 9 項全過"),
        ],
        fig_path=str(FIGS / "fold_stability.png"),
    )
    build_phase_slide(
        prs, 8, TOTAL,
        tag="S8 · PHASE 3",
        phase_num=3,
        zh_title="擴充分析",
        en_subtitle="4 分析 · 9 models × 3 cost scenarios · cross-horizon · pillar · case",
        headline="在三種手續費情境下，D+20 是唯一「全情境皆正報酬」的地平線",
        color=BRAND["indigo"],
        bullets=[
            ("成本敏感度", "9 模型 × 3 情境（std 0.58% / disc 0.36% / cons 0.73%）"),
            ("跨地平線", "3 引擎 × 3 地平線 · D+20 皆正 Sharpe 0.70-0.81"),
            ("支柱貢獻", "normalized importance 跨 6 模型平均 · trend+risk 約 82%"),
            ("個案研究", "2330 / 2317 / 2454 / 2303 · 2454 最強命中 +17.13% edge"),
        ],
        fig_path=str(FIGS / "phase3_cost_sensitivity_heatmap.png"),
    )
    build_phase_slide(
        prs, 9, TOTAL,
        tag="S9 · PHASE 4",
        phase_num=4,
        zh_title="模型治理與風控",
        en_subtitle="Model Card × 6 · Drift Monitoring · DSR · Signal Decay",
        headline="把黑盒打開：每個模型都有完整的 Model Card 與漂移監控",
        color=BRAND["violet"],
        bullets=[
            ("Model Card", "6 joblib (lgb/xgb × D1/D5/D20) · 每個都有 card JSON 描述訓練資料/超參/限制"),
            ("Drift Detection", "特徵漂移、標籤漂移、概念漂移 三層監控"),
            ("DSR (Data Subject Rights)", "完整資料來源紀錄 + 使用授權 + 資料保留政策"),
            ("Signal Decay", "OOS 期間訊號品質隨時間衰減的量化追蹤"),
        ],
        fig_path=None,
    )
    build_phase_slide(
        prs, 10, TOTAL,
        tag="S10 · PHASE 5",
        phase_num=5,
        zh_title="文本與情緒因子",
        en_subtitle="jieba + SnowNLP · 9,655 keywords → 500 · 30 txt + 7 sent features",
        headline="用中文 NLP 把新聞文本變成 37 個可用特徵",
        color=BRAND["amber"],
        bullets=[
            ("文本處理", "jieba 分詞 · stopword 清理 · PMI 詞彙擴充"),
            ("Chi²/MI/Lift 選詞", "9,655 候選關鍵字 → 500 高訊號關鍵字"),
            ("情緒因子", "SnowNLP 情感分數 · 7 個聚合特徵（日/週/月三尺度）"),
            ("結果", "覆蓋率 heatmap / 平台佔比 / 詞雲 / 關鍵字 top"),
        ],
        fig_path=str(FIGS / "text_top_keywords.png"),
    )
    build_phase_slide(
        prs, 11, TOTAL,
        tag="S11 · PHASE 6",
        phase_num=6,
        zh_title="深度驗證",
        en_subtitle="LOPO (9 pillars) · Threshold sweep (21 points) · Single-stock 2454",
        headline="三層深度驗證：從「支柱級」到「個股級」全部拆開看",
        color=BRAND["rose"],
        bullets=[
            ("LOPO 驗證", "9 支柱逐一移除重訓 · 量化真實邊際貢獻（risk 最大 +1.39 bps）"),
            ("閾值敏感度", "t ∈ [0.30, 0.50] 21 點掃描 + Top-K precision"),
            ("個股深案", "2454 聯發科 OOS 213 日 · 月度命中拆解"),
            ("補件驗證", "修復 Phase 2 feature_stability FAIL bug · 最終 9/9 gates"),
        ],
        fig_path=str(FIGS / "lopo_pillar_contribution_D20.png"),
    )

    # Evidence slides
    build_evidence_lopo(prs, 12, TOTAL)
    build_evidence_threshold(prs, 13, TOTAL)
    build_evidence_2454(prs, 14, TOTAL)

    # Product
    build_product_dashboard(prs, 15, TOTAL)

    # Governance
    build_governance(prs, 16, TOTAL)

    # Limitations
    build_limitations(prs, 17, TOTAL)

    # Next Steps
    build_next_steps(prs, 18, TOTAL)

    # Takeaways
    build_takeaways(prs, 19, TOTAL)

    # Thank You
    build_thank_you(prs, 20, TOTAL)

    # Save
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"[OK] Saved: {OUT}")
    print(f"     Slides: {len(prs.slides)}")
    print(f"     Size:   {OUT.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
